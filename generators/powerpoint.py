"""
GFCC-style full Mass deck (community footer, poster dividers). Readings from API/USCCB.

1920×1080 landscape. Slide fill uses the liturgical calendar color; body/muted/emphasis
text colors are chosen for contrast (never matching the background).
"""

from __future__ import annotations

import colorsys
import datetime as _dt
import math
import re
from copy import deepcopy
from dataclasses import dataclass, replace
from pathlib import Path
from typing import Any, List, Mapping, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from services.community_config import get_community_name, get_logo_path
from services.hymn_library import get_hymn
from services.hymn_typography import HymnTypographySettings, typography_for_hymn_slide
from services.mass_text_format import (
    clean_lyrics_for_projection,
    ensure_lyric_section_breaks,
    parse_structured_lyric_sections_typed,
    pick_hymn_lyrics_for_slides,
    strip_reading_verse_markers,
)
from services.prayer_service import get_our_father, get_prayer
from services.prayer_templates import PENITENTIAL_ACT
from services.responsorial_reading import responsorial_section_title
from . import gfcc_flow_content as GFCC
from .deck_template import (
    COLLECTION_FOOTER_TOP,
    COLLECTION_TITLE_H,
    CONTENT_BOTTOM_GAP,
    DIALOGUE_BOTTOM_GAP,
    FOOTER_HEIGHT,
    FOOTER_TOP_OFFSET,
    HYMN_BODY_TOP_OFFSET,
    HYMN_TITLE_BOX_H,
    LOGO_GAP,
    LOGO_TOP,
    MARGIN_SIDE,
    MARGIN_TOP,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    _BRAND_BAND,
    _EMU_PER_INCH,
    _HYMN_DUAL_BOTTOM_CONT,
    _HYMN_DUAL_BOTTOM_FIRST,
    _HYMN_DUAL_CONT_BOX_H,
    _HYMN_DUAL_FIRST_BOX_H,
    _HYMN_DUAL_TOP_CONT,
    _HYMN_DUAL_TOP_FIRST,
    _HYMN_TITLE_TOP,
    _LYRIC_BODY_BOTTOM_MARGIN,
    _LOGO_MAX_H,
    _LOGO_MAX_W,
    _LYRIC_SAFE_SIDE_RATIO,
    _LYRIC_TEXTBOX_WIDTH_RATIO,
    _LYRIC_TF_SIDE_MARGIN,
)

_BG = RGBColor(18, 18, 22)
_GOLD_FALLBACK = RGBColor(220, 170, 90)
_BODY = RGBColor(245, 245, 245)
_MUTED = RGBColor(155, 155, 165)

# --- Theme 1 (the app's default deck theme — "LiturgyFlowTemplate1") -----------
# Built-in deck themes (selected via ``custom_theme.id`` from Mass Setup).
# Theme 1 keeps amber accents; Themes 2–3 are strict mono (bg ↔ text swapped),
# including titles, hymns, and Mass dividers (no liturgical-season tint).
_THEME1_BG = RGBColor(0x00, 0x00, 0x00)        # black background (all non-divider slides)
_THEME1_PRIMARY = RGBColor(0xF0, 0xFD, 0xF4)   # off-white body text
_THEME1_MUTED = RGBColor(0xBC, 0xBA, 0xC6)     # muted secondary text
_THEME1_EMPHASIS = RGBColor(0xFF, 0xB8, 0x00)  # amber/gold accent / highlighted text
_THEME1_FONT = "Arial"                          # dialogue & reading body font

_THEME2_BG = RGBColor(0x00, 0x00, 0x00)        # Midnight — black bg
_THEME2_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)   # pure white text (titles + body)
_THEME2_MUTED = RGBColor(0xC8, 0xC8, 0xC8)
_THEME2_EMPHASIS = RGBColor(0xFF, 0xFF, 0xFF)

_THEME3_BG = RGBColor(0xFF, 0xFF, 0xFF)        # Paper — white bg
_THEME3_PRIMARY = RGBColor(0x00, 0x00, 0x00)   # pure black text (titles + body)
_THEME3_MUTED = RGBColor(0x4A, 0x4A, 0x4A)
_THEME3_EMPHASIS = RGBColor(0x00, 0x00, 0x00)

_DECK_THEME_IDS = frozenset({"theme1", "theme2", "theme3"})

_SLIDE_TEXT_PT = 55
_FOOTER_PT = 13

_MAX_CHARS_READING = 820
_MAX_MARKED_BODY = 2600
# Hymn / lyrics slides (black screen, gold title, white ALL CAPS body — projector style)
_LYRIC_MAX_LINES_PER_SLIDE = 6
_HYMN_DUAL_SOLO_PARAGRAPH_THRESHOLD = 6
_HYMN_TITLE_PT = 38.5
_HYMN_BODY_PT = 56.0
_HYMN_REF_TITLE_PT = 36.0
_HYMN_REF_BODY_PT = 75.0
_HYMN_REF_BODY_PT_MIN = 68.0
_HYMN_REF_BODY_FONT = "Poppins"
_HYMN_REF_LINE_SPACING = 0.7
_LYRIC_TITLE_DISPLAY_PT = _HYMN_TITLE_PT
_LYRIC_BODY_DISPLAY_PT = _HYMN_BODY_PT
_HYMN_BG = RGBColor(0, 0, 0)
_HYMN_GOLD_TITLE = RGBColor(255, 204, 77)
_HYMN_BODY_WHITE = RGBColor(255, 255, 255)
_HYMN_CHORUS_COLOR = RGBColor(0xFF, 0xB8, 0x00)
_HYMN_PAREN_COLOR = RGBColor(0x87, 0xCE, 0xEB)
_PAREN_IN_LYRICS_RE = re.compile(r"\([^)]*\)")
_HYMN_BRAND_WHITE = RGBColor(255, 255, 255)
_HYMN_FOOTER_MUTED = RGBColor(140, 140, 145)
_HYMN_TITLE_FONT = "Georgia"
_HYMN_BODY_FONT = "Poppins Bold"
# Section titles (LiturgyFlowTemplate1): Georgia bold underline, accent color, pinned
# to the top of the slide. Two sizes are used across the deck.
_SECTION_TITLE_FONT = "Georgia"
_SECTION_TITLE_TOP_IN = 0.12
_SECTION_TITLE_H_IN = 0.95
_SECTION_TITLE_PT_SMALL = 36.0
_SECTION_TITLE_PT_LARGE = 50.0
_GOSPEL_ACCLAMATION_BODY_PT = 69.0
_GOSPEL_ACCLAMATION_BODY_FONT = "Poppins Bold"
_COMMUNITY_HEADER_PT = 15
_LYRIC_MIN_WORDS_PER_LINE = 3
_LYRIC_MIN_PT = 40
_LYRIC_MAX_PT = int(_HYMN_REF_BODY_PT)
_LYRIC_SOFT_WRAP_CHARS = 46
_LYRIC_FIT_WIDTH_SAFETY = 0.96
_LYRIC_FIT_HEIGHT_SAFETY = 0.90
# Em->inch width fudge for Poppins Bold ALL CAPS on the projector surface. The
# previous 0.80 under-measured wide lines, so the auto-fit never shrank them and
# the text wrapped/overflowed the box (covering the title on dense slides). 0.92
# still under-measured a few long verse lines (e.g. Your Heart Today dual slide 2),
# so PowerPoint soft-wrapped one more visual line than the model predicted; with
# MIDDLE anchor on a top-flush dual box that spill reads as top overflow. 0.94
# matches observed Poppins Bold cap widths (manual 72 pt fit on that slide).
_LYRIC_WIDTH_CALIBRATION = 0.94
# Rendered line pitch vs. nominal at the configured line spacing (Poppins runs a
# little taller than the bare spacing factor implies).
_LYRIC_LINE_PITCH_FACTOR = 1.20

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
_REFERENCE_DECK_FILENAME = "GCCC24May2026_Eastertide.pptx"
_REFERENCE_SLIDE_PRE_MASS = 0
_REFERENCE_SLIDE_PENITENTIAL = (7, 8)
_REFERENCE_SLIDE_KYRIE = 9
_LAMB_OF_GOD_TEMPLATE_FILENAME = "lamb_of_god_slide.pptx"
_LAMB_OF_GOD_SLIDE_INDEX = 0
_SIGN_OF_PEACE_TEMPLATE_FILENAME = "sign_of_peace_slide.pptx"
_SIGN_OF_PEACE_SLIDE_INDEX = 0
_GLORIA_TEMPLATE_FILENAME = "gloria_slides.pptx"
_KYRIE_TEMPLATE_FILENAME = "kyrie_slide.pptx"
_KYRIE_SLIDE_INDEX = 0
_LOTW_TITLE_IMAGE_FILENAME = "liturgy_of_the_word_title.png"
_LOTW_TITLE_TEMPLATE_FILENAME = "liturgy_of_the_word_title_slide.pptx"
_LOTW_TITLE_SLIDE_INDEX = 0

# Bundled full-bleed poster art for the Liturgy of the Word / Eucharist divider
# slides. Four selectable designs each; chosen in Mass setup, default to design 1.
_POSTER_REFERENCE_DIR = _PROJECT_ROOT / "data" / "reference" / "posters"
_LOTW_POSTER_IDS = ("lotw1", "lotw2", "lotw3", "lotw4")
_LOTE_POSTER_IDS = ("lote1", "lote2", "lote3", "lote4")
_LOTW_POSTER_DEFAULT = "lotw1"
_LOTE_POSTER_DEFAULT = "lote1"
_GOSPEL_ACCLAMATION_TEMPLATE_FILENAME = "gospel_acclamation_slides.pptx"
_APOSTLES_CREED_TEMPLATE_FILENAME = "apostles_creed_slides.pptx"
_APOSTLES_CREED_TITLE = "Apostles' Creed"
_NICENE_CREED_TEMPLATE_FILENAME = "nicene_creed_slides.pptx"
_NICENE_CREED_TITLE = "Nicene Creed"

# --- Master template (per-theme source of truth for fixed liturgical slides) -----
# Theme 1 (LFTemplate1.pptx) is the authored deck. Themes 2–3 are baked offline
# from Theme 1 via ``scripts/bake_deck_theme_templates.py`` so generation can clone
# them verbatim without a fragile live color-remapping pass.
# Fixed (week-invariant) slides are cloned by 0-based index; parish footer/branding
# is swapped in. Dynamic slides inject text into the matching layout.
_MASTER_TEMPLATE_BY_THEME = {
    "theme1": "LFTemplate1.pptx",
    "theme2": "LFTemplate2-midnight.pptx",
    "theme3": "LFTemplate3-paper.pptx",
}
_MASTER_TEMPLATE_FILENAME = _MASTER_TEMPLATE_BY_THEME["theme1"]
_MASTER_SLIDE = {
    "pre_mass": 0,
    "introductory_rites": 6,
    "penitential": (8, 9, 10),
    "kyrie": 11,
    "gloria": (12, 13, 14, 15),
    "lotw_prayer": 16,
    "first_reading": 18,
    "psalm": 19,
    "second_reading": 20,
    "gospel_acclamation_alleluia": 21,
    "gospel_acclamation_dialogue": 22,
    "gospel_acclamation_end": 23,
    "nicene_creed": (25, 26, 27, 28),
    "prayer_faithful": (30, 31),
    "lote_pray_brethren": 36,
    "preface_dialogue": 38,
    # A "Liturgy of the Eucharist" poster slide was inserted before the Sanctus,
    # shifting every section from the Sanctus onward by +1.
    "sanctus": 40,
    "mystery_of_faith": 42,
    "great_amen": 44,
    "sign_of_peace": 49,
    "lamb_of_god": 50,
    "communion_rite": 51,
    "welcoming_newcomers": 60,
    "mass_collection": 61,
    "post_communion": 58,
    "confession": 62,
    "final_blessing": 63,
}
_master_templates: dict[str, Presentation] = {}
_ACTIVE_DECK_THEME_ID = "theme1"

# The new master dropped the Food Sponsors slide; reuse the previous template's
# layout (same Georgia-underline title + Arial Black name style) as the donor.
_FOOD_SPONSORS_TEMPLATE_FILENAME = "food_sponsors_template.pptx"

# Placeholder strings baked into the welcoming / collection / sponsor cards,
# replaced in-place when cloning so the authored format is preserved 1:1.
_TPL_CHURCH_NAME = "[insert church name]"
_TPL_COLLECTION_AMOUNT = "[ENTER AMOUNT]"
_TPL_COLLECTION_DATE = "[DATE]"
_TPL_SPONSOR_EXAMPLE = "KESHI GONZALES"
_TPL_SPONSOR_PLACEHOLDER = "[insert sponsor name]"

# Per-language Our Father slide sets extracted from the authored templates. The
# language picker clones the matching deck verbatim (recolored to amber/black);
# languages without a bundled deck fall back to the procedural builder.
_OUR_FATHER_DECKS = {
    "english": "our_father_english.pptx",
    "tagalog": "our_father_tagalog.pptx",
}

# Placeholder strings baked into the master reading cards; replaced in-place with
# the week's dynamic references when cloning (layout/typography preserved 1:1).
_TPL_FIRST_CITATION = "Jeremiah 20:10-13"
_TPL_SECOND_CITATION = "Romans 5:12-15"
_TPL_PSALM_SECTION = "Responsorial Psalm"
_TPL_PSALM_REF = "Psalm 69:8-10, 14, 17, 33-35"
_TPL_PSALM_ANTIPHON = "LORD, IN YOUR GREAT LOVE, ANSWER ME."

_REFERENCE_FOOTER_ZONE_TOP = int(SLIDE_HEIGHT * 0.78)
# When cloning master/reference slides we only want to drop the very-bottom footer
# tag (parish/section label, anchored ~10.3in). Some templates place legitimate
# liturgical responses (e.g. "All: It is right and just.", "All: Thanks be to
# God.") as low as ~9.4in, so the clone uses this tighter band — anything at/above
# it is treated as footer, everything above is kept as real content.
_CLONE_FOOTER_ZONE_TOP = int(SLIDE_HEIGHT * 0.90)
_reference_mass_deck: Optional[Presentation] = None
_lamb_of_god_template: Optional[Presentation] = None
_sign_of_peace_template: Optional[Presentation] = None
_gloria_template: Optional[Presentation] = None
_kyrie_template: Optional[Presentation] = None
_lotw_title_template: Optional[Presentation] = None
_gospel_acclamation_template: Optional[Presentation] = None
_apostles_creed_template: Optional[Presentation] = None
_nicene_creed_template: Optional[Presentation] = None


@dataclass(frozen=True)
class DeckBrandingOptions:
    include_logo: bool = True
    include_name: bool = True
    # Developer option: when False, the bottom community/section footer tag (the
    # small ~12pt line) is omitted from every slide. Off by default.
    include_footer: bool = False
    # When set (e.g. guest demo watermark), footer brand uses this instead of
    # the parish community name — without enabling parish logo/name branding.
    footer_brand: str = ""


_deck_branding = DeckBrandingOptions()
_OUTPUT_DIR = _PROJECT_ROOT / "outputs"


def _footer_brand_line() -> str:
    brand = (_deck_branding.footer_brand or "").strip()
    if brand:
        return brand
    if _deck_branding.include_name:
        return get_community_name()
    return ""


@dataclass(frozen=True)
class SlideTheme:
    """Single source of truth for every deck color.

    Two surfaces share one theme:
    - The *liturgical surface* (readings, prayers, dialogue, collection) is filled
      with ``bg``; its text uses ``primary``/``muted``/``emphasis``.
    - The *projector surface* (hymn, lyric, gospel-acclamation slides) uses the
      ``hymn_*``/``chorus_accent``/``paren_accent`` roles.

    Theme 1 tints only Mass dividers by liturgical season. Themes 2–3 set
    ``mono_surfaces=True`` so dividers match the same black/white palette.
    """

    bg: RGBColor
    primary: RGBColor
    muted: RGBColor
    emphasis: RGBColor
    font_name: str = "Calibri"
    # Projector surface (hymn / lyric / gospel-acclamation slides).
    hymn_bg: RGBColor = _HYMN_BG
    hymn_title: RGBColor = _HYMN_GOLD_TITLE
    hymn_body: RGBColor = _HYMN_BODY_WHITE
    chorus_accent: RGBColor = _HYMN_CHORUS_COLOR
    paren_accent: RGBColor = _HYMN_PAREN_COLOR
    hymn_brand: RGBColor = _HYMN_BRAND_WHITE
    footer_muted: RGBColor = _HYMN_FOOTER_MUTED
    # Mass-divider surface (season-tinted for Theme 1; mono for Themes 2–3).
    divider_bg: RGBColor = _BG
    divider_primary: RGBColor = _BODY
    divider_muted: RGBColor = _MUTED
    divider_emphasis: RGBColor = _GOLD_FALLBACK
    mono_surfaces: bool = False


_ACTIVE_FONT = "Calibri"


def _clamp_byte(x: float) -> int:
    return max(0, min(255, int(round(x))))


def _hex_to_rgb(value: Any) -> Optional[RGBColor]:
    text = str(value or "").strip().lstrip("#")
    if len(text) != 6:
        return None
    try:
        return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except ValueError:
        return None


def _rel_lum(r: int, g: int, b: int) -> float:
    def lin(c: int) -> float:
        cs = c / 255.0
        return cs / 12.92 if cs <= 0.03928 else ((cs + 0.055) / 1.055) ** 2.4

    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def _accent_on_dark(rgb: RGBColor) -> RGBColor:
    """Brighten a season hue so it stays legible on the dark projector surface.

    Near-neutral seasons (white/cream) have no usable hue, so they fall back to
    the warm gold the projector slides have always used.
    """
    r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    h, _l, s = colorsys.rgb_to_hls(r / 255.0, g / 255.0, b / 255.0)
    if s < 0.12:
        return _HYMN_GOLD_TITLE
    rr, gg, bb = colorsys.hls_to_rgb(h, 0.62, max(s, 0.55))
    return RGBColor(_clamp_byte(rr * 255), _clamp_byte(gg * 255), _clamp_byte(bb * 255))


def _projector_accent(emphasis: RGBColor, bg: RGBColor) -> RGBColor:
    """Season-aware accent for the dark projector surface (hymn titles, chorus).

    Reuses the theme ``emphasis`` when it is already bright enough for a black
    screen; otherwise derives a brightened version of the season color.
    """
    if _rel_lum(int(emphasis[0]), int(emphasis[1]), int(emphasis[2])) >= 0.32:
        return emphasis
    return _accent_on_dark(bg)


def _theme_with_roles(
    bg: RGBColor,
    primary: RGBColor,
    muted: RGBColor,
    emphasis: RGBColor,
    font_name: str = "Calibri",
) -> SlideTheme:
    """Build a SlideTheme, deriving the season-aware projector roles from the
    liturgical surface colors. The dark-screen body/paren/brand/footer tones stay
    at their proven projector defaults (white on black); only the accent shifts."""
    accent = _projector_accent(emphasis, bg)
    return SlideTheme(
        bg=bg,
        primary=primary,
        muted=muted,
        emphasis=emphasis,
        font_name=font_name,
        hymn_title=accent,
        chorus_accent=accent,
    )


def _mono_slide_theme(*, dark: bool) -> SlideTheme:
    """Strict black↔white palette for Themes 2–3 (all text roles match, incl. dividers)."""
    if dark:
        bg, primary, muted, emphasis = (
            _THEME2_BG,
            _THEME2_PRIMARY,
            _THEME2_MUTED,
            _THEME2_EMPHASIS,
        )
        footer = RGBColor(0xA0, 0xA0, 0xA0)
    else:
        bg, primary, muted, emphasis = (
            _THEME3_BG,
            _THEME3_PRIMARY,
            _THEME3_MUTED,
            _THEME3_EMPHASIS,
        )
        footer = RGBColor(0x6B, 0x6B, 0x6B)
    return SlideTheme(
        bg=bg,
        primary=primary,
        muted=muted,
        emphasis=emphasis,
        font_name=_THEME1_FONT,
        hymn_bg=bg,
        hymn_title=primary,
        hymn_body=primary,
        chorus_accent=primary,
        paren_accent=muted,
        hymn_brand=primary,
        footer_muted=footer,
        divider_bg=bg,
        divider_primary=primary,
        divider_muted=muted,
        divider_emphasis=emphasis,
        mono_surfaces=True,
    )


def _resolve_deck_theme_id(custom_theme: Optional[Mapping[str, Any]]) -> str:
    """Map ``custom_theme`` payload to a built-in deck theme id."""
    if not custom_theme:
        return "theme1"
    tid = str(custom_theme.get("id") or "").strip().lower()
    if tid in _DECK_THEME_IDS:
        return tid
    name = str(custom_theme.get("name") or "").strip().lower()
    if "midnight" in name or name in {"theme 2", "theme2", "black", "noir"}:
        return "theme2"
    if "paper" in name or name in {"theme 3", "theme3", "white", "daylight"}:
        return "theme3"
    if "liturgy" in name or name in {"theme 1", "theme1"}:
        return "theme1"
    return "theme1"


def _season_surface_roles(
    liturgical_color: Optional[Mapping[str, Any]],
) -> Tuple[RGBColor, RGBColor, RGBColor, RGBColor]:
    """Season-derived ``(bg, primary, muted, emphasis)`` for the Mass divider only.

    Cream/white seasons use near-black body text; darker greens/purples/reds use warm
    off-white and a gold emphasis so roles never mirror the fill. Falls back to the
    Theme 1 palette when no liturgical color is available.
    """
    if not (liturgical_color and "rgb" in liturgical_color):
        return _THEME1_BG, _THEME1_PRIMARY, _THEME1_MUTED, _THEME1_EMPHASIS
    r, g, b = (int(c) for c in liturgical_color["rgb"])
    bg = RGBColor(r, g, b)
    if _rel_lum(r, g, b) > 0.55:
        primary = RGBColor(26, 26, 30)
        muted = RGBColor(95, 93, 104)
        emphasis = RGBColor(
            _clamp_byte(r * 0.28 + 18),
            _clamp_byte(g * 0.27 + 16),
            _clamp_byte(b * 0.30 + 20),
        )
    else:
        primary = RGBColor(250, 248, 244)
        muted = RGBColor(188, 186, 198)
        emphasis = RGBColor(255, 218, 145)
    return bg, primary, muted, emphasis


def _build_slide_theme(
    liturgical_color: Optional[Mapping[str, Any]],
    custom_theme: Optional[Mapping[str, Any]] = None,
) -> SlideTheme:
    """Build the deck palette for the selected built-in theme.

    - ``theme1`` (Liturgy Flow): fixed black + amber; season color tints Mass dividers only.
    - ``theme2`` (Midnight): black background, all-white text including titles & dividers.
    - ``theme3`` (Paper): white background, all-black text including titles & dividers.

    ``custom_theme.id`` (preferred) or ``custom_theme.name`` selects the theme.
    """
    theme_id = _resolve_deck_theme_id(custom_theme)
    if theme_id == "theme2":
        return _mono_slide_theme(dark=True)
    if theme_id == "theme3":
        return _mono_slide_theme(dark=False)
    base = _theme_with_roles(
        _THEME1_BG, _THEME1_PRIMARY, _THEME1_MUTED, _THEME1_EMPHASIS, _THEME1_FONT
    )
    d_bg, d_primary, d_muted, d_emphasis = _season_surface_roles(liturgical_color)
    return replace(
        base,
        divider_bg=d_bg,
        divider_primary=d_primary,
        divider_muted=d_muted,
        divider_emphasis=d_emphasis,
    )


def _accent(liturgical_color: Optional[Mapping[str, Any]]) -> RGBColor:
    """Backward-compatible single accent RGB (emphasis tone for callers that only need one color)."""
    return _build_slide_theme(liturgical_color).emphasis


# Active per-deck theme. Set at the start of generate_mass_ppt so the projector
# helpers (hymn / lyric / gospel-acclamation) can read role colors without
# threading `theme` through every signature, mirroring the _ACTIVE_FONT pattern.
_ACTIVE_THEME: SlideTheme = _build_slide_theme(None)


def _content_top():
    """Vertical start for slide body text (below logo + parish name)."""
    if not (_deck_branding.include_logo or _deck_branding.include_name):
        return MARGIN_TOP
    return MARGIN_TOP + _BRAND_BAND


def _apply_slide_branding(slide, theme: SlideTheme) -> None:
    """Top-left logo and parish name (GFCC reference deck style)."""
    if not _deck_branding.include_logo and not _deck_branding.include_name:
        return
    logo = get_logo_path() if _deck_branding.include_logo else None
    name = get_community_name() if _deck_branding.include_name else ""
    top = LOGO_TOP
    cursor_left = MARGIN_SIDE
    if logo and logo.is_file():
        pic = slide.shapes.add_picture(str(logo), cursor_left, top, width=_LOGO_MAX_W)
        if pic.height > _LOGO_MAX_H:
            scale = _LOGO_MAX_H / pic.height
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
        cursor_left = pic.left + pic.width + LOGO_GAP

    if _deck_branding.include_name and name:
        name_w = SLIDE_WIDTH - cursor_left - MARGIN_SIDE
        nb = slide.shapes.add_textbox(cursor_left, top, name_w, _LOGO_MAX_H)
        tf = nb.text_frame
        _prep_tf(tf)
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = BODY_VERTICAL_ANCHOR
        p0 = tf.paragraphs[0]
        p0.text = name
        _style_para(p0, size_pt=_COMMUNITY_HEADER_PT, color=theme.primary, bold=True)
        p0.alignment = PP_ALIGN.LEFT


def _apply_hymn_branding(slide) -> None:
    """Top-left logo + parish name on black hymn slides (white text)."""
    if not _deck_branding.include_logo and not _deck_branding.include_name:
        return
    logo = get_logo_path() if _deck_branding.include_logo else None
    name = get_community_name() if _deck_branding.include_name else ""
    top = LOGO_TOP
    cursor_left = MARGIN_SIDE
    if logo and logo.is_file():
        pic = slide.shapes.add_picture(str(logo), cursor_left, top, width=_LOGO_MAX_W)
        if pic.height > _LOGO_MAX_H:
            scale = _LOGO_MAX_H / pic.height
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
        cursor_left = pic.left + pic.width + LOGO_GAP
    if _deck_branding.include_name and name:
        name_w = SLIDE_WIDTH - cursor_left - MARGIN_SIDE
        nb = slide.shapes.add_textbox(cursor_left, top, name_w, _LOGO_MAX_H)
        tf = nb.text_frame
        _prep_tf(tf)
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = BODY_VERTICAL_ANCHOR
        p0 = tf.paragraphs[0]
        p0.text = name
        _style_para(p0, size_pt=_COMMUNITY_HEADER_PT, color=_ACTIVE_THEME.hymn_brand, bold=True)
        p0.alignment = PP_ALIGN.LEFT


def _add_hymn_footer(slide, footer_section: str) -> None:
    if not _deck_branding.include_footer:
        return
    lx = MARGIN_SIDE
    w = SLIDE_WIDTH - 2 * MARGIN_SIDE
    y = SLIDE_HEIGHT - FOOTER_TOP_OFFSET
    foot = slide.shapes.add_textbox(lx, y, w, FOOTER_HEIGHT)
    tf = foot.text_frame
    _prep_tf(tf)
    tf.clear()
    p0 = tf.paragraphs[0]
    brand = _footer_brand_line()
    if brand:
        p0.text = brand
        _style_para(p0, size_pt=_FOOTER_PT, color=_ACTIVE_THEME.footer_muted, bold=True)
    p1 = tf.add_paragraph()
    p1.text = footer_section
    _style_para(p1, size_pt=_FOOTER_PT - 1, color=_ACTIVE_THEME.hymn_title, bold=False)
    p1.space_before = Pt(2)


def _layout_blank(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in (layout.name or "").lower():
            return layout
    return prs.slide_layouts[-1]


def _video_poster_png(title: str, theme: Optional[SlideTheme] = None) -> Path:
    """Solid poster frame for embedded movies (python-pptx requires one)."""
    import tempfile

    from PIL import Image, ImageDraw, ImageFont

    w, h = 1920, 1080
    bg = (16, 18, 22)
    if theme is not None:
        try:
            c = theme.bg
            bg = (int(c[0]), int(c[1]), int(c[2]))
        except Exception:
            pass
    img = Image.new("RGB", (w, h), bg)
    draw = ImageDraw.Draw(img)
    label = (title or "Video").strip() or "Video"
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 56)
    except Exception:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), label, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((w - tw) / 2, (h - th) / 2 - 40), label, fill=(240, 240, 240), font=font)
    hint = "▶  Video for Mass"
    try:
        font2 = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
    except Exception:
        font2 = font
    bbox2 = draw.textbbox((0, 0), hint, font=font2)
    tw2, th2 = bbox2[2] - bbox2[0], bbox2[3] - bbox2[1]
    draw.text(((w - tw2) / 2, (h - th2) / 2 + 40), hint, fill=(200, 200, 200), font=font2)
    tmp = tempfile.NamedTemporaryFile(prefix="verbum_vid_poster_", suffix=".png", delete=False)
    img.save(tmp.name, format="PNG")
    tmp.close()
    return Path(tmp.name)


def _add_video_replacement_slide(
    prs: Presentation,
    video_path: Path | str,
    *,
    title: str,
    theme: Optional[SlideTheme] = None,
) -> bool:
    """Full-bleed single slide with an embedded MP4 (replaces lyric/prayer slides)."""
    p = Path(video_path)
    if not p.is_file():
        return False
    poster: Optional[Path] = None
    try:
        poster = _video_poster_png(title, theme)
        slide = prs.slides.add_slide(_layout_blank(prs))
        slide.shapes.add_movie(
            str(p.resolve()),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
            poster_frame_image=str(poster),
            mime_type="video/mp4",
        )
        return True
    except Exception:
        return False
    finally:
        if poster is not None:
            try:
                poster.unlink(missing_ok=True)
            except Exception:
                pass


def _add_liturgical_poster_full_slide(prs: Presentation, png_path: Path) -> None:
    """
    Embed the 16×9 liturgical poster (PNG) as one full-bleed slide for projection.

    Expects ``png_path`` to match the slide aspect ratio (1920×1080 → 20″×11.25″ deck).
    """
    p = Path(png_path)
    if not p.is_file():
        return
    slide = prs.slides.add_slide(_layout_blank(prs))
    slide.shapes.add_picture(
        str(p.resolve()),
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def _set_slide_bg(slide, rgb: RGBColor):
    fi = slide.background.fill
    fi.solid()
    fi.fore_color.rgb = rgb


def _prep_tf(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)


# Deck-wide vertical anchor for body text. Every slide body centers its text
# block in the available space so short and long content sit consistently
# instead of some slides hugging the top and others centering.
BODY_VERTICAL_ANCHOR = MSO_ANCHOR.MIDDLE
BODY_PARAGRAPH_GAP_PT = 5


def _prep_body_tf(tf):
    """Prep + clear a body textframe and apply the deck-wide vertical anchor."""
    _prep_tf(tf)
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = BODY_VERTICAL_ANCHOR


def _fill_centered_box(
    tf,
    text: str,
    *,
    size_pt: int,
    color: RGBColor,
    bold: bool = False,
    space_after_pt: int = BODY_PARAGRAPH_GAP_PT,
) -> None:
    """Fill a textframe with centered paragraphs (split on blank lines).

    Single shared path so spacing and alignment are identical across reading,
    prayer, and other plain body slides. Does not set the vertical anchor; callers
    prep the box with ``_prep_body_tf`` so the anchor policy lives in one place.
    """
    tf.clear()
    raw = (text or "").strip()
    parts = [b.strip() for b in raw.split("\n\n") if b.strip()] or ([raw] if raw else [""])
    first = True
    for block in parts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = block
        _style_para(p, size_pt=size_pt, color=color, bold=bold)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(space_after_pt)


def _prep_hymn_lyric_tf(tf):
    """Full-width lyric textbox on title slides (side inset 0; vertical padding from _prep_tf)."""
    _prep_tf(tf)
    tf.margin_left = _LYRIC_TF_SIDE_MARGIN
    tf.margin_right = _LYRIC_TF_SIDE_MARGIN


def _prep_hymn_lyric_tf_full_bleed(tf):
    """Title-less lyric slides: textbox flush to all four slide edges."""
    _prep_tf(tf)
    tf.margin_left = _LYRIC_TF_SIDE_MARGIN
    tf.margin_right = _LYRIC_TF_SIDE_MARGIN
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)


def _lyric_textbox_geometry(slide_width: int) -> Tuple[int, int]:
    """Return (left, width) so the textbox spans the full slide width."""
    width = int(slide_width * _LYRIC_TEXTBOX_WIDTH_RATIO)
    left = int(slide_width * _LYRIC_SAFE_SIDE_RATIO)
    return left, width


def _lyric_continuation_textbox_geometry(slide_width: int, slide_height: int) -> Tuple[int, int, int, int]:
    """Return (left, top, width, height) at 0% inset for title-less lyric slides."""
    return 0, 0, int(slide_width), int(slide_height)


def _word_count(line: str) -> int:
    return len((line or "").split())


def _paren_span_containing(text: str, idx: int) -> Optional[Tuple[int, int]]:
    """Return ``(start, end)`` of a ``(...)`` group that would be split at ``idx``."""
    depth = 0
    start: Optional[int] = None
    for i, ch in enumerate(text):
        if ch == "(":
            if depth == 0:
                start = i
            depth += 1
        elif ch == ")":
            if depth <= 0:
                continue
            depth -= 1
            if depth == 0 and start is not None:
                end = i + 1
                if start < idx < end:
                    return (start, end)
                start = None
    return None


def _break_splits_parens(text: str, idx: int) -> bool:
    return _paren_span_containing(text, idx) is not None


def _split_line_preserving_parens(text: str, idx: int) -> Tuple[str, str]:
    """Split at ``idx``, moving any bracketed lyric entirely to the next line."""
    span = _paren_span_containing(text, idx)
    if span is None:
        return text[:idx].strip(), text[idx:].strip(" ,-;:.")
    start, _end = span
    return text[:start].strip(), text[start:].strip(" ,-;:.")


def _first_toplevel_paren_span(text: str) -> Optional[Tuple[int, int]]:
    """Return ``(start, end)`` of the first top-level ``(...)`` group, or None."""
    depth = 0
    start: Optional[int] = None
    for i, ch in enumerate(text):
        if ch == "(":
            if depth == 0:
                start = i
            depth += 1
        elif ch == ")":
            if depth > 0:
                depth -= 1
                if depth == 0 and start is not None:
                    return (start, i + 1)
    return None


def _split_paren_echo_segments(line: str) -> List[Tuple[str, bool]]:
    """Peel a trailing parenthetical echo onto its own row.

    A line like ``main phrase (echo of the phrase)`` becomes two segments so the
    parenthetical always moves to the next row and is never broken mid-way. Short
    annotations (``(2x)``, ``(bis)``) and parentheticals with trailing text are
    left inline. Returns ``(segment, is_paren)`` pairs; ``is_paren`` segments must
    not be re-wrapped.
    """
    span = _first_toplevel_paren_span(line)
    if span is None:
        return [(line, False)]
    start, end = span
    before = line[:start].strip()
    inner = line[start + 1 : end - 1]
    # Only treat as an echo when there is lead-in text, the group runs to the end
    # of the line, and the parenthetical is a real phrase (not a tiny annotation).
    if not before or line[end:].strip() or len(inner.split()) < 2:
        return [(line, False)]
    return [(before, False), (line[start:end].strip(), True)]


def _lyric_paragraph_count(text: str) -> int:
    return len([ln for ln in (text or "").splitlines() if ln.strip()])


def _enforce_min_words_per_line(
    lines: List[str],
    min_words: int = _LYRIC_MIN_WORDS_PER_LINE,
) -> List[str]:
    """Merge orphan lines so each rendered line has at least ``min_words`` when possible."""
    if not lines:
        return lines
    merged: List[str] = []
    pending = ""
    for line in lines:
        text = (line or "").strip()
        if not text:
            continue
        if pending:
            text = f"{pending} {text}".strip()
            pending = ""
        while _word_count(text) < min_words:
            # keep short source lines intact if nothing else to merge with
            if not merged:
                break
            prev = merged.pop()
            text = f"{prev} {text}".strip()
        if _word_count(text) < min_words:
            pending = text
            continue
        merged.append(text)
    if pending:
        if merged:
            merged[-1] = f"{merged[-1]} {pending}".strip()
        else:
            merged.append(pending)
    return merged


def _style_para(
    p, *, size_pt, color, bold=False, italic=False, font_name=None, underline=False
):
    p.font.name = font_name or _ACTIVE_FONT
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.italic = italic
    p.font.underline = underline
    p.font.color.rgb = color
    _set_def_rpr_rgb(p, color)


def _font_rgb_or_none(font) -> Optional[RGBColor]:
    """Read a run/paragraph font's explicit RGB, or None when inherited/theme-based."""
    try:
        color = font.color
        if color is not None and color.type == MSO_COLOR_TYPE.RGB:
            return color.rgb
    except (AttributeError, KeyError, TypeError, ValueError):
        return None
    return None


def _rgb_hex(rgb: RGBColor) -> str:
    return f"{int(rgb[0]):02X}{int(rgb[1]):02X}{int(rgb[2]):02X}"


def _set_solid_srgb_on_rpr(r_pr, rgb: RGBColor) -> None:
    """Replace any color fill on ``a:rPr`` / ``a:defRPr`` with a single sRGB solidFill."""
    for child in list(r_pr):
        tag = child.tag
        if tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill"), qn("a:pattFill")):
            r_pr.remove(child)
    solid = etree.SubElement(r_pr, qn("a:solidFill"))
    srgb = etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", _rgb_hex(rgb))


def _set_run_rpr_rgb(run, rgb: RGBColor) -> None:
    """Write RGB onto the run's ``a:rPr`` (creates rPr when missing)."""
    try:
        r_el = run._r
    except (AttributeError, TypeError, ValueError):
        return
    r_pr = r_el.find(qn("a:rPr"))
    if r_pr is None:
        r_pr = etree.Element(qn("a:rPr"))
        r_el.insert(0, r_pr)
    _set_solid_srgb_on_rpr(r_pr, rgb)


def _set_def_rpr_rgb(para, rgb: RGBColor) -> None:
    """Write RGB onto paragraph ``a:defRPr`` and ``a:endParaRPr``.

    Master-template slides bake body color into ``defRPr`` (often #F0FDF4) and
    sometimes ``endParaRPr``. Recoloring only the run leaves those defaults intact;
    empty or partial runs then inherit white and vanish on the Paper theme.
    """
    try:
        p_el = para._p
    except (AttributeError, TypeError, ValueError):
        return
    p_pr = p_el.find(qn("a:pPr"))
    if p_pr is None:
        p_pr = etree.Element(qn("a:pPr"))
        p_el.insert(0, p_pr)
    for tag in (qn("a:defRPr"),):
        def_r_pr = p_pr.find(tag)
        if def_r_pr is None:
            def_r_pr = etree.SubElement(p_pr, tag)
        _set_solid_srgb_on_rpr(def_r_pr, rgb)
    # endParaRPr lives on the paragraph itself (sibling of pPr), not inside pPr.
    end_r_pr = p_el.find(qn("a:endParaRPr"))
    if end_r_pr is not None:
        _set_solid_srgb_on_rpr(end_r_pr, rgb)


def _apply_run_theme_color(run_or_para, rgb: RGBColor) -> None:
    """Set API color + underlying rPr/defRPr solidFill so nothing inherits white."""
    run_or_para.font.color.rgb = rgb
    # Paragraph objects have no ``_r``; runs do.
    if hasattr(run_or_para, "_r"):
        _set_run_rpr_rgb(run_or_para, rgb)


def _iter_slide_text_shapes(slide):
    """Yield text-bearing shapes, including those nested in groups."""

    def _walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _walk(shape.shapes)
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                yield shape

    yield from _walk(slide.shapes)


def _source_text_role(font) -> str:
    """Classify a cloned reference run as an accent (warm/gold) or body color.

    Reference decks bake warm-gold titles and light body text for their original
    (dark) background. We map those buckets onto the active liturgical roles so the
    cloned slide follows the season theme instead of the baked-in colors.
    """
    rgb = _font_rgb_or_none(font)
    if rgb is None:
        # Theme/scheme colors (e.g. BACKGROUND_1) have no usable RGB — treat as body
        # so Paper/Midnight themes don't leave white-on-white scheme text behind.
        return "primary"
    r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    h, lightness, s = colorsys.rgb_to_hls(r / 255.0, g / 255.0, b / 255.0)
    # A genuine accent color (the template's green highlight ~L0.73) maps to the
    # deck emphasis tone. Body text is either neutral (low saturation) or a
    # near-white off-white (e.g. #F0FDF4) that is technically saturated but very
    # light — those map to ``primary`` so only true highlights become amber.
    if s >= 0.25 and lightness <= 0.85:
        return "emphasis"
    return "primary"


def _theme_color_for_font(font, theme: SlideTheme) -> RGBColor:
    if getattr(font, "underline", False):
        return theme.emphasis
    role = _source_text_role(font)
    return theme.emphasis if role == "emphasis" else theme.primary


def _recolor_cloned_text_to_theme(slide, theme: SlideTheme) -> None:
    """Re-theme a cloned reference slide's text to the active liturgical surface.

    Cloned rite slides sit on ``theme.bg`` but keep colors baked for the original
    (usually dark) deck. We remap each run *and* paragraph ``defRPr`` so rites stay
    legible on Theme 1, Midnight, and Paper alike.

    Only the true bottom footer band (``_CLONE_FOOTER_ZONE_TOP``, ~90%) is skipped —
    the wider ``_REFERENCE_FOOTER_ZONE_TOP`` (~78%) still holds real dialogue such as
    "All: Thanks be to God." that must be recolored (otherwise Paper keeps white text).
    """
    for shape in _iter_slide_text_shapes(slide):
        if shape.top is not None and int(shape.top) >= _CLONE_FOOTER_ZONE_TOP:
            continue
        for para in shape.text_frame.paragraphs:
            if para.runs:
                for run in para.runs:
                    color = _theme_color_for_font(run.font, theme)
                    _apply_run_theme_color(run, color)
                # Paragraph default must not stay template-white for empty/partial runs.
                _set_def_rpr_rgb(para, theme.primary)
            else:
                color = _theme_color_for_font(para.font, theme)
                _apply_run_theme_color(para, color)
                _set_def_rpr_rgb(para, color)


def _ensure_slide_text_contrast(slide, theme: SlideTheme) -> None:
    """Force low-contrast / scheme / inherited text onto theme roles.

    Also rewrites paragraph ``defRPr`` fills. Master templates leave off-white
    defaults there; runs without their own fill inherit them and disappear on Paper.
    """
    bg_lum = _rel_lum(int(theme.bg[0]), int(theme.bg[1]), int(theme.bg[2]))
    light_bg = bg_lum > 0.55

    def _illegible(rgb: Optional[RGBColor]) -> bool:
        if rgb is None:
            return True
        lum = _rel_lum(int(rgb[0]), int(rgb[1]), int(rgb[2]))
        if light_bg and lum > 0.55:
            return True
        if (not light_bg) and lum < 0.28:
            return True
        return False

    for shape in _iter_slide_text_shapes(slide):
        if shape.top is not None and int(shape.top) >= _CLONE_FOOTER_ZONE_TOP:
            continue
        for para in shape.text_frame.paragraphs:
            if para.runs:
                for run in para.runs:
                    rgb = _font_rgb_or_none(run.font)
                    if theme.mono_surfaces:
                        # Strict mono: every run is black or white (titles use emphasis,
                        # which equals primary for Themes 2–3).
                        color = theme.emphasis if run.font.underline else theme.primary
                        _apply_run_theme_color(run, color)
                    elif _illegible(rgb):
                        color = _theme_color_for_font(run.font, theme)
                        _apply_run_theme_color(run, color)
                    elif rgb is not None:
                        _set_run_rpr_rgb(run, rgb)
                _set_def_rpr_rgb(para, theme.primary)
            else:
                rgb = _font_rgb_or_none(para.font)
                if theme.mono_surfaces or _illegible(rgb):
                    color = theme.emphasis if para.font.underline else theme.primary
                    _apply_run_theme_color(para, color)
                    _set_def_rpr_rgb(para, color)
                elif rgb is not None:
                    _set_def_rpr_rgb(para, rgb)


def _finalize_deck_text_colors(prs: Presentation, theme: SlideTheme) -> None:
    """Last pass over every slide so Paper/Midnight never keep template whites."""
    for slide in prs.slides:
        _ensure_slide_text_contrast(slide, theme)


def _bake_slide_text_colors(slide, theme: SlideTheme) -> None:
    """Bake every text run on ``slide`` to ``theme`` (no footer-zone skip).

    Used offline by ``bake_master_theme_template`` so Theme 2/3 masters ship with
    correct colors already in the PPTX XML.
    """
    for shape in _iter_slide_text_shapes(slide):
        for para in shape.text_frame.paragraphs:
            if para.runs:
                for run in para.runs:
                    color = _theme_color_for_font(run.font, theme)
                    if theme.mono_surfaces:
                        color = theme.emphasis if run.font.underline else theme.primary
                    _apply_run_theme_color(run, color)
                _set_def_rpr_rgb(para, theme.primary)
            else:
                color = _theme_color_for_font(para.font, theme)
                if theme.mono_surfaces:
                    color = theme.emphasis if para.font.underline else theme.primary
                _apply_run_theme_color(para, color)
                _set_def_rpr_rgb(para, color)


def bake_master_theme_template(theme_id: str) -> Path:
    """Create a baked master PPTX for ``theme2`` / ``theme3`` from Theme 1.

    Writes ``data/reference/LFTemplate2-midnight.pptx`` or
    ``LFTemplate3-paper.pptx``. Clears the in-memory cache for that file.
    """
    tid = (theme_id or "").strip().lower()
    if tid not in ("theme2", "theme3"):
        raise ValueError(f"Can only bake theme2/theme3, got {theme_id!r}")
    src_name = _MASTER_TEMPLATE_BY_THEME["theme1"]
    dest_name = _MASTER_TEMPLATE_BY_THEME[tid]
    src = _PROJECT_ROOT / "data" / "reference" / src_name
    dest = _PROJECT_ROOT / "data" / "reference" / dest_name
    if not src.is_file():
        raise FileNotFoundError(f"Theme 1 master missing: {src}")
    theme = _mono_slide_theme(dark=(tid == "theme2"))
    prs = Presentation(str(src.resolve()))
    for slide in prs.slides:
        _set_slide_bg(slide, theme.bg)
        _bake_slide_text_colors(slide, theme)
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    _master_templates.pop(dest_name, None)
    return dest.resolve()


def _style_shape_font(
    shape,
    *,
    font_name: str,
    size_pt: float,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    color: Optional[RGBColor] = None,
) -> None:
    """Apply font (and optional color) to every run (or paragraph) on a text shape."""
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        if para.runs:
            for run in para.runs:
                run.font.name = font_name
                run.font.size = Pt(size_pt)
                if bold is not None:
                    run.font.bold = bold
                if italic is not None:
                    run.font.italic = italic
                if color is not None:
                    run.font.color.rgb = color
        else:
            para.font.name = font_name
            para.font.size = Pt(size_pt)
            if bold is not None:
                para.font.bold = bold
            if italic is not None:
                para.font.italic = italic
            if color is not None:
                para.font.color.rgb = color


def _normalize_rite_title(text: str) -> str:
    t = (text or "").strip().lower()
    return t.replace("\u2019", "'").replace("\u2018", "'").replace("'", "'")


def _is_rite_slide_title_text(text: str, title: str) -> bool:
    t = (text or "").strip()
    return "\n" not in t and _normalize_rite_title(t) == _normalize_rite_title(title or "")


def _is_apostles_creed_title_text(text: str) -> bool:
    n = _normalize_rite_title(text)
    return n == _normalize_rite_title(_APOSTLES_CREED_TITLE) or n == "apostles creed"


def _is_nicene_creed_title_text(text: str) -> bool:
    n = _normalize_rite_title(text)
    return n == _normalize_rite_title(_NICENE_CREED_TITLE) or n == "nicene creed"


def _style_section_title_shape(shape, size_pt: float) -> None:
    """Pin a title shape to the top and style it Georgia bold underline, accent color."""
    try:
        shape.left = 0
        shape.top = Inches(_SECTION_TITLE_TOP_IN)
        shape.width = SLIDE_WIDTH
        shape.height = Inches(_SECTION_TITLE_H_IN)
    except (AttributeError, TypeError, ValueError):
        pass
    try:
        shape.text_frame.word_wrap = True
        shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    except (AttributeError, TypeError, ValueError):
        pass
    for para in shape.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        targets = list(para.runs) if para.runs else [para]
        for target in targets:
            target.font.name = _SECTION_TITLE_FONT
            target.font.size = Pt(size_pt)
            target.font.bold = True
            target.font.italic = False
            target.font.underline = True
            target.font.color.rgb = _ACTIVE_THEME.emphasis


def _add_section_title(slide, text: str, size_pt: float) -> None:
    """Add a top-of-slide Georgia bold underline section title (template style)."""
    box = slide.shapes.add_textbox(
        0, Inches(_SECTION_TITLE_TOP_IN), SLIDE_WIDTH, Inches(_SECTION_TITLE_H_IN)
    )
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    _style_para(
        p,
        size_pt=size_pt,
        color=_ACTIVE_THEME.emphasis,
        bold=True,
        font_name=_SECTION_TITLE_FONT,
        underline=True,
    )


def _apply_rite_slide_title_typography(
    slide, section_title: str, *, size_pt: float = _SECTION_TITLE_PT_SMALL
) -> None:
    """Section header: Georgia bold underline, accent color, pinned to the slide top."""
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_rite_slide_title_text(text, section_title):
            _style_section_title_shape(shape, size_pt)


def _is_lamb_of_god_header_text(text: str) -> bool:
    return _is_rite_slide_title_text(text, "Lamb of God")


def _is_lamb_of_god_lyric_text(text: str) -> bool:
    """Main prayer blocks start with ALL CAPS ``LAMB OF GOD`` (not the footer tag line)."""
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    return len(lines) >= 2 and lines[0].upper() == "LAMB OF GOD"


def _apply_lamb_of_god_typography(slide) -> None:
    """Match cloned Lamb slide to hymn title/body typography (Georgia 38.5 / Poppins Bold 56)."""
    _apply_rite_slide_title_typography(slide, "Lamb of God")
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_lamb_of_god_header_text(text) or _is_lamb_of_god_lyric_text(text):
            if _is_lamb_of_god_lyric_text(text):
                _style_shape_font(
                    shape, font_name=_ACTIVE_FONT, size_pt=_HYMN_BODY_PT, bold=True
                )


def _add_community_footer(slide, footer_section: str, theme: SlideTheme):
    if not _deck_branding.include_footer:
        return
    lx = MARGIN_SIDE
    w = SLIDE_WIDTH - 2 * MARGIN_SIDE
    y = SLIDE_HEIGHT - FOOTER_TOP_OFFSET
    foot = slide.shapes.add_textbox(lx, y, w, FOOTER_HEIGHT)
    tf = foot.text_frame
    _prep_tf(tf)
    tf.clear()
    brand = _footer_brand_line()
    if brand:
        p0 = tf.paragraphs[0]
        p0.text = brand
        _style_para(p0, size_pt=_FOOTER_PT, color=theme.muted, bold=True)
        p0.alignment = PP_ALIGN.LEFT
    p1 = tf.add_paragraph()
    p1.text = footer_section
    _style_para(p1, size_pt=_FOOTER_PT - 1, color=theme.emphasis, bold=False)
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(2)


def _parse_marked_lines(marked: str) -> List[Tuple[str, str]]:
    out: List[Tuple[str, str]] = []
    for raw in (marked or "").split("\n"):
        raw = raw.strip()
        if not raw:
            continue
        role = "plain"
        if raw.startswith("<<P>>"):
            role, raw = "priest", raw[5:].strip()
        elif raw.startswith("<<A>>"):
            role, raw = "all", raw[5:].strip()
        elif raw.startswith("<<D>>"):
            role, raw = "direction", raw[5:].strip()
        elif raw.startswith("<<H>>"):
            role, raw = "hymn", raw[5:].strip()
        out.append((role, raw))
    return out


def _strip_marked_rubrics(marked: str) -> str:
    """Remove ``<<D>>`` liturgical rubrics before building projection slides."""
    kept: List[str] = []
    for raw in (marked or "").split("\n"):
        if raw.strip().startswith("<<D>>"):
            continue
        kept.append(raw.rstrip())
    return "\n".join(kept).strip()


def _marked_has_projectable_content(marked: str) -> bool:
    for role, line in _parse_marked_lines(_strip_marked_rubrics(marked)):
        if role != "direction" and (line or "").strip():
            return True
    return False


def _suppress_all_role_prefix(footer_section: str) -> bool:
    """Kyrie / Gloria / Sanctus / Our Father / Lamb: congregation lines without an ``All:`` prefix."""
    f = (footer_section or "").strip().lower()
    keys = ("kyrie eleison", "gloria", "sanctus", "our father", "lamb of god")
    return any(f.startswith(k) for k in keys)


def _is_projection_dialogue_slide(footer_section: str) -> bool:
    """Priest/assembly dialogue: centered body, gold ``Priest:`` label, white text."""
    f = (footer_section or "").strip().lower()
    return f.startswith("final blessing") or f.startswith("the communion rite") or f == "communion rite"


def _is_prayer_rite_slide(footer_section: str) -> bool:
    """Mass prayer rites rendered large for projection (Kyrie, Gloria, Creed, etc.)."""
    f = (footer_section or "").strip().lower()
    keys = (
        "kyrie",
        "gloria",
        "sanctus",
        "our father",
        "lamb of god",
        "penitential act",
        "nicene creed",
        "apostles' creed",
        "apostles creed",
    )
    return any(k in f for k in keys)


def _marked_body_height_inches() -> float:
    """Usable vertical space for marked prayer/body text (inches)."""
    top = _length_to_inches(_content_top())
    return float(SLIDE_HEIGHT.inches) - top - 1.1


def _rite_display_line(role: str, line: str, *, strip_all: bool) -> str:
    if role == "priest":
        return f"Priest: {line}"
    if role == "all":
        return line if strip_all else f"All: {line}"
    return line


def _rite_wrapped_line_units(role: str, line: str, *, strip_all: bool) -> float:
    """Estimate how many slide lines one marked line consumes at rite font size."""
    text = _rite_display_line(role, line, strip_all=strip_all)
    chars_per_line = 24 if role in ("priest", "all", "plain") else 30
    if role == "direction":
        chars_per_line = 34
    wrapped = max(1, math.ceil(len(text) / chars_per_line))
    extra = 0.25 if role in ("priest", "all") else 0.12
    return wrapped + extra


def _rite_line_height_inches(font_pt: float = _SLIDE_TEXT_PT) -> float:
    return (font_pt * 1.22 + 8) / 72.0


def _serialize_marked_lines(items: List[Tuple[str, str]]) -> str:
    out: List[str] = []
    for role, line in items:
        if role == "priest":
            out.append(f"<<P>>{line}")
        elif role == "all":
            out.append(f"<<A>>{line}")
        elif role == "direction":
            out.append(f"<<D>>{line}")
        elif role == "hymn":
            out.append(f"<<H>>{line}")
        else:
            out.append(line)
    return "\n".join(out)


def _chunk_marked_rite_by_fit(
    marked: str,
    footer_section: str,
    body_h_inches: Optional[float] = None,
) -> List[str]:
    """Split prayer-rite marked text across slides when body text would overflow."""
    items = [(r, ln) for r, ln in _parse_marked_lines(marked) if r != "direction"]
    if not items:
        return [marked]

    strip_all = _suppress_all_role_prefix(footer_section)
    body_h = float(body_h_inches or _marked_body_height_inches())
    capacity = max(3.0, body_h / _rite_line_height_inches(_SLIDE_TEXT_PT))

    grouped: List[List[Tuple[str, str]]] = []
    current: List[Tuple[str, str]] = []
    used = 0.0

    for role, line in items:
        units = _rite_wrapped_line_units(role, line, strip_all=strip_all)
        if current and used + units > capacity:
            grouped.append(current)
            current = []
            used = 0.0
        current.append((role, line))
        used += units

    if current:
        grouped.append(current)

    if not grouped:
        return [marked]
    return [_serialize_marked_lines(part) for part in grouped]


def _render_templated_prayer_slide(
    prs: Presentation,
    template: Mapping[str, Any],
    slide_spec: Mapping[str, Any],
    theme: SlideTheme,
    *,
    footer_section: str,
    slide_index: int = 0,
    slide_total: int = 1,
) -> None:
    """Render one prayer-template slide (fixed line breaks and rite typography)."""
    footer = template.get("footer") or footer_section
    if slide_total > 1:
        footer = f"{footer} ({slide_index + 1}/{slide_total})"

    body_pt = _SLIDE_TEXT_PT

    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, theme.bg)
    _apply_slide_branding(slide, theme)
    lx, top, w = MARGIN_SIDE, _content_top(), SLIDE_WIDTH - 2 * MARGIN_SIDE
    body_h = SLIDE_HEIGHT - top - CONTENT_BOTTOM_GAP

    box = slide.shapes.add_textbox(lx, top, w, body_h)
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    tf.vertical_anchor = BODY_VERTICAL_ANCHOR

    first = True
    for line_spec in slide_spec.get("lines") or []:
        style = str(line_spec.get("style") or "plain").strip().lower()
        raw = str(line_spec.get("text") or "").strip()
        if not raw:
            continue
        # Skip sidenote-style directions (kept italic in the template) from projection.
        if style == "direction":
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        if style == "priest":
            p.text = f"Priest: {raw}"
            _style_para(p, size_pt=body_pt, color=theme.emphasis, bold=True)
            p.space_before = Pt(6)
        elif style == "all_lead":
            p.text = f"ALL: {raw}"
            _style_para(p, size_pt=body_pt, color=theme.primary, bold=True)
            p.space_before = Pt(6)
        elif style == "all_body":
            p.text = raw
            _style_para(p, size_pt=body_pt, color=theme.primary, bold=False)
            p.space_before = Pt(4)
        else:
            p.text = raw
            _style_para(p, size_pt=body_pt, color=theme.primary, bold=False)

        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    _add_community_footer(slide, footer, theme)


def _add_templated_prayer(prs: Presentation, template: Mapping[str, Any], theme: SlideTheme) -> None:
    slides = list(template.get("slides") or [])
    total = len(slides)
    footer_base = str(template.get("footer") or "Prayer")
    for i, slide_spec in enumerate(slides):
        _render_templated_prayer_slide(
            prs,
            template,
            slide_spec,
            theme,
            footer_section=footer_base,
            slide_index=i,
            slide_total=total,
        )


def _reference_mass_deck_path() -> Optional[Path]:
    candidates = (
        _PROJECT_ROOT / "data" / "reference" / _REFERENCE_DECK_FILENAME,
        _PROJECT_ROOT / "outputs" / "GFCC24May2026_Eastertide.pptx",
        Path.home() / "Downloads" / _REFERENCE_DECK_FILENAME,
    )
    for path in candidates:
        if path.is_file():
            return path.resolve()
    return None


def _load_reference_mass_deck() -> Optional[Presentation]:
    global _reference_mass_deck
    if _reference_mass_deck is not None:
        return _reference_mass_deck
    ref_path = _reference_mass_deck_path()
    if not ref_path:
        return None
    _reference_mass_deck = Presentation(str(ref_path))
    return _reference_mass_deck


def _master_template_filename(theme_id: Optional[str] = None) -> str:
    """Return the on-disk master PPTX filename for ``theme_id`` (falls back to Theme 1)."""
    tid = (theme_id or _ACTIVE_DECK_THEME_ID or "theme1").strip().lower()
    if tid not in _MASTER_TEMPLATE_BY_THEME:
        tid = "theme1"
    name = _MASTER_TEMPLATE_BY_THEME[tid]
    path = _PROJECT_ROOT / "data" / "reference" / name
    if tid != "theme1" and not path.is_file():
        return _MASTER_TEMPLATE_BY_THEME["theme1"]
    return name


def _master_template_is_baked(theme_id: Optional[str] = None) -> bool:
    """True when the active theme has its own baked master (skip live recolor on clone)."""
    tid = (theme_id or _ACTIVE_DECK_THEME_ID or "theme1").strip().lower()
    if tid not in ("theme2", "theme3"):
        return False
    name = _MASTER_TEMPLATE_BY_THEME.get(tid)
    if not name:
        return False
    return (_PROJECT_ROOT / "data" / "reference" / name).is_file()


def _master_template_path(theme_id: Optional[str] = None) -> Optional[Path]:
    path = _PROJECT_ROOT / "data" / "reference" / _master_template_filename(theme_id)
    return path.resolve() if path.is_file() else None


def _load_master_template(theme_id: Optional[str] = None) -> Optional[Presentation]:
    """Load (and cache) the master deck for the given / active theme id."""
    global _master_templates
    tid = (theme_id or _ACTIVE_DECK_THEME_ID or "theme1").strip().lower()
    if tid not in _MASTER_TEMPLATE_BY_THEME:
        tid = "theme1"
    # If the baked file is missing we actually load Theme 1 — cache under the
    # resolved filename so we don't pretend Theme 2 is baked.
    fname = _master_template_filename(tid)
    cache_key = fname
    cached = _master_templates.get(cache_key)
    if cached is not None:
        return cached
    ref_path = _PROJECT_ROOT / "data" / "reference" / fname
    if not ref_path.is_file():
        return None
    tpl = Presentation(str(ref_path.resolve()))
    _master_templates[cache_key] = tpl
    return tpl


def _master_slide_src(key: str, part: int = 0):
    """Resolve a master-template source slide for a section key (see ``_MASTER_SLIDE``)."""
    tpl = _load_master_template()
    if tpl is None:
        return None
    spec = _MASTER_SLIDE.get(key)
    if spec is None:
        return None
    idx = spec[part] if isinstance(spec, tuple) else spec
    if idx < 0 or idx >= len(tpl.slides):
        return None
    return tpl.slides[idx]


def _clone_master_section(
    prs: Presentation,
    key: str,
    theme: SlideTheme,
    footer_section: str,
    *,
    mutate=None,
) -> bool:
    """Clone every master slide for ``key`` into ``prs``.

    Baked Theme 2/3 masters are cloned without a live color pass. Theme 1 (and
    fallback) still recolors. ``mutate(slide, part_index)`` runs after each clone
    for dynamic text injection. Returns ``False`` if the master is unavailable.
    """
    spec = _MASTER_SLIDE.get(key)
    if spec is None or _load_master_template() is None:
        return False
    baked = _master_template_is_baked()
    parts = spec if isinstance(spec, tuple) else (spec,)
    total = len(parts)
    for part_i in range(total):
        src = _master_slide_src(key, part_i)
        if src is None:
            return False
        foot = footer_section if total == 1 else f"{footer_section} ({part_i + 1}/{total})"
        _copy_slide_into_presentation(
            prs,
            src,
            theme,
            foot,
            strip_italic_rubrics=False,
            recolor=not baked,
        )
        if mutate is not None:
            mutate(prs.slides[-1], part_i)
            # Injected runs keep the baked template colors; only re-check when we
            # still rely on live remapping (Theme 1 / missing bake).
            if not baked:
                _ensure_slide_text_contrast(prs.slides[-1], theme)
    return True


def _set_run_text_keep_format(slide, old: str, new: str) -> bool:
    """Replace the first run whose text equals ``old`` with ``new`` (keeps font/color)."""
    want = (old or "").strip()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip() == want:
                    run.text = new
                    return True
    return False


def _color_shapes(
    slide,
    color,
    *,
    exact: Optional[List[str]] = None,
    contains: Optional[str] = None,
    top_below: Optional[int] = None,
) -> None:
    """Recolor every run of shapes whose text matches ``exact``/``contains``.

    ``top_below`` (EMU) restricts the match to shapes above that vertical offset,
    used to target the header title without touching the footer tag.
    """
    exact_set = {(t or "").strip().lower() for t in (exact or []) if (t or "").strip()}
    sub = (contains or "").strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if top_below is not None and (shape.top is None or int(shape.top) >= int(top_below)):
            continue
        full = (shape.text_frame.text or "").strip().lower()
        if not full:
            continue
        if (exact_set and full in exact_set) or (sub and sub in full):
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    for run in para.runs:
                        _apply_run_theme_color(run, color)
                    _set_def_rpr_rgb(para, color)
                else:
                    _apply_run_theme_color(para, color)
                    _set_def_rpr_rgb(para, color)


def _amber_header_title(slide, band_in: float = 0.5) -> None:
    """Force a cloned slide's header title (the top ``band_in`` inches) into amber.

    Used for cloned decks whose title run is neither underlined nor saturated, so
    the standard recolor leaves it white (e.g. the Our Father language decks).
    Body/dialogue text sits well below the band and is left untouched.
    """
    limit = int(Inches(band_in))
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if shape.top is None or int(shape.top) >= limit:
            continue
        if not (shape.text_frame.text or "").strip():
            continue
        for para in shape.text_frame.paragraphs:
            if para.runs:
                for run in para.runs:
                    _apply_run_theme_color(run, _ACTIVE_THEME.emphasis)
                _set_def_rpr_rgb(para, _ACTIVE_THEME.emphasis)
            else:
                _apply_run_theme_color(para, _ACTIVE_THEME.emphasis)
                _set_def_rpr_rgb(para, _ACTIVE_THEME.emphasis)


def _extract_psalm_response(psalm_text: str) -> str:
    """First refrain sentence following the ``R``/``(R)``/``R. (14c)`` marker.

    Falls back to the first sentence of the body when no responsory marker exists.
    Returned uppercase to mirror the template antiphon styling.
    """
    text = (psalm_text or "").strip()
    if not text:
        return ""
    # Responsory marker: "(R)", "R." / "R/" / "R:", or "R(14c)" / "R. (cf. 11)".
    marker = re.compile(
        r"(?:\(R\)|R\s*\([^)]*\)|R[.:/])\s*(?:\([^)]*\)\s*)?",
        re.IGNORECASE,
    )
    m = marker.search(text)
    body = text[m.end():] if m else text
    sent = re.search(r"(.+?[.!?])(?:\s|$)", body, re.S)
    response = (sent.group(1) if sent else body).strip()
    response = re.sub(r"\s+", " ", response)
    return response.upper()


def _add_our_father_from_deck(prs: Presentation, theme: SlideTheme, choice: str) -> bool:
    """Clone the bundled per-language Our Father slide set (recolored). False if none."""
    fname = _OUR_FATHER_DECKS.get(choice)
    if not fname:
        return False
    path = _PROJECT_ROOT / "data" / "reference" / fname
    if not path.is_file():
        return False
    try:
        tpl = Presentation(str(path))
    except Exception:
        return False
    if not tpl.slides:
        return False
    for src in tpl.slides:
        _copy_slide_into_presentation(prs, src, theme, "Our Father", strip_italic_rubrics=False)
        _amber_header_title(prs.slides[-1])
    return True


def _add_food_sponsors_from_template(
    prs: Presentation, theme: SlideTheme, names: List[str]
) -> bool:
    """Clone the authored Food Sponsors card (recolored) and inject sponsor names.

    All names are stacked into the primary name box; the spare placeholder box is
    cleared. Returns ``False`` if the donor deck is unavailable.
    """
    path = _PROJECT_ROOT / "data" / "reference" / _FOOD_SPONSORS_TEMPLATE_FILENAME
    if not path.is_file():
        return False
    try:
        tpl = Presentation(str(path))
    except Exception:
        return False
    if not tpl.slides:
        return False
    _copy_slide_into_presentation(
        prs, tpl.slides[0], theme, "Food Sponsors", strip_italic_rubrics=False
    )
    slide = prs.slides[-1]
    _set_run_text_keep_format(slide, _TPL_SPONSOR_PLACEHOLDER, "")
    _fill_lines_keep_format(slide, _TPL_SPONSOR_EXAMPLE, names)
    _color_shapes(slide, _ACTIVE_THEME.emphasis, exact=["food sponsors"], top_below=Inches(2))
    return True


def _fill_lines_keep_format(slide, anchor: str, lines: List[str]) -> bool:
    """Replace the text box containing ``anchor`` with one paragraph per line,
    reusing the anchor run's font/size/style/color and paragraph alignment."""
    want = (anchor or "").strip()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if (tf.text or "").strip() != want:
            continue
        src = None
        for para in tf.paragraphs:
            if para.runs:
                src = para.runs[0]
                break
        font_name = src.font.name if src else None
        font_size = src.font.size if src else None
        font_bold = src.font.bold if src else None
        font_italic = src.font.italic if src else None
        try:
            font_color = src.font.color.rgb if src and src.font.color and src.font.color.type is not None else None
        except Exception:
            font_color = None
        align = tf.paragraphs[0].alignment if tf.paragraphs else None
        tf.clear()
        first = True
        for line in (lines or [""]):
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if align is not None:
                para.alignment = align
            run = para.add_run()
            run.text = line
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = font_size
            run.font.bold = font_bold
            run.font.italic = font_italic
            if font_color is not None:
                run.font.color.rgb = font_color
        return True
    return False


def _lamb_of_god_template_path() -> Optional[Path]:
    candidates = (
        _PROJECT_ROOT / "data" / "reference" / _LAMB_OF_GOD_TEMPLATE_FILENAME,
        Path.home() / "Downloads" / "GFCC_26APRIL2026 (3).pptx",
    )
    for path in candidates:
        if path.is_file():
            return path.resolve()
    return None


def _load_lamb_of_god_template() -> Optional[Presentation]:
    global _lamb_of_god_template
    if _lamb_of_god_template is not None:
        return _lamb_of_god_template
    ref_path = _lamb_of_god_template_path()
    if not ref_path:
        return None
    _lamb_of_god_template = Presentation(str(ref_path))
    return _lamb_of_god_template


def _sign_of_peace_template_path() -> Optional[Path]:
    candidates = (
        _PROJECT_ROOT / "data" / "reference" / _SIGN_OF_PEACE_TEMPLATE_FILENAME,
        Path.home() / "Downloads" / "GFCC_26APRIL2026 (4).pptx",
    )
    for path in candidates:
        if path.is_file():
            return path.resolve()
    return None


def _load_sign_of_peace_template() -> Optional[Presentation]:
    global _sign_of_peace_template
    if _sign_of_peace_template is not None:
        return _sign_of_peace_template
    ref_path = _sign_of_peace_template_path()
    if not ref_path:
        return None
    _sign_of_peace_template = Presentation(str(ref_path))
    return _sign_of_peace_template


def _gloria_template_path() -> Optional[Path]:
    candidates = (
        _PROJECT_ROOT / "data" / "reference" / _GLORIA_TEMPLATE_FILENAME,
        Path.home() / "Downloads" / "Gloria.pptx",
    )
    for path in candidates:
        if path.is_file():
            return path.resolve()
    return None


def _load_gloria_template() -> Optional[Presentation]:
    global _gloria_template
    if _gloria_template is not None:
        return _gloria_template
    ref_path = _gloria_template_path()
    if not ref_path:
        return None
    _gloria_template = Presentation(str(ref_path))
    return _gloria_template


def _kyrie_template_path() -> Optional[Path]:
    candidates = (
        _PROJECT_ROOT / "data" / "reference" / _KYRIE_TEMPLATE_FILENAME,
        Path.home() / "Downloads" / "GFCC_26APRIL2026 (6).pptx",
    )
    for path in candidates:
        if path.is_file():
            return path.resolve()
    return None


def _load_kyrie_template() -> Optional[Presentation]:
    global _kyrie_template
    if _kyrie_template is not None:
        return _kyrie_template
    ref_path = _kyrie_template_path()
    if not ref_path:
        return None
    _kyrie_template = Presentation(str(ref_path))
    return _kyrie_template


def _is_kyrie_body_text(text: str) -> bool:
    """Main Kyrie blocks: LORD / CHRIST HAVE MERCY (all caps, not the section title)."""
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    if not lines:
        return False
    joined = " ".join(lines).upper()
    if _is_rite_slide_title_text(text, "Kyrie Eleison"):
        return False
    return "HAVE MERCY" in joined or joined.startswith("LORD,") or joined.startswith("CHRIST,")


def _apply_kyrie_typography(slide) -> None:
    """Kyrie title Georgia 38.5 pt (same as Gloria); body Poppins Bold 56 pt ALL CAPS blocks."""
    _apply_rite_slide_title_typography(slide, "Kyrie Eleison")
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_rite_slide_title_text(text, "Kyrie Eleison"):
            continue
        if _is_kyrie_body_text(text):
            _style_shape_font(
                shape, font_name=_ACTIVE_FONT, size_pt=_HYMN_BODY_PT, bold=True
            )


def _gloria_source_slide_indices(slide_count: int) -> Tuple[int, ...]:
    """
    Map Gloria reference deck to four projection slides.

    The bundled deck has five slides; index 1 is a duplicate refrain-only slide
    between the opening and the Christ section, so we use 0, 2, 3, 4.
    """
    if slide_count >= 5:
        return (0, 2, 3, 4)
    if slide_count == 4:
        return (0, 1, 2, 3)
    return tuple(range(max(0, slide_count)))


def _is_reference_branding_shape(shape) -> bool:
    """Skip reference logo groups and baked-in parish footer text boxes."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return True
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return False
    if int(shape.top) >= _CLONE_FOOTER_ZONE_TOP:
        return True
    text = (shape.text_frame.text or "").strip().lower()
    parish = get_community_name().strip().lower()
    if parish and parish in text and int(shape.top) < _CLONE_FOOTER_ZONE_TOP:
        return True
    return False


def _strip_italic_rubric_paragraphs_on_slide(slide) -> None:
    """Remove italic rubric paragraphs from cloned reference-deck shapes."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = (para.text or "").strip()
            if not text:
                continue
            try:
                if para.font.italic:
                    para.text = ""
            except (AttributeError, TypeError):
                continue


def _copy_slide_into_presentation(
    prs: Presentation,
    slide_src,
    theme: SlideTheme,
    footer_section: str,
    *,
    copy_groups: bool = False,
    strip_italic_rubrics: bool = True,
    recolor: bool = True,
) -> None:
    """
    Clone prayer/body shapes from a reference slide.

    Uses liturgical ``theme`` background plus current logo and parish name;
    does not copy reference background, logo group, or footer text.

    When ``recolor`` is False (baked Theme 2/3 masters), text colors from the
    source template are kept as-authored — only bg + branding/footer are applied.
    """
    layout = _layout_blank(prs)
    dest = prs.slides.add_slide(layout)
    for shp in list(dest.shapes):
        el = shp.element
        el.getparent().remove(el)
    parish = get_community_name().strip().lower()
    for shp in slide_src.shapes:
        if copy_groups:
            if int(shp.top) >= _CLONE_FOOTER_ZONE_TOP:
                continue
            if (
                getattr(shp, "has_text_frame", False)
                and shp.has_text_frame
                and parish
            ):
                text = (shp.text_frame.text or "").strip().lower()
                if parish in text:
                    continue
        elif _is_reference_branding_shape(shp):
            continue
        newel = deepcopy(shp.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")
    _set_slide_bg(dest, theme.bg)
    _apply_slide_branding(dest, theme)
    if recolor:
        _recolor_cloned_text_to_theme(dest, theme)
    if strip_italic_rubrics:
        _strip_italic_rubric_paragraphs_on_slide(dest)
    _add_community_footer(dest, footer_section, theme)
    if recolor:
        _ensure_slide_text_contrast(dest, theme)


def _lotw_title_image_path() -> Optional[Path]:
    path = _PROJECT_ROOT / "data" / "reference" / _LOTW_TITLE_IMAGE_FILENAME
    if path.is_file():
        return path.resolve()
    return None


def _lotw_title_template_path() -> Optional[Path]:
    path = _PROJECT_ROOT / "data" / "reference" / _LOTW_TITLE_TEMPLATE_FILENAME
    if path.is_file():
        return path.resolve()
    return None


def _load_lotw_title_template() -> Optional[Presentation]:
    global _lotw_title_template
    if _lotw_title_template is not None:
        return _lotw_title_template
    ref_path = _lotw_title_template_path()
    if not ref_path:
        return None
    _lotw_title_template = Presentation(str(ref_path))
    return _lotw_title_template


def _resolve_bundled_poster(
    selection: Optional[str], valid_ids: Tuple[str, ...], default_id: str
) -> Optional[Path]:
    """Map a poster selection id (e.g. ``"lotw2"``) to a bundled PNG path.

    Unknown/empty selections fall back to the default design. Returns ``None`` when
    the file is missing so callers can use their non-poster fallback layout.
    """
    key = str(selection or "").strip().lower() or default_id
    if key not in valid_ids:
        key = default_id
    path = _POSTER_REFERENCE_DIR / f"{key}.png"
    return path if path.is_file() else None


def _add_full_bleed_poster_slide(
    prs: Presentation, poster_path: Path, footer_section: str, theme: SlideTheme
) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    slide.shapes.add_picture(
        str(poster_path),
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    _add_community_footer(slide, footer_section, theme)


def _add_lotw_title_slide(
    prs: Presentation, theme: SlideTheme, poster_path: Optional[Path] = None
) -> None:
    """Liturgy of the Word divider: full-bleed poster art (replaces plain section card)."""
    footer_section = "Liturgy of the Word"
    if poster_path is not None and Path(poster_path).is_file():
        _add_full_bleed_poster_slide(prs, Path(poster_path), footer_section, theme)
        return
    image_path = _lotw_title_image_path()
    if image_path is not None:
        _add_full_bleed_poster_slide(prs, image_path, footer_section, theme)
        return
    tpl = _load_lotw_title_template()
    if tpl is not None and _LOTW_TITLE_SLIDE_INDEX < len(tpl.slides):
        _copy_slide_into_presentation(
            prs,
            tpl.slides[_LOTW_TITLE_SLIDE_INDEX],
            theme,
            footer_section,
            copy_groups=True,
        )
        return
    _add_section_card(prs, "LITURGY OF\nTHE WORD", footer_section, theme)


def _add_lote_poster_slide(
    prs: Presentation, theme: SlideTheme, poster_path: Optional[Path] = None
) -> None:
    """Liturgy of the Eucharist divider: full-bleed poster art (replaces section card)."""
    footer_section = "Liturgy of the Eucharist"
    if poster_path is not None and Path(poster_path).is_file():
        _add_full_bleed_poster_slide(prs, Path(poster_path), footer_section, theme)
        return
    _add_section_card(prs, "LITURGY OF\nTHE EUCHARIST", footer_section, theme)


def _copy_reference_slides(
    prs: Presentation,
    slide_specs: Tuple[Tuple[int, str], ...],
    theme: SlideTheme,
) -> bool:
    ref = _load_reference_mass_deck()
    if ref is None:
        return False
    for idx, _footer in slide_specs:
        if idx < 0 or idx >= len(ref.slides):
            return False
    total = len(slide_specs)
    for part_i, (idx, footer_base) in enumerate(slide_specs):
        footer = footer_base if total == 1 else f"{footer_base} ({part_i + 1}/{total})"
        _copy_slide_into_presentation(prs, ref.slides[idx], theme, footer)
    return True


def _add_pre_mass_slide(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "pre_mass", theme, "Pre-Mass"):
        return
    if _copy_reference_slides(prs, ((_REFERENCE_SLIDE_PRE_MASS, "Pre-Mass"),), theme):
        return
    _add_marked_slide(prs, "Pre-Mass", GFCC.SILENT_REMINDER, theme)


def _add_penitential_act_slides(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "penitential", theme, "Penitential Act"):
        return
    specs = tuple((idx, "Penitential Act") for idx in _REFERENCE_SLIDE_PENITENTIAL)
    if _copy_reference_slides(prs, specs, theme):
        return
    _add_templated_prayer(prs, PENITENTIAL_ACT, theme)


def _add_kyrie_slide(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "kyrie", theme, "Kyrie Eleison"):
        return
    tpl = _load_kyrie_template()
    if tpl is not None and _KYRIE_SLIDE_INDEX < len(tpl.slides):
        _copy_slide_into_presentation(
            prs, tpl.slides[_KYRIE_SLIDE_INDEX], theme, "Kyrie Eleison"
        )
        _apply_kyrie_typography(prs.slides[-1])
        return
    if _copy_reference_slides(prs, ((_REFERENCE_SLIDE_KYRIE, "Kyrie Eleison"),), theme):
        _apply_kyrie_typography(prs.slides[-1])
        return
    _add_marked_slide(prs, "Kyrie Eleison", GFCC.KYRIE, theme)


def _add_lamb_of_god_slide(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "lamb_of_god", theme, "Lamb of God"):
        return
    tpl = _load_lamb_of_god_template()
    if tpl is not None and _LAMB_OF_GOD_SLIDE_INDEX < len(tpl.slides):
        _copy_slide_into_presentation(
            prs, tpl.slides[_LAMB_OF_GOD_SLIDE_INDEX], theme, "Lamb of God"
        )
        _apply_lamb_of_god_typography(prs.slides[-1])
        return
    _add_marked_slide(prs, "Lamb of God", get_prayer("lamb_of_god"), theme)


def _apply_sign_of_peace_typography(slide) -> None:
    _apply_rite_slide_title_typography(
        slide, "Sign of Peace", size_pt=_SECTION_TITLE_PT_LARGE
    )


def _add_sign_of_peace_slide(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "sign_of_peace", theme, "Sign of Peace"):
        return
    tpl = _load_sign_of_peace_template()
    if tpl is not None and _SIGN_OF_PEACE_SLIDE_INDEX < len(tpl.slides):
        _copy_slide_into_presentation(
            prs, tpl.slides[_SIGN_OF_PEACE_SLIDE_INDEX], theme, "Sign of Peace"
        )
        _apply_sign_of_peace_typography(prs.slides[-1])
        return
    _add_marked_slide(prs, "Sign of Peace", GFCC.SIGN_PEACE, theme)


def _apply_gloria_typography(slide) -> None:
    _apply_rite_slide_title_typography(slide, "Gloria")


def _add_gloria_slides(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "gloria", theme, "Gloria"):
        return
    tpl = _load_gloria_template()
    if tpl is None:
        _add_marked_chunked(prs, "Gloria", get_prayer("gloria"), theme)
        return
    indices = _gloria_source_slide_indices(len(tpl.slides))
    if not indices or any(i < 0 or i >= len(tpl.slides) for i in indices):
        _add_marked_chunked(prs, "Gloria", get_prayer("gloria"), theme)
        return
    total = len(indices)
    for part_i, idx in enumerate(indices):
        footer = "Gloria" if total == 1 else f"Gloria ({part_i + 1}/{total})"
        _copy_slide_into_presentation(prs, tpl.slides[idx], theme, footer)
        _apply_gloria_typography(prs.slides[-1])


def _gospel_acclamation_template_path() -> Optional[Path]:
    path = _PROJECT_ROOT / "data" / "reference" / _GOSPEL_ACCLAMATION_TEMPLATE_FILENAME
    if path.is_file():
        return path.resolve()
    return None


def _load_gospel_acclamation_template() -> Optional[Presentation]:
    global _gospel_acclamation_template
    if _gospel_acclamation_template is not None:
        return _gospel_acclamation_template
    ref_path = _gospel_acclamation_template_path()
    if not ref_path:
        return None
    _gospel_acclamation_template = Presentation(str(ref_path))
    return _gospel_acclamation_template


def _disable_cloned_slide_autofit(slide) -> None:
    """Stop PowerPoint from shrinking text to fit boxes (SHAPE_TO_FIT_TEXT on reference slides)."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE


def _gospel_acclamation_body_shape(slide):
    """Largest non-title text box on a Gospel Acclamation slide."""
    best = None
    best_len = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or _is_rite_slide_title_text(text, "Gospel Acclamation"):
            continue
        if len(text) > best_len:
            best_len = len(text)
            best = shape
    return best


def _format_gospel_acclamation_projection_text(verse: str) -> str:
    """Alleluia + lectionary verse + Alleluia (one projection block, uniform sizing)."""
    raw = (verse or "").strip()
    if not raw:
        return ""
    raw = re.sub(r"^R\.?\s*", "", raw, flags=re.I).strip()
    lines = [ln.strip() for ln in re.split(r"\n+", raw) if ln.strip()]
    core_lines = [
        ln
        for ln in lines
        if not re.match(r"^alleluia[,.!]?\s*$", ln, flags=re.I)
    ]
    core = " ".join(core_lines) if core_lines else raw
    core = re.sub(
        r"^Alleluia[,.!]?\s*(?=.)",
        "",
        core,
        count=1,
        flags=re.I,
    ).strip()
    if not core:
        return ""
    return (
        "Alleluia, alleluia, alleluia!\n"
        f"{core}\n"
        "Alleluia, alleluia, alleluia!"
    )


def _replace_shape_text_preserve_runs(shape, old: str, new: str) -> bool:
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in (run.text or ""):
                run.text = (run.text or "").replace(old, new, 1)
                replaced = True
    if not replaced and old in (shape.text_frame.text or ""):
        shape.text_frame.text = (shape.text_frame.text or "").replace(old, new, 1)
        replaced = True
    return replaced


def _gospel_book_from_reference(gospel_reference: str) -> str:
    """Evangelist or book name for “Gospel according to …” (from lectionary citation)."""
    ref = (gospel_reference or "").strip()
    if not ref or ref.upper() == "N/A":
        return "John"
    low = ref.lower()
    if "according to" in low:
        return ref.split("according to", 1)[-1].strip().rstrip(".")
    m = re.match(r"^(?:\d+\s+)?([A-Za-z][A-Za-z]+)", ref)
    if m:
        return m.group(1)
    parts = ref.split()
    return parts[0] if parts else "John"


def _gospel_acclamation_run_font(
    run,
    *,
    color: RGBColor,
    italic: bool = False,
) -> None:
    run.font.name = _ACTIVE_FONT
    run.font.size = Pt(_GOSPEL_ACCLAMATION_BODY_PT)
    run.font.bold = True
    run.font.italic = italic
    run.font.color.rgb = color


def _gospel_acclamation_role_line(
    tf,
    label: str,
    body: str,
    *,
    first: bool,
    label_color: RGBColor,
) -> None:
    if first:
        para = tf.paragraphs[0]
        para.text = ""
    else:
        para = tf.add_paragraph()
        para.text = ""
    label_run = para.add_run()
    label_run.text = label
    _gospel_acclamation_run_font(label_run, color=label_color)
    body_run = para.add_run()
    body_run.text = body
    _gospel_acclamation_run_font(body_run, color=_ACTIVE_THEME.primary)


def _gospel_acclamation_priest_line(tf, body: str, *, first: bool) -> None:
    _gospel_acclamation_role_line(
        tf,
        "Priest: ",
        body,
        first=first,
        label_color=_ACTIVE_THEME.emphasis,
    )


def _gospel_acclamation_all_line(tf, body: str, *, first: bool) -> None:
    _gospel_acclamation_role_line(
        tf,
        "All: ",
        body,
        first=first,
        label_color=_ACTIVE_THEME.primary,
    )


def _rebuild_gospel_acclamation_intro_shape(shape, gospel_book: str) -> None:
    """Priest label gold; All label and dialogue white; evangelist/book name in italic."""
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return
    raw = (shape.text_frame.text or "").strip()
    if not raw:
        return
    low = raw.lower()
    tf = shape.text_frame
    tf.clear()
    book = (gospel_book or "John").strip().rstrip(".")

    if "gospel according to" in low:
        _gospel_acclamation_priest_line(tf, "A reading from the holy", first=True)
        para = tf.add_paragraph()
        para.text = ""
        lead = para.add_run()
        lead.text = "Gospel according to "
        _gospel_acclamation_run_font(lead, color=_ACTIVE_THEME.primary)
        book_run = para.add_run()
        book_run.text = book
        _gospel_acclamation_run_font(
            book_run,
            color=_ACTIVE_THEME.primary,
            italic=True,
        )
        dot = para.add_run()
        dot.text = "."
        _gospel_acclamation_run_font(dot, color=_ACTIVE_THEME.primary)
        return

    if low.startswith("priest:"):
        body = raw.split(":", 1)[-1].strip()
        _gospel_acclamation_priest_line(tf, body, first=True)
        return

    if low.startswith("all:") or low.startswith("people:"):
        body = raw.split(":", 1)[-1].strip()
        _gospel_acclamation_all_line(tf, body, first=True)
        return


def _apply_gospel_acclamation_intro_typography(slide, gospel_reference: str) -> None:
    """Gospel Acclamation dialogue slide: gold Priest label, white All/dialogue, italic book."""
    _apply_rite_slide_title_typography(slide, "Gospel Acclamation", size_pt=_SECTION_TITLE_PT_LARGE)
    book = _gospel_book_from_reference(gospel_reference)
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_rite_slide_title_text(text, "Gospel Acclamation"):
            continue
        low = text.lower()
        if not (
            low.startswith("priest:")
            or low.startswith("all:")
            or low.startswith("people:")
            or "gospel according to" in low
        ):
            continue
        _rebuild_gospel_acclamation_intro_shape(shape, book)


def _patch_gospel_acclamation_alleluia_slide(slide, verse: str) -> None:
    body = _format_gospel_acclamation_projection_text(verse)
    if not body:
        return
    shape = _gospel_acclamation_body_shape(slide)
    if shape is None:
        return
    shape.text_frame.text = body


def _set_gospel_acclamation_alleluia_body(shape, verse: str) -> None:
    """Inject the week's acclamation verse while preserving the master body format.

    The master template's alleluia body (Arial 69 bold, centered, near-white) is
    mirrored exactly: we read the formatting off the existing first run and reuse
    it for every rebuilt line, instead of forcing a different deck font/alignment.
    """
    body = _format_gospel_acclamation_projection_text(verse)
    if not body:
        return
    tf = shape.text_frame
    ref_run = None
    for para in tf.paragraphs:
        if para.runs:
            ref_run = para.runs[0]
            break
    ref_align = tf.paragraphs[0].alignment if tf.paragraphs else None
    ref_name = ref_run.font.name if ref_run is not None else None
    ref_size = ref_run.font.size if ref_run is not None else None
    ref_bold = ref_run.font.bold if ref_run is not None else None
    ref_color = None
    if ref_run is not None:
        try:
            if ref_run.font.color is not None and ref_run.font.color.type is not None:
                ref_color = ref_run.font.color.rgb
        except (AttributeError, TypeError):
            ref_color = None

    lines = body.split("\n")
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ref_align is not None:
            para.alignment = ref_align
        run = para.add_run()
        run.text = line
        if ref_name:
            run.font.name = ref_name
        if ref_size is not None:
            run.font.size = ref_size
        if ref_bold is not None:
            run.font.bold = ref_bold
        if ref_color is not None:
            run.font.color.rgb = ref_color


def _apply_gospel_acclamation_typography(slide) -> None:
    """Title Georgia 50 pt; all body lines fixed 69 pt (no autofit shrink)."""
    _apply_rite_slide_title_typography(slide, "Gospel Acclamation", size_pt=_SECTION_TITLE_PT_LARGE)
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_rite_slide_title_text(text, "Gospel Acclamation"):
            continue
        _style_shape_font(
            shape,
            font_name=_ACTIVE_FONT,
            size_pt=_GOSPEL_ACCLAMATION_BODY_PT,
            bold=True,
            color=_ACTIVE_THEME.primary,
        )


def _gospel_acclamation_source_slide_indices(slide_count: int) -> Tuple[int, ...]:
    """Alleluia (one slide) + priest/assembly dialogue; skip duplicate alleluia in reference."""
    if slide_count >= 3:
        return (0, 2)
    if slide_count == 2:
        return (0, 1)
    if slide_count == 1:
        return (0,)
    return ()


def _add_gospel_acclamation_slides(
    prs: Presentation,
    theme: SlideTheme,
    *,
    gospel_reference: str = "",
    gospel_acclamation_verse: str = "",
) -> None:
    """Gospel Acclamation: alleluia + priest/assembly dialogue.

    Both slides are cloned 1:1 from the master template so they mirror the closing
    "Gospel of the Lord" slide exactly (same title alignment, fonts, and layout).
    The alleluia keeps the dynamic lectionary verse; the dialogue injects the
    evangelist/book name.
    """
    if _load_master_template() is not None:
        def _finish_alleluia(slide, _i):
            shape = _gospel_acclamation_body_shape(slide)
            if shape is not None:
                _set_gospel_acclamation_alleluia_body(shape, gospel_acclamation_verse or "")
            _apply_rite_slide_title_typography(
                slide, "Gospel Acclamation", size_pt=_SECTION_TITLE_PT_LARGE
            )

        def _finish_dialogue(slide, _i):
            ref = (gospel_reference or "").strip()
            if ref:
                _set_run_text_keep_format(slide, "Matthew", _gospel_book_from_reference(ref))
            _apply_rite_slide_title_typography(
                slide, "Gospel Acclamation", size_pt=_SECTION_TITLE_PT_LARGE
            )

        _clone_master_section(
            prs, "gospel_acclamation_alleluia", theme, "Gospel Acclamation (1/2)",
            mutate=_finish_alleluia,
        )
        _clone_master_section(
            prs, "gospel_acclamation_dialogue", theme, "Gospel Acclamation (2/2)",
            mutate=_finish_dialogue,
        )
        return

    # Fallback: legacy designed deck when the master template is unavailable.
    tpl = _load_gospel_acclamation_template()
    indices = _gospel_acclamation_source_slide_indices(
        len(tpl.slides) if tpl is not None else 0
    )
    if tpl is None or not indices:
        _add_marked_slide(prs, "Gospel Acclamation", GFCC.ALLELUIA_SING, theme)
        _add_marked_slide(prs, "Gospel Acclamation", GFCC.GOSPEL_INTRO, theme)
        return
    total = len(indices)
    for part_i, idx in enumerate(indices):
        footer = (
            "Gospel Acclamation"
            if total == 1
            else f"Gospel Acclamation ({part_i + 1}/{total})"
        )
        _copy_slide_into_presentation(prs, tpl.slides[idx], theme, footer)
        slide = prs.slides[-1]
        _disable_cloned_slide_autofit(slide)
        if idx == 0:
            _patch_gospel_acclamation_alleluia_slide(slide, gospel_acclamation_verse)
            _apply_gospel_acclamation_typography(slide)
        else:
            _apply_gospel_acclamation_intro_typography(slide, gospel_reference)


def _normalize_creed_choice(choice: str) -> str:
    c = (choice or "").strip().lower().replace("-", "_")
    if c in ("apostles", "apostles_creed", "apostle"):
        return "apostles"
    return "nicene"


def _normalize_our_father_choice(choice: str) -> str:
    c = (choice or "").strip().lower().replace("-", "_").replace(" ", "_")
    if c in ("cebuano", "bisaya"):
        return "visaya"
    if c in ("english", "malay", "tagalog", "visaya", "korean"):
        return c
    return "english"


def _apostles_creed_template_path() -> Optional[Path]:
    path = _PROJECT_ROOT / "data" / "reference" / _APOSTLES_CREED_TEMPLATE_FILENAME
    if path.is_file():
        return path.resolve()
    return None


def _load_apostles_creed_template() -> Optional[Presentation]:
    global _apostles_creed_template
    if _apostles_creed_template is not None:
        return _apostles_creed_template
    ref_path = _apostles_creed_template_path()
    if not ref_path:
        return None
    _apostles_creed_template = Presentation(str(ref_path))
    return _apostles_creed_template


def _apply_apostles_creed_typography(slide) -> None:
    """Section title Georgia 38.5 pt; leave body fonts from reference deck."""
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_apostles_creed_title_text(text):
            _style_section_title_shape(shape, _SECTION_TITLE_PT_SMALL)


def _add_apostles_creed_slides(prs: Presentation, theme: SlideTheme) -> None:
    tpl = _load_apostles_creed_template()
    if tpl is None or not tpl.slides:
        _add_marked_slide(
            prs,
            _APOSTLES_CREED_TITLE,
            "<<D>>Apostles' Creed reference slides not found. Add data/reference/apostles_creed_slides.pptx.",
            theme,
        )
        return
    total = len(tpl.slides)
    footer_base = _APOSTLES_CREED_TITLE
    for part_i in range(total):
        footer = footer_base if total == 1 else f"{footer_base} ({part_i + 1}/{total})"
        _copy_slide_into_presentation(prs, tpl.slides[part_i], theme, footer)
        slide = prs.slides[-1]
        _disable_cloned_slide_autofit(slide)
        _apply_apostles_creed_typography(slide)


def _nicene_creed_template_path() -> Optional[Path]:
    path = _PROJECT_ROOT / "data" / "reference" / _NICENE_CREED_TEMPLATE_FILENAME
    if path.is_file():
        return path.resolve()
    return None


def _load_nicene_creed_template() -> Optional[Presentation]:
    global _nicene_creed_template
    if _nicene_creed_template is not None:
        return _nicene_creed_template
    ref_path = _nicene_creed_template_path()
    if not ref_path:
        return None
    _nicene_creed_template = Presentation(str(ref_path))
    return _nicene_creed_template


def _apply_nicene_creed_typography(slide) -> None:
    """Section title Georgia 38.5 pt; leave body fonts from reference deck."""
    parish = get_community_name().strip().lower()
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        if int(shape.top) >= _REFERENCE_FOOTER_ZONE_TOP:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text or (parish and parish in text.lower()):
            continue
        if _is_nicene_creed_title_text(text):
            _style_section_title_shape(shape, _SECTION_TITLE_PT_SMALL)


def _add_nicene_creed_slides(prs: Presentation, theme: SlideTheme) -> None:
    if _clone_master_section(prs, "nicene_creed", theme, _NICENE_CREED_TITLE):
        return
    tpl = _load_nicene_creed_template()
    if tpl is None or not tpl.slides:
        _add_marked_chunked(prs, _NICENE_CREED_TITLE, get_prayer("nicene_creed"), theme)
        return
    total = len(tpl.slides)
    footer_base = _NICENE_CREED_TITLE
    for part_i in range(total):
        footer = footer_base if total == 1 else f"{footer_base} ({part_i + 1}/{total})"
        _copy_slide_into_presentation(
            prs,
            tpl.slides[part_i],
            theme,
            footer,
            strip_italic_rubrics=False,
        )
        slide = prs.slides[-1]
        _disable_cloned_slide_autofit(slide)
        _apply_nicene_creed_typography(slide)


def _add_creed_slides(prs: Presentation, theme: SlideTheme, *, creed_choice: str = "nicene") -> None:
    """Nicene or Apostles' Creed — same place in the Mass, never both."""
    if _normalize_creed_choice(creed_choice) == "apostles":
        _add_apostles_creed_slides(prs, theme)
        return
    _add_nicene_creed_slides(prs, theme)


def _style_dialogue_run(
    run,
    *,
    color: RGBColor,
    size_pt: float = _SLIDE_TEXT_PT,
    bold: bool = True,
) -> None:
    run.font.name = _ACTIVE_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _append_projection_dialogue_paragraph(
    tf,
    role: str,
    line: str,
    *,
    first: bool,
    strip_all: bool,
    size_pt: float = _SLIDE_TEXT_PT,
) -> None:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = ""
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
    p.space_after = Pt(8)

    if role == "priest":
        label = p.add_run()
        label.text = "Priest: "
        _style_dialogue_run(label, color=_ACTIVE_THEME.emphasis, size_pt=size_pt)
        body = p.add_run()
        body.text = line
        _style_dialogue_run(body, color=_ACTIVE_THEME.primary, size_pt=size_pt, bold=True)
        return

    if role == "all":
        if not strip_all:
            label = p.add_run()
            label.text = "All: "
            _style_dialogue_run(label, color=_ACTIVE_THEME.primary, size_pt=size_pt)
        body = p.add_run()
        body.text = line
        _style_dialogue_run(body, color=_ACTIVE_THEME.primary, size_pt=size_pt, bold=True)
        return

    body = p.add_run()
    body.text = line
    _style_dialogue_run(
        body,
        color=_ACTIVE_THEME.primary,
        size_pt=size_pt,
        bold=(role == "hymn"),
    )


def _render_projection_dialogue_slide(
    prs: Presentation,
    footer_section: str,
    marked_text: str,
    theme: SlideTheme,
) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, theme.bg)
    _apply_slide_branding(slide, theme)
    lx, w = MARGIN_SIDE, SLIDE_WIDTH - 2 * MARGIN_SIDE
    zone_top = _content_top()
    zone_bottom = SLIDE_HEIGHT - DIALOGUE_BOTTOM_GAP
    body_h = zone_bottom - zone_top

    box = slide.shapes.add_textbox(lx, zone_top, w, body_h)
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    tf.vertical_anchor = BODY_VERTICAL_ANCHOR

    strip_all = _suppress_all_role_prefix(footer_section)
    first = True
    for role, line in _parse_marked_lines(marked_text):
        if role == "direction":
            continue
        _append_projection_dialogue_paragraph(
            tf,
            role,
            line,
            first=first,
            strip_all=strip_all,
        )
        first = False

    _add_community_footer(slide, footer_section, theme)


def _render_marked_slide(
    prs: Presentation,
    footer_section: str,
    marked_text: str,
    theme: SlideTheme,
    *,
    title: Optional[str] = None,
    title_pt: float = _SECTION_TITLE_PT_LARGE,
) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, theme.bg)
    _apply_slide_branding(slide, theme)
    if title:
        _add_section_title(slide, title, title_pt)
    lx, top, w = MARGIN_SIDE, _content_top(), SLIDE_WIDTH - 2 * MARGIN_SIDE
    body_h = SLIDE_HEIGHT - top - CONTENT_BOTTOM_GAP

    box = slide.shapes.add_textbox(lx, top, w, body_h)
    tf = box.text_frame
    _prep_body_tf(tf)
    first = True
    strip_all = _suppress_all_role_prefix(footer_section)
    rite_slide = _is_prayer_rite_slide(footer_section)
    for role, line in _parse_marked_lines(marked_text):
        if role == "direction":
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if role == "priest":
            p.text = f"Priest: {line}"
            _style_para(p, size_pt=_SLIDE_TEXT_PT, color=theme.emphasis, bold=True)
            p.space_before = Pt(4)
        elif role == "all":
            p.text = line if strip_all else f"All: {line}"
            _style_para(p, size_pt=_SLIDE_TEXT_PT, color=theme.primary, bold=True)
            p.space_before = Pt(4)
        elif role == "hymn":
            p.text = line
            _style_para(p, size_pt=_SLIDE_TEXT_PT, color=theme.primary, bold=True)
        else:
            p.text = line
            _style_para(p, size_pt=_SLIDE_TEXT_PT, color=theme.primary, bold=False)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8 if rite_slide else 5)

    _add_community_footer(slide, footer_section, theme)


def _add_marked_slide(
    prs: Presentation,
    footer_section: str,
    marked_text: str,
    theme: SlideTheme,
    *,
    title: Optional[str] = None,
    title_pt: float = _SECTION_TITLE_PT_LARGE,
) -> None:
    marked_text = _strip_marked_rubrics(marked_text)
    if not _marked_has_projectable_content(marked_text):
        return
    if _is_projection_dialogue_slide(footer_section):
        _render_projection_dialogue_slide(prs, footer_section, marked_text, theme)
        return
    if _is_prayer_rite_slide(footer_section):
        chunks = _chunk_marked_rite_by_fit(marked_text, footer_section)
        total = len(chunks)
        for i, ch in enumerate(chunks):
            foot = footer_section if total == 1 else f"{footer_section} ({i + 1}/{total})"
            _render_marked_slide(prs, foot, ch, theme, title=title, title_pt=title_pt)
        return
    _render_marked_slide(prs, footer_section, marked_text, theme, title=title, title_pt=title_pt)


def _chunk_marked_body(marked: str, limit: int = _MAX_MARKED_BODY) -> List[str]:
    if len(marked) <= limit:
        return [marked]
    parts, buf, n = [], [], 0
    for line in marked.split("\n"):
        line_len = len(line) + 1
        if n + line_len > limit and buf:
            parts.append("\n".join(buf))
            buf, n = [line], line_len
        else:
            buf.append(line)
            n += line_len
    if buf:
        parts.append("\n".join(buf))
    return parts if parts else [marked[:limit]]


def _add_marked_chunked(prs: Presentation, footer: str, marked: str, theme: SlideTheme) -> None:
    marked = _strip_marked_rubrics(marked)
    if not _marked_has_projectable_content(marked):
        return
    if _is_prayer_rite_slide(footer):
        chunks = _chunk_marked_rite_by_fit(marked, footer)
    else:
        chunks = _chunk_marked_body(marked)
    for i, ch in enumerate(chunks):
        foot = footer if len(chunks) == 1 else f"{footer} ({i + 1}/{len(chunks)})"
        _add_marked_slide(prs, foot, ch, theme)


def _is_wide_char(ch: str) -> bool:
    """True for CJK / fullwidth glyphs that occupy roughly double Latin width."""
    o = ord(ch)
    return (
        0x1100 <= o <= 0x115F  # Hangul Jamo
        or 0x2E80 <= o <= 0xA4CF  # CJK radicals .. Yi
        or 0xAC00 <= o <= 0xD7A3  # Hangul syllables
        or 0xF900 <= o <= 0xFAFF  # CJK compatibility ideographs
        or 0xFF00 <= o <= 0xFF60  # fullwidth forms
    )


def _rite_text_width_units(text: str) -> float:
    return float(sum(2 if _is_wide_char(c) else 1 for c in text))


def _fit_rite_font_pt(
    items: List[Tuple[str, str]],
    *,
    strip_all: bool,
    body_h: float,
    max_pt: float = _SLIDE_TEXT_PT,
    min_pt: float = 20.0,
) -> float:
    """Largest font (<= rite default) that fits all lines on one slide.

    Width model: a glyph advances ~0.55 em on average; CJK glyphs count as
    two width units (see ``_rite_text_width_units``). Wide-format slides are
    20in across, so this avoids over-shrinking text into a narrow column.
    """
    usable_w_in = _length_to_inches(SLIDE_WIDTH - 2 * MARGIN_SIDE) * 0.96
    for font_pt in range(int(max_pt), int(min_pt) - 1, -1):
        line_h = _rite_line_height_inches(font_pt)
        spacing = (font_pt * 0.14 + 4) / 72.0
        em_in = font_pt / 72.0
        cpl = max(1.0, usable_w_in / (em_in * 0.55))
        total = 0.0
        for role, line in items:
            text = _rite_display_line(role, line, strip_all=strip_all)
            wrapped = max(1, math.ceil(_rite_text_width_units(text) / cpl))
            total += wrapped * line_h + spacing
        if total <= body_h:
            return float(font_pt)
    return float(min_pt)


_OUR_FATHER_TITLES = {
    "english": "Our Father",
    "tagalog": "Ama Namin",
    "visaya": "Amahan Namo",
    "malay": "Bapa Kami",
    "korean": "주님의 기도",
}


def _add_our_father_slide(
    prs: Presentation,
    footer: str,
    marked: str,
    theme: SlideTheme,
    *,
    title: Optional[str] = None,
) -> None:
    """Render the complete Our Father on a single slide, auto-fitting the font."""
    marked = _strip_marked_rubrics(marked)
    if not _marked_has_projectable_content(marked):
        return
    items = [(r, ln) for r, ln in _parse_marked_lines(marked) if r != "direction"]
    if not items:
        return

    strip_all = _suppress_all_role_prefix(footer)
    body_h = _marked_body_height_inches()
    font_pt = _fit_rite_font_pt(items, strip_all=strip_all, body_h=body_h)

    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, theme.bg)
    _apply_slide_branding(slide, theme)
    if title:
        _add_section_title(slide, title, _SECTION_TITLE_PT_SMALL)
    lx, top, w = MARGIN_SIDE, _content_top(), SLIDE_WIDTH - 2 * MARGIN_SIDE
    box_h = SLIDE_HEIGHT - top - CONTENT_BOTTOM_GAP

    box = slide.shapes.add_textbox(lx, top, w, box_h)
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    tf.vertical_anchor = BODY_VERTICAL_ANCHOR

    space_after = Pt(max(2.0, font_pt * 0.14))
    first = True
    for role, line in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if role == "priest":
            p.text = f"Priest: {line}"
            _style_para(p, size_pt=font_pt, color=theme.emphasis, bold=True)
            p.space_before = Pt(4)
        elif role == "all":
            p.text = line if strip_all else f"All: {line}"
            _style_para(p, size_pt=font_pt, color=theme.primary, bold=True)
            p.space_before = Pt(4)
        elif role == "hymn":
            p.text = line
            _style_para(p, size_pt=font_pt, color=theme.primary, bold=True)
        else:
            p.text = line
            _style_para(p, size_pt=font_pt, color=theme.primary, bold=False)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = space_after

    _add_community_footer(slide, footer, theme)


# Mass divider poster (GFCC layout — rounded panels, season-themed palette)
_DIVIDER_FONT = "Arial"
_DIVIDER_CORNER_ADJ = 0.1667  # matches reference freeform corner radius (~16.67% of height)
_DIVIDER_PANEL_ALPHA = 44706  # ~44.7% opaque (reference right panel)
_DIVIDER_BAR_ALPHA = 60784  # ~60.8% opaque (reference bottom bar)
_DIVIDER_LINE_W = 14299  # reference outline weight (EMU)


@dataclass(frozen=True)
class _DividerPalette:
    grad_start: RGBColor
    grad_end: RGBColor
    panel_fill: RGBColor
    panel_border: RGBColor
    bar_fill: RGBColor
    bar_border: RGBColor
    label: RGBColor
    primary: RGBColor
    quote: RGBColor
    gospel_label: RGBColor


def _rgb_channels(color: RGBColor) -> tuple[int, int, int]:
    return int(color[0]), int(color[1]), int(color[2])


def _rgb_from_channels(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(_clamp_byte(r), _clamp_byte(g), _clamp_byte(b))


def _mix_rgb(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
    ar, ag, ab = _rgb_channels(a)
    br, bg, bb = _rgb_channels(b)
    t = max(0.0, min(1.0, t))
    return _rgb_from_channels(
        ar + (br - ar) * t,
        ag + (bg - ag) * t,
        ab + (bb - ab) * t,
    )


def _divider_palette(theme: SlideTheme) -> _DividerPalette:
    """Divider colors for the Mass section cards.

    Theme 1 follows the liturgical season via ``divider_*``. Mono themes (2–3) keep a
    solid black or white divider with matching text — no season tint.
    """
    bg = theme.divider_bg
    if theme.mono_surfaces:
        primary = theme.divider_primary
        muted = theme.divider_muted
        emphasis = theme.divider_emphasis
        panel_fill = _mix_rgb(bg, primary, 0.06)
        panel_border = _mix_rgb(bg, primary, 0.18)
        bar_fill = _mix_rgb(bg, primary, 0.12)
        return _DividerPalette(
            grad_start=bg,
            grad_end=bg,
            panel_fill=panel_fill,
            panel_border=panel_border,
            bar_fill=bar_fill,
            bar_border=bar_fill,
            label=emphasis,
            primary=primary,
            quote=muted,
            gospel_label=emphasis,
        )
    black = RGBColor(8, 8, 10)
    white = RGBColor(255, 255, 255)
    grad_start = _mix_rgb(bg, black, 0.82)
    grad_end = _mix_rgb(bg, white, 0.28)
    panel_fill = bg
    panel_border = _mix_rgb(bg, white, 0.34)
    bar_fill = _mix_rgb(bg, black, 0.42)
    bar_border = bar_fill
    label = _mix_rgb(theme.divider_emphasis, white, 0.42)
    quote = _mix_rgb(theme.divider_primary, theme.divider_muted, 0.18)
    gospel_label = _mix_rgb(theme.divider_emphasis, white, 0.12)
    return _DividerPalette(
        grad_start=grad_start,
        grad_end=grad_end,
        panel_fill=panel_fill,
        panel_border=panel_border,
        bar_fill=bar_fill,
        bar_border=bar_border,
        label=label,
        primary=theme.divider_primary,
        quote=quote,
        gospel_label=gospel_label,
    )


def _apply_solid_fill_alpha(shape, rgb: RGBColor, alpha_val: int) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    sp_pr = shape._element.spPr
    solid = sp_pr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(alpha_val))


def _apply_line_alpha(shape, rgb: RGBColor, alpha_val: int, *, width_emu: int = _DIVIDER_LINE_W) -> None:
    sp_pr = shape._element.spPr
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(sp_pr, qn("a:ln"))
    ln.set("w", str(width_emu))
    ln.set("cap", "sq")
    for child in list(ln):
        ln.remove(child)
    solid_fill = etree.SubElement(ln, qn("a:solidFill"))
    r, g, b = _rgb_channels(rgb)
    srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(alpha_val))


def _divider_add_rounded_panel(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_rgb: RGBColor,
    border_rgb: RGBColor,
    alpha_val: int,
) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = _DIVIDER_CORNER_ADJ
    _apply_solid_fill_alpha(shp, fill_rgb, alpha_val)
    _apply_line_alpha(shp, border_rgb, alpha_val)


def _set_divider_gradient_bg(slide, start: RGBColor, end: RGBColor) -> None:
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 0.0
    stops = fill.gradient_stops
    stops[0].color.rgb = start
    stops[1].color.rgb = end


def _divider_date_display(date: str) -> str:
    raw = (date or "").strip()
    if not raw:
        return ""
    for fmt in ("%Y-%m-%d", "%B %d, %Y", "%b %d, %Y", "%d %B %Y"):
        try:
            d = _dt.datetime.strptime(raw, fmt).date()
            return f"{d.day} {d.strftime('%B').upper()} {d.year}"
        except ValueError:
            continue
    return raw.upper()


def _divider_year_date_line(lectionary_cycle: str, date: str) -> str:
    cycle = (lectionary_cycle or "—").strip().upper()
    date_line = _divider_date_display(date)
    if date_line:
        return f"YEAR {cycle} · {date_line}"
    return f"YEAR {cycle}"


def _divider_gospel_heading(gospel_reference: str) -> str:
    ref = (gospel_reference or "").strip() or "—"
    ref = ref.replace("–", "–").replace("-", "–")
    return f"GOSPEL ({ref.upper()})"


def _divider_quote_lines(quote: str) -> List[str]:
    q = (quote or "").strip()
    if not q:
        return []
    if q[0] not in "\"“‘":
        q = f"\u201c{q}"
    if q[-1] not in "\"”’":
        q = f"{q}\u201d"
    parts = re.split(r"(?<=[;.])\s+", q)
    return [p.strip() for p in parts if p.strip()]


def _divider_est_lines(text: str, width_in: float, pt: float) -> int:
    plain = (text or "").strip()
    if not plain:
        return 0
    chars_per_line = max(6, int(width_in * 72 / max(pt * 0.52, 1)))
    return max(1, math.ceil(len(plain) / chars_per_line))


def _divider_fit_font_pt(
    chunks: List[str],
    *,
    width_in: float,
    height_in: float,
    max_pt: float,
    min_pt: float,
) -> float:
    lines = [c for c in chunks if (c or "").strip()]
    if not lines:
        return max_pt
    pt = max_pt
    while pt > min_pt:
        total_lines = sum(_divider_est_lines(line, width_in, pt) for line in lines)
        line_h = (pt * 1.14) / 72.0
        if total_lines * line_h <= height_in * 0.92:
            break
        pt -= 1
    return max(min_pt, pt)


def _divider_fit_single_line_pt(
    text: str,
    *,
    width_in: float,
    max_pt: float,
    min_pt: float,
) -> float:
    """Largest point size at which ``text`` fits on a single line within ``width_in``."""
    plain = (text or "").strip()
    if not plain:
        return max_pt
    pt = max_pt
    while pt > min_pt and _divider_est_lines(plain, width_in, pt) > 1:
        pt -= 1
    return max(min_pt, pt)


def _divider_add_textbox(
    slide,
    *,
    left,
    top,
    width,
    height,
    lines: List[Tuple[str, dict]],
    anchor_middle: bool = False,
    no_wrap: bool = False,
) -> None:
    """Add a textbox; each item is (text, style kwargs for _style_para)."""
    if not any((t or "").strip() for t, _ in lines):
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    if no_wrap:
        tf.word_wrap = False
    if anchor_middle:
        tf.vertical_anchor = BODY_VERTICAL_ANCHOR
    first = True
    for text, style in lines:
        if not (text or "").strip():
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text.strip()
        _style_para(
            p,
            size_pt=style.get("size_pt", _SLIDE_TEXT_PT),
            color=style.get("color", RGBColor(255, 255, 255)),
            bold=style.get("bold", False),
            italic=style.get("italic", False),
            font_name=style.get("font_name", _DIVIDER_FONT),
        )
        p.alignment = style.get("align", PP_ALIGN.CENTER)
        p.line_spacing = style.get("line_spacing", 1.08)
        if style.get("space_after"):
            p.space_after = Pt(style["space_after"])


def _render_default_divider_cover(
    slide,
    *,
    celebrant: str,
    co_celebrant: str = "",
    date: str,
    mass_title: str,
    season: str,
    lectionary_cycle: str,
    gospel_reference: str,
    gospel_quote: str,
    quote_max_chars: int,
    theme: SlideTheme,
) -> None:
    pal = _divider_palette(theme)
    _set_divider_gradient_bg(slide, pal.grad_start, pal.grad_end)
    _divider_add_rounded_panel(
        slide,
        Inches(7.498),
        Inches(1.388),
        Inches(11.968),
        Inches(6.702),
        fill_rgb=pal.panel_fill,
        border_rgb=pal.panel_border,
        alpha_val=_DIVIDER_PANEL_ALPHA,
    )
    _divider_add_rounded_panel(
        slide,
        Inches(0.985),
        Inches(9.299),
        Inches(18.03),
        Inches(1.217),
        fill_rgb=pal.bar_fill,
        border_rgb=pal.bar_border,
        alpha_val=_DIVIDER_BAR_ALPHA,
    )

    celebrant_name = (celebrant or "").strip() or "—"
    celebrant_pt = _divider_fit_single_line_pt(
        celebrant_name,
        width_in=7.206,
        max_pt=61,
        min_pt=20,
    )
    _divider_add_textbox(
        slide,
        left=Inches(1.405),
        top=Inches(2.614),
        width=Inches(4.604),
        height=Inches(0.698),
        lines=[("MASS CELEBRANT:", {"size_pt": 37, "color": pal.label, "bold": True})],
        anchor_middle=True,
    )
    _divider_add_textbox(
        slide,
        left=Inches(0.292),
        top=Inches(3.339),
        width=Inches(7.206),
        height=Inches(1.156),
        lines=[(celebrant_name, {"size_pt": celebrant_pt, "color": pal.primary, "bold": True})],
        anchor_middle=True,
        no_wrap=True,
    )

    year_date_line = _divider_year_date_line(lectionary_cycle, date)
    year_date_pt = _divider_fit_font_pt(
        [year_date_line],
        width_in=7.206,
        height_in=0.95,
        max_pt=49,
        min_pt=28,
    )
    _divider_add_textbox(
        slide,
        left=Inches(0.292),
        top=Inches(4.844),
        width=Inches(7.206),
        height=Inches(0.95),
        lines=[
            (
                year_date_line,
                {"size_pt": year_date_pt, "color": pal.primary, "bold": True, "italic": True},
            )
        ],
        anchor_middle=True,
    )

    # Optional co-celebrant block (template TextBox 9 + 10). Only rendered when a
    # co-celebrant name is supplied; otherwise the placeholder is omitted entirely.
    co_name = (co_celebrant or "").strip()
    if co_name:
        _divider_add_textbox(
            slide,
            left=Inches(1.405),
            top=Inches(6.208),
            width=Inches(4.604),
            height=Inches(0.698),
            lines=[("CO-CELEBRANT:", {"size_pt": 37, "color": pal.label, "bold": True})],
            anchor_middle=True,
        )
        co_pt = _divider_fit_single_line_pt(
            co_name,
            width_in=7.206,
            max_pt=52,
            min_pt=20,
        )
        _divider_add_textbox(
            slide,
            left=Inches(0.104),
            top=Inches(6.966),
            width=Inches(7.206),
            height=Inches(1.156),
            lines=[(co_name, {"size_pt": co_pt, "color": pal.primary, "bold": True})],
            anchor_middle=True,
            no_wrap=True,
        )

    g_line = (gospel_quote or "").strip()
    if quote_max_chars and len(g_line) > quote_max_chars:
        g_line = g_line[: quote_max_chars - 1].rstrip() + "\u2026"
    quote_parts = _divider_quote_lines(g_line)
    quote_w_in = 11.784
    quote_h_in = 2.649
    quote_pt = _divider_fit_font_pt(
        quote_parts or [g_line],
        width_in=quote_w_in,
        height_in=quote_h_in,
        max_pt=39,
        min_pt=22,
    )
    quote_style = {"size_pt": quote_pt, "color": pal.quote, "bold": False}
    _divider_add_textbox(
        slide,
        left=Inches(7.546),
        top=Inches(2.614),
        width=Inches(quote_w_in),
        height=Inches(quote_h_in),
        lines=[(part, quote_style) for part in (quote_parts or ([g_line] if g_line else []))],
        anchor_middle=True,
    )

    g_heading = _divider_gospel_heading(gospel_reference)
    g_head_pt = _divider_fit_font_pt(
        [g_heading],
        width_in=6.364,
        height_in=0.771,
        max_pt=41,
        min_pt=24,
    )
    _divider_add_textbox(
        slide,
        left=Inches(10.256),
        top=Inches(6.345),
        width=Inches(6.364),
        height=Inches(0.771),
        lines=[(g_heading, {"size_pt": g_head_pt, "color": pal.gospel_label, "bold": True})],
        anchor_middle=True,
    )

    bottom_title = (mass_title or season or "Sunday Mass").strip()
    bottom_title = bottom_title.replace(" Celebration", "").strip() or "Sunday Mass"
    bottom_pt = _divider_fit_font_pt(
        [bottom_title],
        width_in=13.52,
        height_in=1.271,
        max_pt=68,
        min_pt=32,
    )
    _divider_add_textbox(
        slide,
        left=Inches(3.495),
        top=Inches(9.199),
        width=Inches(13.52),
        height=Inches(1.271),
        lines=[(bottom_title, {"size_pt": bottom_pt, "color": pal.primary, "bold": True})],
        anchor_middle=True,
    )


def _add_divider_cover(
    prs: Presentation,
    *,
    celebrant: str,
    co_celebrant: str = "",
    date: str,
    season: str,
    lectionary_cycle: str,
    gospel_reference: str,
    gospel_quote: str,
    quote_max_chars: int,
    theme: SlideTheme,
    mass_title: str = "",
    background_image_path: Optional[Path] = None,
    divider_poster_path: Optional[Path] = None,
) -> None:
    del background_image_path  # GFCC divider uses season theme; poster is not overlaid here
    slide = prs.slides.add_slide(_layout_blank(prs))

    image_divider = bool(
        divider_poster_path and Path(divider_poster_path).is_file()
    )
    if image_divider:
        cover = Path(divider_poster_path).resolve()
        slide.shapes.add_picture(
            str(cover),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        return

    _render_default_divider_cover(
        slide,
        celebrant=celebrant,
        co_celebrant=co_celebrant,
        date=date,
        mass_title=mass_title,
        season=season,
        lectionary_cycle=lectionary_cycle,
        gospel_reference=gospel_reference,
        gospel_quote=gospel_quote,
        quote_max_chars=quote_max_chars,
        theme=theme,
    )


def _add_section_card(prs: Presentation, big_lines: str, footer_section: str, theme: SlideTheme) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, theme.bg)
    _apply_slide_branding(slide, theme)
    lx, top, w = MARGIN_SIDE, _content_top() + Inches(0.25), SLIDE_WIDTH - 2 * MARGIN_SIDE
    box = slide.shapes.add_textbox(lx, top, w, Inches(4.5))
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = big_lines
    _style_para(p, size_pt=_SLIDE_TEXT_PT, color=theme.emphasis, bold=True)
    p.alignment = PP_ALIGN.CENTER
    _add_community_footer(slide, footer_section, theme)


def _length_to_inches(value: Any) -> float:
    """Convert python-pptx Length or raw EMU int to inches."""
    if hasattr(value, "inches"):
        return float(value.inches)
    return float(value) / _EMU_PER_INCH


def _lyric_fit_height_inches(box_height_inches: float) -> float:
    """Usable lyric height after text-frame top/bottom padding from ``_prep_tf``."""
    pad = _length_to_inches(Inches(0.08)) * 2
    return max(0.5, float(box_height_inches) - pad)


def split_lyrics(lines: List[str], max_lines: int = _LYRIC_MAX_LINES_PER_SLIDE) -> List[str]:
    """Group lyric lines into slide blocks (each block is lines joined with newlines)."""
    if not lines:
        return []
    blocks: List[str] = []
    for i in range(0, len(lines), max_lines):
        blocks.append("\n".join(lines[i : i + max_lines]))
    return blocks


def _token_width_units(token: str) -> float:
    """Approximate rendered token width in relative units."""
    if not token:
        return 0.0
    narrow = set(".,:;!'|ijlIt`")
    wide = set("MW@#%&QGOD")
    units = 0.0
    for ch in token:
        if ch == " ":
            units += 0.33
        elif ch in narrow:
            units += 0.38
        elif ch in wide:
            units += 0.9
        else:
            units += 0.62
    return units


def measureRenderedText(lines: List[str], font_size_pt: float) -> Mapping[str, float]:
    """
    Estimate text footprint for full-width lyric fitting.

    Uses a conservative width model to keep projector readability.
    """
    safe_width_inches = float(SLIDE_WIDTH.inches) * _LYRIC_TEXTBOX_WIDTH_RATIO
    line_height_inches = (
        font_size_pt * _HYMN_REF_LINE_SPACING * _LYRIC_LINE_PITCH_FACTOR
    ) / 72.0
    avail_width_inches = safe_width_inches * _LYRIC_FIT_WIDTH_SAFETY
    max_line_inches = 0.0
    visual_line_count = 0
    for line in lines:
        units = _token_width_units(line.strip().upper())
        estimated = (units * font_size_pt * _LYRIC_WIDTH_CALIBRATION) / 72.0
        if estimated > max_line_inches:
            max_line_inches = estimated
        # Account for soft word-wrap: a line wider than the box becomes multiple
        # visual lines, which is what actually pushes text past the box bounds.
        if avail_width_inches > 0 and estimated > avail_width_inches:
            visual_line_count += int(math.ceil(estimated / avail_width_inches))
        else:
            visual_line_count += 1
    return {
        "max_line_width_inches": max_line_inches,
        "text_height_inches": visual_line_count * line_height_inches,
        "available_width_inches": safe_width_inches,
    }


def detectOverflow(lines: List[str], font_size_pt: float, box_height_inches: float) -> bool:
    measured = measureRenderedText(lines, font_size_pt)
    max_w = measured["available_width_inches"] * _LYRIC_FIT_WIDTH_SAFETY
    max_h = box_height_inches * _LYRIC_FIT_HEIGHT_SAFETY
    return bool(
        measured["max_line_width_inches"] > max_w
        or measured["text_height_inches"] > max_h
    )


def _lyric_lines_from_chunk(lyrics_text: str) -> List[str]:
    """
    Preserve line breaks from Lyrics Studio / saved hymn blocks.

    Only re-wraps individual lines that exceed the soft character limit.
    """
    raw_lines = [ln.strip() for ln in (lyrics_text or "").splitlines() if ln.strip()]
    if not raw_lines:
        return []
    out: List[str] = []
    for raw in raw_lines:
        for seg, is_paren in _split_paren_echo_segments(raw):
            # A parenthetical echo always occupies its own row and is never broken,
            # regardless of length (the auto-fit shrinks the font instead).
            if is_paren or len(seg) <= _LYRIC_SOFT_WRAP_CHARS:
                out.append(seg)
                continue
            wrapped = optimizeLineBreaks(seg)
            out.extend(wrapped if wrapped else [seg])
    return out


def optimizeLineBreaks(lyrics_text: str) -> List[str]:
    """
    Phrase-aware line optimizer for worship lyrics.

    Prefers natural breathing punctuation and conjunction boundaries.
    Each line keeps at least ``_LYRIC_MIN_WORDS_PER_LINE`` words when split.
    """
    min_words = _LYRIC_MIN_WORDS_PER_LINE
    raw_lines = [ln.strip() for ln in (lyrics_text or "").splitlines() if ln.strip()]
    out: List[str] = []
    break_re = re.compile(r"\s+(,|;|:|\.|—|-|and|but|for|that|with|to)\s+", flags=re.IGNORECASE)
    for raw in raw_lines:
        words = raw.split()
        if len(words) <= (min_words * 2 - 1):
            out.append(raw)
            continue
        candidate = None
        midpoint = len(raw) // 2
        for m in break_re.finditer(raw):
            idx = m.start()
            if _break_splits_parens(raw, idx):
                continue
            left_n = _word_count(raw[:idx])
            right_n = _word_count(raw[idx:])
            if left_n < min_words or right_n < min_words:
                continue
            if candidate is None or abs(idx - midpoint) < abs(candidate - midpoint):
                candidate = idx
        if candidate is not None:
            first, second = _split_line_preserving_parens(raw, candidate)
            if first and second and _word_count(first) >= min_words and _word_count(second) >= min_words:
                out.extend([first, second])
                continue
        # Balanced split with minimum words on each line.
        cut = max(min_words, min(len(words) - min_words, len(words) // 2))
        if cut < min_words or len(words) - cut < min_words:
            out.append(raw)
            continue
        prefix = " ".join(words[:cut])
        idx = raw.find(prefix)
        split_at = idx + len(prefix) if idx >= 0 else len(prefix)
        first, second = _split_line_preserving_parens(raw, split_at)
        if first and second and _word_count(first) >= min_words and _word_count(second) >= min_words:
            out.extend([first, second])
        else:
            out.append(raw)
    return _enforce_min_words_per_line(out)


def calculateOptimalFontSize(lines: List[str], box_height_inches: float) -> int:
    """Choose largest lyric size in 68–75 pt that fits the textbox (reference hymn slides)."""
    line_count = len(lines)
    height_cap = _LYRIC_MAX_PT
    if line_count >= 8:
        height_cap = min(height_cap, int(_HYMN_REF_BODY_PT_MIN))
    elif line_count >= 6:
        height_cap = min(height_cap, 72)
    elif line_count >= 5:
        height_cap = min(height_cap, 74)
    floor_pt = int(_HYMN_REF_BODY_PT_MIN)
    for pt in range(height_cap, floor_pt - 1, -1):
        if not detectOverflow(lines, float(pt), box_height_inches):
            return pt
    for pt in range(floor_pt - 1, _LYRIC_MIN_PT - 1, -1):
        if not detectOverflow(lines, float(pt), box_height_inches):
            return pt
    return _LYRIC_MIN_PT


def fitLyricsToFullWidthTextbox(lyrics_text: str, box_height_inches: float) -> Tuple[List[str], int]:
    """
    Fit lyrics into fixed near-edge-to-edge textbox.

    Keeps constant width and adjusts font size; preserves structured block line breaks.
    """
    lines = _lyric_lines_from_chunk(lyrics_text)
    if not lines:
        return [""], _LYRIC_MIN_PT

    fitted = lines
    size_pt = calculateOptimalFontSize(fitted, box_height_inches)
    while size_pt > _LYRIC_MIN_PT and detectOverflow(fitted, float(size_pt), box_height_inches):
        size_pt -= 2
    return fitted, size_pt


def _chunk_section_for_slides(section_text: str, max_lines: int = _LYRIC_MAX_LINES_PER_SLIDE) -> List[str]:
    """One structured block (verse/chorus); sub-chunk only when that block exceeds line budget."""
    cleaned = clean_lyrics_for_projection(section_text)
    if not cleaned:
        return []
    line_list = [ln for ln in cleaned.splitlines() if ln.strip()]
    if not line_list:
        return [cleaned]
    if len(line_list) > _HYMN_DUAL_SOLO_PARAGRAPH_THRESHOLD:
        return ["\n".join(line_list)]
    if len(line_list) <= max_lines:
        return ["\n".join(line_list)]
    return split_lyrics(line_list, max_lines=max_lines)


def _chunk_lyrics_display(text: str, max_lines: int = _LYRIC_MAX_LINES_PER_SLIDE) -> List[Tuple[str, str]]:
    """
    Split lyrics for hymn slides by structured-editor blocks (verse, chorus, bridge, etc.).

    Each blank-line-separated section from Lyrics Studio becomes its own slide group.
    Long sections may span multiple slides, but chunks never cross section boundaries.
    Returns ``(chunk_text, block_kind)`` pairs (``block_kind`` is verse|chorus|…).
    """
    t = ensure_lyric_section_breaks((text or "").strip())
    if not t:
        return []

    chunks: List[Tuple[str, str]] = []
    for block_kind, section in parse_structured_lyric_sections_typed(t):
        for chunk in _chunk_section_for_slides(section, max_lines=max_lines):
            chunks.append((chunk, block_kind))

    return chunks if chunks else [(t, "verse")]


def _pp_align(name: str) -> PP_ALIGN:
    key = (name or "center").strip().lower()
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(key, PP_ALIGN.CENTER)


def _auto_body_pt(line_count: int, longest: int, base_pt: float) -> float:
    size_pt = base_pt
    if line_count >= 8 or longest >= 40:
        size_pt = min(size_pt, 50.0)
    if line_count >= 10 or longest >= 48:
        size_pt = min(size_pt, 44.0)
    if line_count >= 12 or longest >= 56:
        size_pt = min(size_pt, 38.0)
    if line_count >= 14 or longest >= 64:
        size_pt = min(size_pt, 34.0)
    return size_pt


def _is_chorus_block_kind(block_kind: str) -> bool:
    kind = (block_kind or "").strip().lower().replace("_", "-")
    return kind in ("chorus", "refrain", "post-chorus", "hook", "tag")


def _dual_second_verse_italic(
    group: List[Tuple[str, str]],
    block_index: int,
) -> bool:
    """Italicize the lower block when a dual slide pairs two non-chorus verses."""
    if len(group) != 2 or block_index != 1:
        return False
    return all(not _is_chorus_block_kind(kind) for _text, kind in group)


def _fill_hymn_body_caps(
    tf,
    chunk: str,
    *,
    typography: Optional[HymnTypographySettings] = None,
    box_height_inches: Optional[float] = None,
    block_kind: str = "verse",
    italic_body: bool = False,
) -> None:
    """ALL CAPS Poppins on black; reference deck uses 75 pt and 0.7 line spacing."""
    box_h = float(box_height_inches or 0.0) or float(SLIDE_HEIGHT.inches * 0.72)
    fit_h = _lyric_fit_height_inches(box_h)
    lines, auto_fit_pt = fitLyricsToFullWidthTextbox(chunk, fit_h)
    size_pt = int(max(_HYMN_REF_BODY_PT_MIN, min(_LYRIC_MAX_PT, auto_fit_pt)))
    if typography:
        requested = int(round(typography.body_pt))
        if requested >= _HYMN_REF_BODY_PT_MIN:
            size_pt = min(size_pt, requested)
    while size_pt > _LYRIC_MIN_PT and detectOverflow(lines, float(size_pt), fit_h):
        size_pt -= 2
    align = _pp_align(typography.body_align if typography else "center")
    is_chorus = _is_chorus_block_kind(block_kind)
    body_color = _ACTIVE_THEME.chorus_accent if is_chorus else _ACTIVE_THEME.hymn_body
    use_italic = is_chorus or italic_body

    def _style_hymn_body_run(run, *, is_paren: bool) -> None:
        font = run.font
        font.name = _HYMN_REF_BODY_FONT
        font.size = Pt(size_pt)
        font.bold = True
        font.italic = use_italic or is_paren
        font.color.rgb = _ACTIVE_THEME.paren_accent if is_paren else body_color

    def _apply_body_para(p, text: str) -> None:
        p.alignment = align
        p.line_spacing = _HYMN_REF_LINE_SPACING
        segments: List[Tuple[str, bool]] = []
        last = 0
        for match in _PAREN_IN_LYRICS_RE.finditer(text):
            if match.start() > last:
                segments.append((text[last : match.start()], False))
            segments.append((match.group(0), True))
            last = match.end()
        if last < len(text):
            segments.append((text[last:], False))
        if not segments:
            segments = [(text, False)]

        first_seg, first_paren = segments[0]
        p.text = first_seg
        if p.runs:
            _style_hymn_body_run(p.runs[0], is_paren=first_paren)
        for seg, is_paren in segments[1:]:
            run = p.add_run()
            run.text = seg
            _style_hymn_body_run(run, is_paren=is_paren)

    tf.clear()
    first = True
    for raw in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _apply_body_para(p, raw.upper())
    if first:
        _apply_body_para(tf.paragraphs[0], (chunk or "").strip().upper())


def _normalize_hymn_lyrics_layout(layout: str) -> str:
    key = (layout or "").strip().lower().replace("-", "_")
    if key in ("dual", "two", "2", "double", "pair", "pairs"):
        return "dual"
    return "single"


def _lyric_blocks_for_slides(text: str) -> List[Tuple[str, str]]:
    """Structured lyric blocks (verse, chorus, …) for slide pairing."""
    t = ensure_lyric_section_breaks((text or "").strip())
    if not t:
        return []
    blocks: List[Tuple[str, str]] = []
    for block_kind, section in parse_structured_lyric_sections_typed(t):
        for chunk in _chunk_section_for_slides(section):
            blocks.append((chunk, block_kind))
    return blocks if blocks else [(t, "verse")]


def _pair_blocks_for_dual_slides(
    blocks: List[Tuple[str, str]],
) -> List[List[Tuple[str, str]]]:
    """Pair two blocks per slide unless a block exceeds the solo paragraph threshold."""
    slides: List[List[Tuple[str, str]]] = []
    i = 0
    n = len(blocks)
    while i < n:
        if _lyric_paragraph_count(blocks[i][0]) > _HYMN_DUAL_SOLO_PARAGRAPH_THRESHOLD:
            slides.append([blocks[i]])
            i += 1
            continue
        if (
            i + 1 < n
            and _lyric_paragraph_count(blocks[i + 1][0]) <= _HYMN_DUAL_SOLO_PARAGRAPH_THRESHOLD
        ):
            slides.append([blocks[i], blocks[i + 1]])
            i += 2
        else:
            slides.append([blocks[i]])
            i += 1
    return slides


def _apply_hymn_song_title(
    para,
    title: str,
    *,
    title_pt: float,
    title_align: PP_ALIGN,
) -> None:
    para.text = title
    para.alignment = title_align
    targets = list(para.runs) if para.runs else [para]
    for target in targets:
        font = target.font
        font.name = _HYMN_TITLE_FONT
        font.size = Pt(title_pt)
        font.bold = True
        font.underline = True
        font.color.rgb = _ACTIVE_THEME.hymn_title


def _add_hymn_lyric_box(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    chunk: str,
    block_kind: str,
    *,
    typography: Optional[HymnTypographySettings] = None,
    italic_body: bool = False,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    _prep_hymn_lyric_tf(tf)
    tf.word_wrap = True
    tf.vertical_anchor = BODY_VERTICAL_ANCHOR
    _fill_hymn_body_caps(
        tf,
        chunk,
        typography=typography,
        box_height_inches=_length_to_inches(height),
        block_kind=block_kind,
        italic_body=italic_body,
    )


def _add_hymn_lyric_slides_single(
    prs: Presentation,
    footer_section: str,
    hymn_title: str,
    lyrics: str,
    *,
    hymn_typography: Optional[Mapping[str, Any]] = None,
    section: str = "",
) -> None:
    """One structured block per slide."""
    title = (hymn_title or "Hymn").strip()
    raw_lyrics = (lyrics or "").strip() or "(No lyrics in library for this hymn.)"
    chunks = _chunk_lyrics_display(raw_lyrics)
    if not chunks:
        chunks = [(raw_lyrics, "verse")]

    first_chunk, first_kind = chunks[0]
    rest_chunks = chunks[1:]

    slide0 = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide0, _ACTIVE_THEME.hymn_bg)
    _apply_hymn_branding(slide0)

    w = SLIDE_WIDTH - 2 * MARGIN_SIDE
    title_top = _HYMN_TITLE_TOP
    title_box = slide0.shapes.add_textbox(MARGIN_SIDE, title_top, w, HYMN_TITLE_BOX_H)
    tft = title_box.text_frame
    _prep_tf(tft)
    tft.clear()
    typo0 = typography_for_hymn_slide(hymn_typography, section, 0)
    title_pt = max(typo0.title_pt, _HYMN_REF_TITLE_PT)
    _apply_hymn_song_title(
        tft.paragraphs[0],
        title,
        title_pt=title_pt,
        title_align=_pp_align(typo0.title_align),
    )

    body_top = title_top + HYMN_BODY_TOP_OFFSET
    body_h = SLIDE_HEIGHT - body_top - _LYRIC_BODY_BOTTOM_MARGIN
    lyric_left, lyric_w = _lyric_textbox_geometry(prs.slide_width)
    _add_hymn_lyric_box(
        slide0,
        lyric_left,
        body_top,
        lyric_w,
        body_h,
        first_chunk,
        first_kind,
        typography=typo0,
    )
    _add_hymn_footer(slide0, footer_section)

    for slide_idx, (chunk, block_kind) in enumerate(rest_chunks, start=1):
        typo_n = typography_for_hymn_slide(hymn_typography, section, slide_idx)
        slide = prs.slides.add_slide(_layout_blank(prs))
        _set_slide_bg(slide, _ACTIVE_THEME.hymn_bg)
        _apply_hymn_branding(slide)
        cont_left, cont_top, cont_w, cont_h = _lyric_continuation_textbox_geometry(
            prs.slide_width, prs.slide_height
        )
        _add_hymn_lyric_box(
            slide,
            cont_left,
            cont_top,
            cont_w,
            cont_h,
            chunk,
            block_kind,
            typography=typo_n,
        )
        _add_hymn_footer(slide, footer_section)


def _add_hymn_lyric_slides_dual(
    prs: Presentation,
    footer_section: str,
    hymn_title: str,
    lyrics: str,
    *,
    hymn_typography: Optional[Mapping[str, Any]] = None,
    section: str = "",
) -> None:
    """Two blocks per slide (reference slides 3–5); odd remainder is full-bleed."""
    title = (hymn_title or "Hymn").strip()
    raw_lyrics = (lyrics or "").strip() or "(No lyrics in library for this hymn.)"
    blocks = _lyric_blocks_for_slides(raw_lyrics)
    slide_groups = _pair_blocks_for_dual_slides(blocks)
    full_w = int(prs.slide_width)

    for group_i, group in enumerate(slide_groups):
        typo = typography_for_hymn_slide(hymn_typography, section, group_i)
        slide = prs.slides.add_slide(_layout_blank(prs))
        _set_slide_bg(slide, _ACTIVE_THEME.hymn_bg)
        _apply_hymn_branding(slide)

        if len(group) == 1 and group_i == 0:
            title_top = _HYMN_TITLE_TOP
            w = SLIDE_WIDTH - 2 * MARGIN_SIDE
            title_box = slide.shapes.add_textbox(MARGIN_SIDE, title_top, w, HYMN_TITLE_BOX_H)
            tft = title_box.text_frame
            _prep_tf(tft)
            tft.clear()
            title_pt = max(typo.title_pt, _HYMN_REF_TITLE_PT)
            _apply_hymn_song_title(
                tft.paragraphs[0],
                title,
                title_pt=title_pt,
                title_align=_pp_align(typo.title_align),
            )
            body_top = title_top + HYMN_BODY_TOP_OFFSET
            body_h = SLIDE_HEIGHT - body_top - _LYRIC_BODY_BOTTOM_MARGIN
            lyric_left, lyric_w = _lyric_textbox_geometry(prs.slide_width)
            _add_hymn_lyric_box(
                slide,
                lyric_left,
                body_top,
                lyric_w,
                body_h,
                group[0][0],
                group[0][1],
                typography=typo,
            )
        elif len(group) == 1:
            _add_hymn_lyric_box(
                slide,
                0,
                0,
                full_w,
                int(SLIDE_HEIGHT),
                group[0][0],
                group[0][1],
                typography=typo,
            )
        elif group_i == 0:
            w = SLIDE_WIDTH - 2 * MARGIN_SIDE
            title_box = slide.shapes.add_textbox(MARGIN_SIDE, _HYMN_TITLE_TOP, w, HYMN_TITLE_BOX_H)
            tft = title_box.text_frame
            _prep_tf(tft)
            tft.clear()
            title_pt = max(typo.title_pt, _HYMN_REF_TITLE_PT)
            _apply_hymn_song_title(
                tft.paragraphs[0],
                title,
                title_pt=title_pt,
                title_align=_pp_align(typo.title_align),
            )
            _add_hymn_lyric_box(
                slide,
                0,
                int(_HYMN_DUAL_TOP_FIRST),
                full_w,
                int(_HYMN_DUAL_FIRST_BOX_H),
                group[0][0],
                group[0][1],
                typography=typo,
            )
            _add_hymn_lyric_box(
                slide,
                0,
                int(_HYMN_DUAL_BOTTOM_FIRST),
                full_w,
                int(_HYMN_DUAL_FIRST_BOX_H),
                group[1][0],
                group[1][1],
                typography=typo,
                italic_body=_dual_second_verse_italic(group, 1),
            )
        else:
            _add_hymn_lyric_box(
                slide,
                0,
                int(_HYMN_DUAL_TOP_CONT),
                full_w,
                int(_HYMN_DUAL_CONT_BOX_H),
                group[0][0],
                group[0][1],
                typography=typo,
            )
            _add_hymn_lyric_box(
                slide,
                0,
                int(_HYMN_DUAL_BOTTOM_CONT),
                full_w,
                int(_HYMN_DUAL_CONT_BOX_H),
                group[1][0],
                group[1][1],
                typography=typo,
                italic_body=_dual_second_verse_italic(group, 1),
            )

        _add_hymn_footer(slide, footer_section)


def _add_hymn_lyric_slides(
    prs: Presentation,
    footer_section: str,
    hymn_title: str,
    lyrics: str,
    theme: SlideTheme,
    *,
    hymn_typography: Optional[Mapping[str, Any]] = None,
    section: str = "",
    hymn_lyrics_layout: str = "dual",
) -> None:
    del theme  # hymn slides use fixed projector palette, not liturgical theme
    if _normalize_hymn_lyrics_layout(hymn_lyrics_layout) == "dual":
        _add_hymn_lyric_slides_dual(
            prs,
            footer_section,
            hymn_title,
            lyrics,
            hymn_typography=hymn_typography,
            section=section,
        )
        return
    _add_hymn_lyric_slides_single(
        prs,
        footer_section,
        hymn_title,
        lyrics,
        hymn_typography=hymn_typography,
        section=section,
    )


def _try_library_hymn(
    prs: Presentation,
    section: str,
    hymn_id: str,
    footer: str,
    theme: SlideTheme,
    *,
    hymn_typography: Optional[Mapping[str, Any]] = None,
    hymn_lyric_overrides: Optional[Mapping[str, Any]] = None,
    hymn_lyrics_layout: str = "dual",
    hymn_layout_overrides: Optional[Mapping[str, Any]] = None,
) -> bool:
    h = get_hymn(section, hymn_id)
    if not h:
        return False
    title = str(h.get("title") or "Hymn")
    library_lyrics = str(h.get("lyrics") or "")
    lyrics = library_lyrics
    if hymn_lyric_overrides:
        sec_block = hymn_lyric_overrides.get(section)
        if isinstance(sec_block, Mapping):
            ov = sec_block.get(hymn_id) or sec_block.get(str(hymn_id))
            if ov:
                lyrics = pick_hymn_lyrics_for_slides(library_lyrics, str(ov))
    layout = hymn_lyrics_layout
    if hymn_layout_overrides:
        layout_block = hymn_layout_overrides.get(section)
        if isinstance(layout_block, Mapping):
            ov_layout = layout_block.get(hymn_id) or layout_block.get(str(hymn_id))
            if ov_layout:
                layout = str(ov_layout)
    _add_hymn_lyric_slides(
        prs,
        footer,
        title,
        lyrics,
        theme,
        hymn_typography=hymn_typography,
        section=section,
        hymn_lyrics_layout=layout,
    )
    return True


def chunk_plain_text(text: str, limit: int = _MAX_CHARS_READING) -> List[str]:
    if not (text or "").strip():
        return []
    norm = " ".join(text.split())
    if len(norm) <= limit:
        return [norm]
    sentences = re.split(r"(?<=[.!?])\s+", norm)
    out: List[str] = []
    buf = ""
    for s in sentences:
        w = s.strip()
        if not w:
            continue
        spacer = " " if buf else ""
        if len(buf) + len(spacer) + len(w) <= limit:
            buf += spacer + w
        else:
            if buf:
                out.append(buf.strip())
            if len(w) <= limit:
                buf = w
            else:
                for i in range(0, len(w), limit):
                    piece = w[i : i + limit].strip()
                    if piece:
                        out.append(piece)
                buf = ""
    if buf:
        out.append(buf.strip())
    return out if out else [norm[:limit]]


def _paragraphs(tf, *, size_pt, color, bold=False):
    tf.clear()
    p = tf.paragraphs[0]
    _style_para(p, size_pt=size_pt, color=color, bold=bold)
    p.alignment = PP_ALIGN.CENTER


def _fill_multipara(tf, text: str, *, size_pt: int, color: RGBColor):
    _fill_centered_box(tf, text, size_pt=size_pt, color=color)


def _add_reading_block(
    prs: Presentation,
    *,
    section: str,
    reference: str,
    body: str,
    unavailable_note: str,
    lotw_banner: bool,
    footer_tag: str,
    theme: SlideTheme,
) -> None:
    ref = (reference or "").strip() or "—"
    body = strip_reading_verse_markers((body or "").strip())

    def one_slide(head: str, sub: str, main: str):
        slide = prs.slides.add_slide(_layout_blank(prs))
        _set_slide_bg(slide, theme.bg)
        _apply_slide_branding(slide, theme)
        lx, top, w = MARGIN_SIDE, _content_top(), SLIDE_WIDTH - 2 * MARGIN_SIDE
        title_h = Inches(1.12) if lotw_banner else Inches(0.92)
        title_box = slide.shapes.add_textbox(lx, top, w, title_h)
        tf_t = title_box.text_frame
        _prep_tf(tf_t)
        tf_t.clear()
        if lotw_banner:
            p0 = tf_t.paragraphs[0]
            p0.text = "Liturgy of the Word"
            _style_para(p0, size_pt=_SLIDE_TEXT_PT, color=theme.emphasis, bold=True)
            p0.alignment = PP_ALIGN.CENTER
            p1 = tf_t.add_paragraph()
            p1.text = head if "continued" in head.lower() else f"{section} ({ref})"
            _style_para(p1, size_pt=_SLIDE_TEXT_PT, color=theme.muted, bold=False)
            p1.alignment = PP_ALIGN.CENTER
        else:
            _paragraphs(tf_t, size_pt=_SLIDE_TEXT_PT, color=theme.emphasis, bold=True)
            tf_t.paragraphs[0].text = head

        sub_top = top + title_h + Inches(0.06)
        sub_h = Inches(0.48)
        sub_box = slide.shapes.add_textbox(lx, sub_top, w, sub_h)
        _prep_tf(sub_box.text_frame)
        _paragraphs(sub_box.text_frame, size_pt=_SLIDE_TEXT_PT, color=theme.muted)
        sub_box.text_frame.paragraphs[0].text = sub

        body_top = sub_top + sub_h + Inches(0.12)
        body_h = SLIDE_HEIGHT - body_top - Inches(1.0)
        bsh = slide.shapes.add_textbox(lx, body_top, w, body_h)
        _prep_tf(bsh.text_frame)
        bsh.text_frame.vertical_anchor = BODY_VERTICAL_ANCHOR
        _paragraphs(bsh.text_frame, size_pt=_SLIDE_TEXT_PT, color=theme.primary)
        _fill_multipara(bsh.text_frame, main, size_pt=_SLIDE_TEXT_PT, color=theme.primary)
        _add_community_footer(slide, footer_tag, theme)

    if not body:
        one_slide(section, ref, unavailable_note)
        return
    chunks = chunk_plain_text(body)
    total = len(chunks)
    for i, chunk in enumerate(chunks):
        head = section if i == 0 else f"{section} (continued)"
        if lotw_banner:
            sub = "" if total <= 1 else f"Slide {i + 1} of {total}"
        else:
            sub = ref if total <= 1 else f"{ref}  ·  slide {i + 1} of {total}"
        one_slide(head, sub, chunk)


def _add_lotw_reading_slide(
    prs: Presentation,
    *,
    section: str,
    reference: str,
    full_text: str,
    theme: SlideTheme,
    reference_only: bool = False,
) -> None:
    """Liturgy of the Word: full text for psalm/canticle; citation only for other readings."""
    ref = (reference or "").strip() or "—"

    if reference_only:
        slide = prs.slides.add_slide(_layout_blank(prs))
        _set_slide_bg(slide, theme.bg)
        _apply_slide_branding(slide, theme)
        _add_section_title(slide, "Liturgy of the Word", _SECTION_TITLE_PT_LARGE)
        lx, top, w = MARGIN_SIDE, _content_top(), SLIDE_WIDTH - 2 * MARGIN_SIDE
        body_h = SLIDE_HEIGHT - top - CONTENT_BOTTOM_GAP
        box = slide.shapes.add_textbox(lx, top, w, body_h)
        tf = box.text_frame
        _prep_tf(tf)
        tf.clear()
        tf.vertical_anchor = BODY_VERTICAL_ANCHOR

        p1 = tf.paragraphs[0]
        p1.text = section
        _style_para(p1, size_pt=_SLIDE_TEXT_PT, color=theme.emphasis, bold=True)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(16)

        p2 = tf.add_paragraph()
        p2.text = ref
        _style_para(p2, size_pt=_SLIDE_TEXT_PT, color=theme.primary, bold=False)
        p2.alignment = PP_ALIGN.CENTER

        _add_community_footer(slide, "Liturgy of the Word", theme)
        return

    body = (full_text or "").strip()
    chunks = chunk_plain_text(body, limit=220) if body else []
    total = max(1, len(chunks))

    for i in range(total):
        chunk = chunks[i] if chunks else ""
        slide = prs.slides.add_slide(_layout_blank(prs))
        _set_slide_bg(slide, theme.bg)
        _apply_slide_branding(slide, theme)
        _add_section_title(slide, "Liturgy of the Word", _SECTION_TITLE_PT_LARGE)
        lx, top, w = MARGIN_SIDE, _content_top(), SLIDE_WIDTH - 2 * MARGIN_SIDE
        body_h = SLIDE_HEIGHT - top - CONTENT_BOTTOM_GAP
        box = slide.shapes.add_textbox(lx, top, w, body_h)
        tf = box.text_frame
        _prep_tf(tf)
        tf.clear()
        tf.vertical_anchor = BODY_VERTICAL_ANCHOR

        head = section if i == 0 else f"{section} (continued)"
        p1 = tf.paragraphs[0]
        p1.text = f"{head}\n({ref})"
        _style_para(p1, size_pt=_SLIDE_TEXT_PT, color=theme.emphasis, bold=True)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(12)

        if total > 1:
            p_cnt = tf.add_paragraph()
            p_cnt.text = f"Slide {i + 1} of {total}"
            _style_para(p_cnt, size_pt=_SLIDE_TEXT_PT, color=theme.muted, bold=False)
            p_cnt.alignment = PP_ALIGN.CENTER
            p_cnt.space_after = Pt(10)

        if chunk:
            p2 = tf.add_paragraph()
            p2.text = chunk
            _style_para(p2, size_pt=_SLIDE_TEXT_PT, color=theme.primary, bold=False)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_after = Pt(8)

        foot = "Liturgy of the Word" if total == 1 else f"Liturgy of the Word ({i + 1}/{total})"
        _add_community_footer(slide, foot, theme)


def _add_full_bleed_png_slides(prs: Presentation, paths: List[Optional[Path]]) -> None:
    for raw in paths or []:
        if not raw:
            continue
        p = Path(raw)
        if not p.is_file():
            continue
        slide = prs.slides.add_slide(_layout_blank(prs))
        slide.shapes.add_picture(
            str(p.resolve()),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )


_COLLECTION_CURRENCY_SPECS: dict[str, dict[str, Any]] = {
    "PHP": {"symbol": "₱", "prefix": True, "space": False},
    "KRW": {"symbol": "₩", "prefix": True, "space": False},
    "MYR": {"symbol": "RM", "prefix": True, "space": True},
}


def _normalize_collection_currency(code: str) -> str:
    cur = (code or "PHP").strip().upper()
    return cur if cur in _COLLECTION_CURRENCY_SPECS else "PHP"


def _format_collection_amount(amount: str, currency: str = "PHP") -> str:
    """Apply currency symbol placement for the Mass Collection slide."""
    raw = (amount or "").strip()
    if not raw:
        return ""
    cur = _normalize_collection_currency(currency)
    spec = _COLLECTION_CURRENCY_SPECS[cur]
    sym = spec["symbol"]
    gap = " " if spec.get("space") else ""
    if raw.startswith(sym) or (gap and raw.startswith(sym + gap)):
        return raw
    digits = re.sub(r"[^\d]", "", raw)
    if not digits:
        return raw
    formatted_num = f"{int(digits):,}"
    if spec.get("prefix", True):
        return f"{sym}{gap}{formatted_num}"
    return f"{formatted_num}{gap}{sym}"


def _add_mass_collection_slide(
    prs: Presentation,
    theme: SlideTheme,
    *,
    amount: str,
    date_label: str,
    currency: str = "PHP",
) -> None:
    """Title at top, amount centered, date + thank-you above the deck footer."""
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, theme.bg)
    _apply_slide_branding(slide, theme)

    lx = MARGIN_SIDE
    w = SLIDE_WIDTH - 2 * MARGIN_SIDE
    title_top = _content_top()
    title_h = COLLECTION_TITLE_H

    _add_section_title(slide, "Mass Collection", _SECTION_TITLE_PT_LARGE)

    foot_top = SLIDE_HEIGHT - COLLECTION_FOOTER_TOP
    mid_top = title_top + title_h
    mid_h = foot_top - mid_top - Inches(0.2)
    formatted_amount = _format_collection_amount(amount, currency)
    amount_text = formatted_amount or "(Enter collection amount in Mass Builder.)"
    mid_box = slide.shapes.add_textbox(lx, mid_top, w, mid_h)
    mid_tf = mid_box.text_frame
    _prep_tf(mid_tf)
    mid_tf.clear()
    mid_tf.vertical_anchor = BODY_VERTICAL_ANCHOR
    mid_p = mid_tf.paragraphs[0]
    mid_p.text = amount_text
    _style_para(mid_p, size_pt=_SLIDE_TEXT_PT, color=theme.primary, bold=True)
    mid_p.alignment = PP_ALIGN.CENTER

    foot_box = slide.shapes.add_textbox(lx, foot_top, w, Inches(0.95))
    foot_tf = foot_box.text_frame
    _prep_tf(foot_tf)
    foot_tf.clear()
    foot_lines: List[str] = []
    if (date_label or "").strip():
        foot_lines.append(date_label.strip())
    foot_lines.append("Thank you for your generosity.")
    first_foot = True
    for line in foot_lines:
        fp = foot_tf.paragraphs[0] if first_foot else foot_tf.add_paragraph()
        first_foot = False
        fp.text = line
        _style_para(fp, size_pt=_FOOTER_PT + 2, color=theme.muted, bold=False)
        fp.alignment = PP_ALIGN.CENTER
        fp.space_after = Pt(2)

    _add_community_footer(slide, "Mass Collection", theme)


def _add_food_sponsors_slide(prs: Presentation, theme: SlideTheme, sponsors: List[str]) -> None:
    names = [(s or "").strip() for s in (sponsors or [])]
    names = [n for n in names if n]
    if not names:
        return
    lines: List[str] = ["The community thanks our food sponsors."]
    for ss in names:
        lines.append(f"• {ss}")
    _add_marked_slide(
        prs, "Food Sponsors", "\n".join(lines), theme,
        title="Food Sponsors", title_pt=_SECTION_TITLE_PT_LARGE,
    )


def generate_mass_ppt(
    title: str,
    gospel_reference: str,
    gospel_quote: str,
    season: str,
    lectionary_cycle: str,
    celebrant: str,
    date: str,
    *,
    co_celebrant: str = "",
    gospel_full_text: str = "",
    first_reading_ref: str = "",
    first_reading_text: str = "",
    psalm_ref: str = "",
    psalm_text: str = "",
    second_reading_ref: str = "",
    second_reading_text: str = "",
    quote_attribution=None,
    quote_max_chars: int = 400,
    liturgical_color: Optional[Mapping[str, Any]] = None,
    custom_theme: Optional[Mapping[str, Any]] = None,
    song_selections: Optional[Mapping[str, Any]] = None,
    output_stem: str = "mass_presentation",
    liturgical_poster_png: Optional[Path] = None,
    divider_poster_png: Optional[Path] = None,
    lotw_poster: str = _LOTW_POSTER_DEFAULT,
    lote_poster: str = _LOTE_POSTER_DEFAULT,
    announcement_image_paths: Optional[List[Path]] = None,
    mass_collection_amount: str = "",
    mass_collection_date_label: str = "",
    mass_collection_currency: str = "PHP",
    food_sponsors: Optional[List[str]] = None,
    hymn_typography: Optional[Mapping[str, Any]] = None,
    include_church_logo: bool = False,
    include_church_name: bool = False,
    include_footer: bool = False,
    footer_brand: str = "",
    hymn_lyric_overrides: Optional[Mapping[str, Any]] = None,
    gospel_acclamation_verse: str = "",
    creed_choice: str = "nicene",
    our_father_choice: str = "english",
    hymn_lyrics_layout: str = "dual",
    hymn_layout_overrides: Optional[Mapping[str, Any]] = None,
    video_replacements: Optional[Mapping[str, Any]] = None,
) -> tuple[int, Path]:
    global _ACTIVE_FONT, _ACTIVE_THEME, _ACTIVE_DECK_THEME_ID, _deck_branding
    _deck_branding = DeckBrandingOptions(
        include_logo=bool(include_church_logo),
        include_name=bool(include_church_name),
        include_footer=bool(include_footer) or bool(str(footer_brand or "").strip()),
        footer_brand=str(footer_brand or "").strip(),
    )
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    theme_id = _resolve_deck_theme_id(custom_theme)
    _ACTIVE_DECK_THEME_ID = theme_id
    theme = _build_slide_theme(liturgical_color, custom_theme)
    _ACTIVE_FONT = theme.font_name
    _ACTIVE_THEME = theme

    lotw_poster_path = _resolve_bundled_poster(
        lotw_poster, _LOTW_POSTER_IDS, _LOTW_POSTER_DEFAULT
    )
    lote_poster_path = _resolve_bundled_poster(
        lote_poster, _LOTE_POSTER_IDS, _LOTE_POSTER_DEFAULT
    )

    g_line = (gospel_quote or "").strip()
    if quote_max_chars and len(g_line) > quote_max_chars:
        g_line = g_line[: quote_max_chars - 1].rstrip() + "\u2026"

    unavail = (
        "Full text was not loaded from bible.usccb.org. "
        "Open today’s readings for this date and paste if needed."
    )

    ctx = dict(
        celebrant=celebrant,
        co_celebrant=co_celebrant,
        date=date,
        season=season,
        lectionary_cycle=lectionary_cycle,
        gospel_reference=gospel_reference or "",
        gospel_quote=g_line,
        quote_max_chars=quote_max_chars,
        theme=theme,
        mass_title=title,
        background_image_path=liturgical_poster_png,
        divider_poster_path=divider_poster_png,
    )

    sel = song_selections or {}
    videos: dict[str, Path] = {}
    if isinstance(video_replacements, Mapping):
        for key, val in video_replacements.items():
            k = str(key or "").strip().lower()
            if not k:
                continue
            if isinstance(val, Path):
                path = val
            else:
                path = Path(str(val or "").strip())
            if path.is_file():
                videos[k] = path

    def _use_video(slot: str, label: str) -> bool:
        path = videos.get(slot)
        return bool(path and _add_video_replacement_slide(prs, path, title=label, theme=theme))

    # --- Pre-Mass (reference deck slide) ---
    _add_pre_mass_slide(prs, theme)

    _add_divider_cover(prs, **ctx)

    ent_id = str(sel.get("entrance") or "").strip()
    if _use_video("entrance", "Entrance"):
        pass
    elif not ent_id or not _try_library_hymn(
        prs, "entrance", ent_id, "Entrance", theme, hymn_typography=hymn_typography, hymn_lyric_overrides=hymn_lyric_overrides, hymn_lyrics_layout=hymn_lyrics_layout, hymn_layout_overrides=hymn_layout_overrides
    ):
        _add_marked_slide(
            prs,
            "Entrance",
            "No Entrance hymn lyrics were selected. Choose one Entrance song in Mass Flow or save lyrics in Lyrics Studio before generating.",
            theme,
        )
    _add_divider_cover(prs, **ctx)

    # --- Introductory Rites ---
    if not _clone_master_section(prs, "introductory_rites", theme, "Introductory Rites"):
        _add_marked_slide(prs, "Introductory Rites", GFCC.SIGN_CROSS, theme)
    _add_penitential_act_slides(prs, theme)
    if not _use_video("kyrie", "Kyrie"):
        _add_kyrie_slide(prs, theme)
    if not _use_video("gloria", "Gloria"):
        _add_gloria_slides(prs, theme)
    if not _clone_master_section(prs, "lotw_prayer", theme, "Liturgy of the Word"):
        _add_marked_slide(
            prs, "Liturgy of the Word", GFCC.OPENING_PRAYER, theme,
            title="Liturgy of the Word", title_pt=_SECTION_TITLE_PT_LARGE,
        )

    # --- Liturgy of the Word ---
    _add_lotw_title_slide(prs, theme, lotw_poster_path)

    # --- Readings: citation-only cards cloned 1:1 from the master template.
    # Scripture bodies are intentionally dropped (matches the template card).
    # Reading cards: "Liturgy of the Word" header + section label render amber;
    # the scripture citation (verse number) stays white. The psalm keeps section
    # + verse number amber, with only the responsory refrain in white.
    def _color_reading_card(slide, label: str) -> None:
        _color_shapes(
            slide, _ACTIVE_THEME.emphasis,
            exact=["liturgy of the word"], top_below=Inches(2),
        )
        _color_shapes(slide, _ACTIVE_THEME.emphasis, exact=[label])

    first_ref = (first_reading_ref or "—").strip() or "—"

    def _inject_first(slide, _i):
        _set_run_text_keep_format(slide, _TPL_FIRST_CITATION, first_ref)
        _color_reading_card(slide, "First Reading")

    if not _clone_master_section(
        prs, "first_reading", theme, "Liturgy of the Word", mutate=_inject_first,
    ):
        _add_lotw_reading_slide(
            prs, section="First Reading", reference=first_ref,
            full_text="", theme=theme, reference_only=True,
        )

    psalm_section = (responsorial_section_title(psalm_ref or "") or _TPL_PSALM_SECTION).strip()
    psalm_ref_clean = (psalm_ref or "").strip()
    psalm_response = _extract_psalm_response(psalm_text)

    def _inject_psalm(slide, _i):
        _set_run_text_keep_format(slide, _TPL_PSALM_SECTION, psalm_section)
        if psalm_ref_clean:
            _set_run_text_keep_format(slide, _TPL_PSALM_REF, psalm_ref_clean)
        _set_run_text_keep_format(slide, _TPL_PSALM_ANTIPHON, psalm_response)
        _color_shapes(
            slide, _ACTIVE_THEME.emphasis,
            exact=["liturgy of the word"], top_below=Inches(2),
        )
        # Section title + verse number (the "Responsorial Psalm (…)" block) render amber.
        _color_shapes(
            slide, _ACTIVE_THEME.emphasis,
            contains="responsorial", top_below=Inches(5),
        )

    if not _clone_master_section(prs, "psalm", theme, "Liturgy of the Word", mutate=_inject_psalm):
        _add_lotw_reading_slide(
            prs, section=psalm_section, reference=psalm_ref_clean or "—",
            full_text="", theme=theme, reference_only=True,
        )

    if (second_reading_ref or "").strip():
        second_ref = second_reading_ref.strip()

        def _inject_second(slide, _i):
            _set_run_text_keep_format(slide, _TPL_SECOND_CITATION, second_ref)
            _color_reading_card(slide, "Second Reading")

        if not _clone_master_section(
            prs, "second_reading", theme, "Liturgy of the Word", mutate=_inject_second,
        ):
            _add_lotw_reading_slide(
                prs, section="Second Reading", reference=second_ref,
                full_text="", theme=theme, reference_only=True,
            )

    _add_gospel_acclamation_slides(
        prs,
        theme,
        gospel_reference=gospel_reference or "",
        gospel_acclamation_verse=gospel_acclamation_verse or "",
    )

    if not _clone_master_section(
        prs, "gospel_acclamation_end", theme, "Gospel Acclamation",
        mutate=lambda s, _i: _apply_rite_slide_title_typography(
            s, "Gospel Acclamation", size_pt=_SECTION_TITLE_PT_LARGE
        ),
    ):
        _add_marked_slide(prs, "Gospel Acclamation", GFCC.GOSPEL_END, theme)
    _add_divider_cover(prs, **ctx)

    # --- Creed (Nicene or Apostles' — never both) ---
    _add_creed_slides(prs, theme, creed_choice=creed_choice)
    _add_divider_cover(prs, **ctx)

    # --- Prayer of the Faithful ---
    if not _clone_master_section(prs, "prayer_faithful", theme, "Prayer of the Faithful"):
        _add_marked_slide(
            prs, "Prayer of the Faithful", GFCC.PRAYER_FAITHFUL_1, theme,
            title="Prayer of the Faithful", title_pt=_SECTION_TITLE_PT_LARGE,
        )
        _add_marked_slide(
            prs, "Prayer of the Faithful", GFCC.PRAYER_FAITHFUL_2, theme,
            title="Prayer of the Faithful", title_pt=_SECTION_TITLE_PT_LARGE,
        )
    _add_divider_cover(prs, **ctx)

    # --- Liturgy of the Eucharist ---
    off_id = str(sel.get("offertory") or "").strip()
    if _use_video("offertory", "Offertory"):
        pass
    elif not off_id or not _try_library_hymn(
        prs, "offertory", off_id, "Offertory", theme, hymn_typography=hymn_typography, hymn_lyric_overrides=hymn_lyric_overrides, hymn_lyrics_layout=hymn_lyrics_layout, hymn_layout_overrides=hymn_layout_overrides
    ):
        _add_marked_slide(
            prs,
            "Offertory",
            "No Offertory hymn lyrics were selected. Choose one Offertory song in Mass Flow or save lyrics in Lyrics Studio before generating.",
            theme,
        )
    _add_lote_poster_slide(prs, theme, lote_poster_path)
    if not _clone_master_section(prs, "lote_pray_brethren", theme, "Liturgy of the Eucharist"):
        _add_marked_slide(prs, "Liturgy of the Eucharist", GFCC.PRAY_BRETHREN, theme)
    _add_lote_poster_slide(prs, theme, lote_poster_path)
    if not _clone_master_section(prs, "preface_dialogue", theme, "Liturgy of the Eucharist"):
        _add_marked_slide(prs, "Liturgy of the Eucharist", GFCC.PREFACE_DIALOGUE, theme)
        _add_marked_slide(prs, "Liturgy of the Eucharist", GFCC.PREFACE_ACCLAIM, theme)
    _add_lote_poster_slide(prs, theme, lote_poster_path)
    if not _use_video("sanctus", "Sanctus"):
        if not _clone_master_section(prs, "sanctus", theme, "Sanctus"):
            _add_marked_chunked(prs, "Sanctus", get_prayer("holy_holy"), theme)
    _add_lote_poster_slide(prs, theme, lote_poster_path)
    if not _clone_master_section(prs, "mystery_of_faith", theme, "The Eucharistic Prayer"):
        _add_marked_slide(
            prs, "The Eucharistic Prayer", get_prayer("mystery_of_faith"), theme,
            title="The Mystery of Faith", title_pt=_SECTION_TITLE_PT_LARGE,
        )
    _add_lote_poster_slide(prs, theme, lote_poster_path)
    if not _clone_master_section(prs, "great_amen", theme, "Great Amen"):
        _add_marked_slide(prs, "Great Amen", GFCC.GREAT_AMEN, theme)
    _of_choice = _normalize_our_father_choice(our_father_choice)
    if not _use_video("our_father", "Our Father"):
        if not _add_our_father_from_deck(prs, theme, _of_choice):
            _add_our_father_slide(
                prs,
                "Our Father",
                get_our_father(_of_choice),
                theme,
                title=_OUR_FATHER_TITLES.get(_of_choice, "Our Father"),
            )
    _add_divider_cover(prs, **ctx)
    _add_sign_of_peace_slide(prs, theme)
    if not _use_video("lamb_of_god", "Lamb of God"):
        _add_lamb_of_god_slide(prs, theme)
    if not _clone_master_section(prs, "communion_rite", theme, "The Communion Rite"):
        _add_marked_slide(prs, "The Communion Rite", GFCC.COMMUNION_DIALOGUE, theme)
    _add_divider_cover(prs, **ctx)
    c1 = str(sel.get("communion_1") or "").strip()
    c2 = str(sel.get("communion_2") or "").strip()
    comm_ok = False
    if _use_video("communion_1", "Communion (1)"):
        comm_ok = True
    elif c1 and _try_library_hymn(
        prs, "communion", c1, "Communion (1)", theme, hymn_typography=hymn_typography, hymn_lyric_overrides=hymn_lyric_overrides, hymn_lyrics_layout=hymn_lyrics_layout, hymn_layout_overrides=hymn_layout_overrides
    ):
        comm_ok = True
    if _use_video("communion_2", "Communion (2)"):
        comm_ok = True
    elif c2 and _try_library_hymn(
        prs, "communion", c2, "Communion (2)", theme, hymn_typography=hymn_typography, hymn_lyric_overrides=hymn_lyric_overrides, hymn_lyrics_layout=hymn_lyrics_layout, hymn_layout_overrides=hymn_layout_overrides
    ):
        comm_ok = True
    if not comm_ok:
        _add_marked_slide(
            prs,
            "Communion",
            "No Communion hymn lyrics were selected. Choose up to two Communion songs in Mass Flow or save lyrics in Lyrics Studio before generating.",
            theme,
        )
    med_id = str(sel.get("meditation") or "").strip()
    if med_id:
        _try_library_hymn(
            prs, "meditation", med_id, "Meditation", theme, hymn_typography=hymn_typography, hymn_lyric_overrides=hymn_lyric_overrides, hymn_lyrics_layout=hymn_lyrics_layout, hymn_layout_overrides=hymn_layout_overrides
        )
    extra_sections = sel.get("extra_sections") or []
    if isinstance(extra_sections, list):
        for item in extra_sections:
            if not isinstance(item, dict):
                continue
            label = str(item.get("label") or "Custom").strip() or "Custom"
            song_id = str(item.get("song_id") or "").strip()
            if song_id:
                _try_library_hymn(
                    prs,
                    "meditation",
                    song_id,
                    label,
                    theme,
                    hymn_typography=hymn_typography,
                    hymn_lyric_overrides=hymn_lyric_overrides,
                    hymn_lyrics_layout=hymn_lyrics_layout,
                    hymn_layout_overrides=hymn_layout_overrides,
                )
    if not _clone_master_section(prs, "post_communion", theme, "The Communion Rite"):
        _add_marked_slide(prs, "The Communion Rite", GFCC.POST_COMMUNION, theme)
    _add_divider_cover(prs, **ctx)

    # --- Stewardship, sponsors, announcements (before final blessing) ---
    # Template order: Welcoming Newcomers -> Mass Collection -> Food Sponsors ->
    # Confession. The "Church Announcements" title and "Updates" slides were removed
    # from the template, so they are no longer generated.
    ann_paths: List[Optional[Path]] = list(announcement_image_paths or [])
    if not ann_paths:
        church_name = (get_community_name() or "").strip() or _TPL_CHURCH_NAME
        if not _clone_master_section(
            prs, "welcoming_newcomers", theme, "Welcoming Newcomers",
            mutate=lambda s, _i: _set_run_text_keep_format(s, _TPL_CHURCH_NAME, church_name),
        ):
            _add_marked_slide(prs, "Welcoming Newcomers", GFCC.WELCOME_NEWCOMERS, theme)

    collection_amount = _format_collection_amount(
        mass_collection_amount or "", mass_collection_currency or "PHP"
    )
    collection_date = (mass_collection_date_label or "").strip()

    def _inject_collection(slide, _i):
        _set_run_text_keep_format(slide, _TPL_COLLECTION_AMOUNT, collection_amount)
        _set_run_text_keep_format(slide, _TPL_COLLECTION_DATE, collection_date)
        _color_shapes(
            slide, _ACTIVE_THEME.emphasis,
            exact=["mass collection"], top_below=Inches(2),
        )

    if not _clone_master_section(prs, "mass_collection", theme, "Mass Collection", mutate=_inject_collection):
        _add_mass_collection_slide(
            prs,
            theme,
            amount=mass_collection_amount or "",
            date_label=mass_collection_date_label or "",
            currency=mass_collection_currency or "PHP",
        )

    _sponsor_names = [(s or "").strip().upper() for s in (food_sponsors or []) if (s or "").strip()]
    if _sponsor_names:
        if not _add_food_sponsors_from_template(prs, theme, _sponsor_names):
            _add_food_sponsors_slide(prs, theme, _sponsor_names)
    if ann_paths:
        _add_full_bleed_png_slides(prs, ann_paths)
    elif not _clone_master_section(prs, "confession", theme, "Announcements"):
        _add_marked_slide(
            prs,
            "Announcements",
            "The Lord never tires of forgiving us; we are the ones who tire of seeking his mercy.\n— Pope Francis",
            theme,
            title="Sacrament of Confession",
            title_pt=_SECTION_TITLE_PT_LARGE,
        )

    if not _clone_master_section(prs, "final_blessing", theme, "Final Blessing"):
        _add_marked_slide(prs, "Final Blessing", GFCC.FINAL_BLESSING, theme)
    rec_id = str(sel.get("recessional") or "").strip()
    if _use_video("recessional", "Recessional"):
        pass
    elif not rec_id or not _try_library_hymn(
        prs, "recessional", rec_id, "Recessional", theme, hymn_typography=hymn_typography, hymn_lyric_overrides=hymn_lyric_overrides, hymn_lyrics_layout=hymn_lyrics_layout, hymn_layout_overrides=hymn_layout_overrides
    ):
        _add_marked_slide(
            prs,
            "Recessional",
            "No Recessional hymn lyrics were selected. Choose one Recessional song in Mass Flow or save lyrics in Lyrics Studio before generating.",
            theme,
        )
    _add_divider_cover(prs, **ctx)

    # Secondary reference decks (Gloria, Kyrie, Our Father, …) still use live
    # remapping. Baked Theme 2/3 masters skip recolor on clone; this pass only
    # needs to scrub leftovers from those secondary clones.
    _finalize_deck_text_colors(prs, theme)

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    stem = (output_stem or "mass_presentation").strip() or "mass_presentation"
    out = _OUTPUT_DIR / f"{stem}.pptx"
    n_slides = len(prs.slides)
    prs.save(out)
    print(f"✅ PowerPoint created: {out} ({n_slides} slides)")
    return n_slides, out
