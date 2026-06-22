"""Extract Our Father slide sets into standalone per-language reference decks."""
from __future__ import annotations

import copy
from pptx import Presentation


def keep_slides(src: str, dst: str, keep_1based) -> None:
    prs = Presentation(src)
    keep = {i - 1 for i in keep_1based}
    sld_lst = prs.slides._sldIdLst
    for i, sld_id in enumerate(list(sld_lst)):
        if i not in keep:
            sld_lst.remove(sld_id)
    prs.save(dst)
    out = Presentation(dst)
    print(f"{dst}: {len(out.slides.__iter__.__self__._sldIdLst)} slide refs")


if __name__ == "__main__":
    # English: new master LFTemplate1.pptx slides 45,46,47 (Our Father / embolism / doxology)
    keep_slides(
        "data/reference/LFTemplate1.pptx",
        "data/reference/our_father_english.pptx",
        [45, 46, 47],
    )
    # Tagalog: previous master LiturgyFlowTemplate1.pptx slides 46,47,48,49 (Ama Namin x4)
    keep_slides(
        "data/reference/LiturgyFlowTemplate1.pptx",
        "data/reference/our_father_tagalog.pptx",
        [46, 47, 48, 49],
    )
    print("done")
