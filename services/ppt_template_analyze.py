"""
Extract colors and typography hints from an uploaded PPTX for slide generation.

Heuristic (no external AI required): samples fills and fonts from early slides.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Optional


def _rgb_to_hex(rgb) -> Optional[str]:
    if rgb is None:
        return None
    try:
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
        return f"#{r:02x}{g:02x}{b:02x}"
    except (TypeError, ValueError, IndexError):
        return None


def analyze_pptx_theme(pptx_path: Path) -> dict[str, Any]:
    """Return a dict usable as ``custom_theme`` merge target plus diagnostics."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
    except ImportError:
        return {"ok": False, "error": "python-pptx is required.", "custom_theme": None}

    p = Path(pptx_path)
    if not p.is_file():
        return {"ok": False, "error": "File not found.", "custom_theme": None}

    prs = Presentation(str(p))
    bg_samples: list[str] = []
    text_samples: list[str] = []
    font_names: list[str] = []

    for si, slide in enumerate(prs.slides[:8]):
        try:
            fi = slide.background.fill
            if fi.type and str(fi.type).endswith("SOLID"):
                rgb = fi.fore_color.rgb
                if isinstance(rgb, RGBColor):
                    hx = _rgb_to_hex((rgb[0], rgb[1], rgb[2]))
                    if hx:
                        bg_samples.append(hx)
        except Exception:
            pass
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_names.append(str(run.font.name).split(",")[0].strip())
                        if run.font.color and run.font.color.type:
                            try:
                                rgb = run.font.color.rgb
                                if isinstance(rgb, RGBColor):
                                    hx = _rgb_to_hex((rgb[0], rgb[1], rgb[2]))
                                    if hx:
                                        text_samples.append(hx)
                            except Exception:
                                pass
            except Exception:
                continue

    def _pick(dark_candidates: list[str], light_candidates: list[str], default: str) -> str:
        for bucket in (dark_candidates, light_candidates, bg_samples, text_samples):
            for hx in bucket:
                if re.match(r"^#[0-9a-fA-F]{6}$", hx):
                    return hx
        return default

    bg = bg_samples[0] if bg_samples else "#050907"
    primary = _pick(text_samples, [], "#ffffff")
    accent = text_samples[1] if len(text_samples) > 1 else "#ffb800"
    font = font_names[0] if font_names else "Poppins"

    custom_theme = {
        "bg": bg,
        "text": primary,
        "primary": accent,
        "accent": accent,
        "font": font,
        "template_source": p.name,
        "slides_sampled": min(len(prs.slides), 8),
    }

    font_preview = ", ".join(sorted(set(font_names))[:5]) or "—"
    return {
        "ok": True,
        "custom_theme": custom_theme,
        "notes": f"Sampled {custom_theme['slides_sampled']} slide(s); fonts: {font_preview}.",
    }
