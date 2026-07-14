"""Extract rite prayer/lyric text from the authored PPTX slides for Mass Builder preview."""

from __future__ import annotations

import logging
import re
from functools import lru_cache
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

_PROJECT_ROOT = Path(__file__).resolve().parent.parent
_REFERENCE = _PROJECT_ROOT / "data" / "reference"
_MASTER = _REFERENCE / "LFTemplate1.pptx"

# Master-template slide indices (see generators.powerpoint._MASTER_SLIDE).
_MASTER_SLIDES: dict[str, tuple[int, ...]] = {
    "penitential": (8, 9, 10),
    "kyrie": (11,),
    "gloria": (12, 13, 14, 15),
    "nicene_creed": (25, 26, 27, 28),
    "sanctus": (40,),
    "lamb_of_god": (50,),
}

_FOOTER_LABEL_RE = re.compile(
    r"^(?:"
    r"penitential act|kyrie eleison|gloria|sanctus|lamb of god|"
    r"our father|nicene creed|apostles'? creed|ama namin"
    r")"
    r"(?:\s*\(\d+\s*/\s*\d+\))?$",
    re.IGNORECASE,
)

_PARISH_BRAND_MARKERS = (
    "CATHOLIC COMMUNITY",
    "FILIPINO CATHOLIC",
    " CATHOLIC CHURCH",
    "PARISH COMMUNITY",
)
_HEADER_ZONE_FRAC = 0.14
_FOOTER_ZONE_FRAC = 0.90

_STATIC_FALLBACKS: dict[str, list[str]] = {
    "penitential::form-b": [
        "Have mercy on us, O Lord.\nFor we have sinned against you.\n"
        "Show us, O Lord, your mercy.\nAnd grant us your salvation."
    ],
    "penitential::form-c": [
        "You were sent to heal the contrite of heart:\nLord, have mercy.\n\n"
        "You came to call sinners:\nChrist, have mercy.\n\n"
        "You are seated at the right hand of the Father to intercede for us:\nLord, have mercy."
    ],
    "kyrie::greek": ["Kyrie eleison.\nChriste eleison.\nKyrie eleison."],
    "kyrie::latin": ["Kyrie eleison.\nChriste eleison.\nKyrie eleison."],
    "gloria::latin": [
        "Gloria in excelsis Deo\net in terra pax hominibus bonae voluntatis.\n"
        "Laudamus te, benedicimus te, adoramus te,\nglorificamus te, gratias agimus tibi\n"
        "propter magnam gloriam tuam,\nDomine Deus, Rex caelestis, Deus Pater omnipotens.\n\n"
        "Domine Fili unigenite, Iesu Christe,\nDomine Deus, Agnus Dei, Filius Patris,\n"
        "qui tollis peccata mundi, miserere nobis;\nqui tollis peccata mundi, suscipe deprecationem nostram.\n"
        "Qui sedes ad dexteram Patris, miserere nobis.\n\n"
        "Quoniam tu solus Sanctus, tu solus Dominus,\ntu solus Altissimus, Iesu Christe,\n"
        "cum Sancto Spiritu: in gloria Dei Patris.\nAmen."
    ],
    "lamb_of_god::latin": [
        "Agnus Dei,\nqui tollis peccata mundi,\nmiserere nobis.\n"
        "Agnus Dei,\nqui tollis peccata mundi,\nmiserere nobis.\n"
        "Agnus Dei,\nqui tollis peccata mundi,\ndona nobis pacem."
    ],
}


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        t = (para.text or "").replace("\u000b", "\n").strip()
        if t:
            lines.append(t)
    return "\n".join(lines)


def _is_parish_branding_text(text: str) -> bool:
    """Drop stacked parish/logo wordmarks baked into reference decks."""
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    if not lines or len(lines) > 8:
        return False
    joined = " ".join(lines).upper()
    if not any(marker in joined for marker in _PARISH_BRAND_MARKERS):
        return False
    # Branding blocks are short ALL-CAPS (or near-caps) label stacks.
    short = all(len(ln) <= 40 for ln in lines)
    caps_ish = sum(1 for ln in lines if ln.upper() == ln) >= max(1, len(lines) - 1)
    return short and caps_ish


def _iter_shape_texts(slide, *, slide_height: int | None = None) -> list[str]:
    chunks: list[str] = []
    height = int(slide_height or 0) or None

    def walk(shapes, abs_top: int = 0) -> None:
        for shape in shapes:
            st = getattr(shape, "shape_type", None)
            top = abs_top + int(getattr(shape, "top", 0) or 0)
            if st == MSO_SHAPE_TYPE.GROUP:
                # Parish logo / name stacks live in header chrome groups.
                if height and top < int(height * _HEADER_ZONE_FRAC):
                    continue
                try:
                    walk(shape.shapes, top)
                except Exception:
                    pass
                continue
            if height and top >= int(height * _FOOTER_ZONE_FRAC):
                continue
            text = _shape_text(shape)
            if not text or _is_parish_branding_text(text):
                continue
            chunks.append(text)

    walk(slide.shapes)
    return chunks


def _clean_slide_text(chunks: list[str]) -> str:
    kept: list[str] = []
    seen: set[str] = set()
    for raw in chunks:
        block = raw.strip()
        if not block or _is_parish_branding_text(block):
            continue
        lines = []
        for line in block.splitlines():
            s = line.strip()
            if not s:
                continue
            if _FOOTER_LABEL_RE.match(s):
                continue
            lines.append(s)
        cleaned = "\n".join(lines).strip()
        if not cleaned or _is_parish_branding_text(cleaned):
            continue
        key = re.sub(r"\s+", " ", cleaned).lower()
        if key in seen:
            continue
        seen.add(key)
        kept.append(cleaned)
    return "\n\n".join(kept).strip()


@lru_cache(maxsize=8)
def _load_presentation(path_str: str) -> Presentation | None:
    path = Path(path_str)
    if not path.is_file():
        return None
    try:
        return Presentation(str(path))
    except Exception:
        logger.exception("Failed to open rite preview deck %s", path)
        return None


def _slides_from_deck(path: Path, indices: tuple[int, ...] | None = None) -> list[str]:
    prs = _load_presentation(str(path.resolve()))
    if prs is None:
        return []
    slide_h = int(getattr(prs, "slide_height", 0) or 0)
    out: list[str] = []
    if indices is None:
        slide_list = list(prs.slides)
    else:
        slide_list = []
        for i in indices:
            if 0 <= i < len(prs.slides):
                slide_list.append(prs.slides[i])
    for slide in slide_list:
        text = _clean_slide_text(_iter_shape_texts(slide, slide_height=slide_h))
        if text:
            out.append(text)
    return out


def _our_father_slides(option: str) -> list[str]:
    decks = {
        "english": "our_father_english.pptx",
        "tagalog": "our_father_tagalog.pptx",
    }
    fname = decks.get(option)
    if fname:
        pages = _slides_from_deck(_REFERENCE / fname)
        if pages:
            return pages
    try:
        from services.prayer_service import get_our_father

        text = get_our_father(option)
        text = re.sub(r"<<[APD]>>", "", text).strip()
        return [text] if text else []
    except Exception:
        return []


def _resolve_slide_pages(section: str, option: str) -> list[str]:
    sec = (section or "").strip().lower()
    opt = (option or "").strip().lower() or "default"
    key = f"{sec}::{opt}"

    if key in _STATIC_FALLBACKS:
        return list(_STATIC_FALLBACKS[key])

    if sec == "penitential":
        # Authored deck is Confiteor (Form A). Forms B/C use static fallbacks above.
        return _slides_from_deck(_MASTER, _MASTER_SLIDES["penitential"])

    if sec == "creed":
        if opt == "apostles":
            pages = _slides_from_deck(_REFERENCE / "apostles_creed_slides.pptx")
            if pages:
                return pages
        pages = _slides_from_deck(_REFERENCE / "nicene_creed_slides.pptx")
        if pages:
            return pages
        return _slides_from_deck(_MASTER, _MASTER_SLIDES["nicene_creed"])

    if sec == "our_father":
        return _our_father_slides(opt)

    if sec == "kyrie":
        pages = _slides_from_deck(_REFERENCE / "kyrie_slide.pptx")
        if pages:
            return pages
        return _slides_from_deck(_MASTER, _MASTER_SLIDES["kyrie"])

    if sec == "gloria":
        pages = _slides_from_deck(_REFERENCE / "gloria_slides.pptx")
        if pages:
            return pages
        return _slides_from_deck(_MASTER, _MASTER_SLIDES["gloria"])

    if sec == "sanctus":
        return _slides_from_deck(_MASTER, _MASTER_SLIDES["sanctus"])

    if sec == "lamb_of_god":
        pages = _slides_from_deck(_REFERENCE / "lamb_of_god_slide.pptx")
        if pages:
            return pages
        return _slides_from_deck(_MASTER, _MASTER_SLIDES["lamb_of_god"])

    return []


def build_rite_slide_preview(section: str, option: str = "") -> dict[str, Any]:
    """Return joined + per-slide text for one rite option card."""
    sec = (section or "").strip().lower()
    opt = (option or "").strip().lower() or "default"
    slides = _resolve_slide_pages(sec, opt)
    text = "\n\n".join(slides).strip()
    return {
        "ok": bool(text),
        "section": sec,
        "option": opt,
        "slide_count": len(slides),
        "slides": slides,
        "text": text,
    }
