"""
Church Media Generator — minimal web UI + JSON API.

Run from project root:
  cd church_media_generator && uvicorn server:app --reload --host 127.0.0.1 --port 8000
"""

from __future__ import annotations

import io
import json
import logging
import base64
import os
import re
import shutil
import subprocess
import tempfile
import threading
import uuid
import zipfile
from pathlib import Path
from typing import Any, Optional

from fastapi import Depends, FastAPI, File, Header, HTTPException, Request, UploadFile
from fastapi.middleware.gzip import GZipMiddleware
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, Response
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel, Field, model_validator
from starlette.concurrency import run_in_threadpool
from urllib.parse import urlencode

from services.env_config import load_project_dotenv

load_project_dotenv()

logger = logging.getLogger("server")

from services.calendar_month import fetch_calendar_month
from services.catholic_news import fetch_catholic_headlines
from services.catholic_radio import radio_catalog

from pipeline import (
    GenerationResult,
    PreviewPayload,
    fetch_preview,
    generate_mass_media,
    regenerate_mass_pptx,
    refresh_all_song_sections,
    refresh_song_section,
)
from services.community_config import (
    LOGO_RELATIVE,
    get_celebrant_names,
    get_community_name,
    load_community,
    logo_file_absolute,
    update_community,
    uploads_dir,
)
from services.gospel_mood import gospel_moods_for_song
from services.hymn_library import get_hymn
from services.hymn_slide_preview import build_hymn_slide_preview
from services.lyrics_fetcher import fetch_and_store_for_selection
from services.ppt_preview_render import render_ppt_preview_pngs
from services.ppt_template_analyze import analyze_pptx_theme
from services.song_catalog import (
    catalog_for_api,
    catalog_lite_response,
    delete_catalog_song,
    find_catalog_row_by_id,
    import_song_rows,
    import_titles,
    save_lyrics_song,
    update_catalog_song,
)
from services.auth_config import auth_enabled, invite_contact_email, supabase_enabled
from services.api_security import (
    AuthSession,
    optional_session,
    register_security_middleware,
    require_approved_membership,
    require_session_when_auth,
    require_superadmin,
)
from services.membership_config import (
    can_edit_logo,
    is_superadmin_user,
    membership_payload,
    parish_name_is_locked,
)
from services.pending_submissions import submit_pending_priest, submit_pending_song
from services.choir_practice_shares import (
    create_practice_share,
    fetch_practice_share,
    get_practice_share_by_token,
    update_practice_share_lyrics,
    verify_practice_share_pin,
)
from services.practice_access import (
    check_pin_unlock_allowed,
    check_practice_fetch_allowed,
    check_practice_share_create_allowed,
    check_practice_token_allowed,
    ensure_practice_secret_configured,
    is_unlocked,
    issue_lead_token,
    issue_unlock_cookie,
    practice_no_store_headers,
    verify_lead_token,
)
from services.catalog_rate_limit import check_catalog_lyric_fetch_allowed
from services.hymn_normalized_store import (
    fetch_lyrics_from_normalized,
    fetch_song_from_normalized,
)
from services.lyric_audit import log_lyric_read
from services.media_ownership import register_owned_files, session_may_access_media
from services.runtime_config import is_production_runtime
from services.practice_qr import practice_qr_data_url
from services.user_church_context import get_church_profile_context
from services.readings_snapshot import readings_snapshot, warm_readings_for_date
from services.image_generation_quota import (
    quota_status_payload,
    reserve_daily_image_generation,
    resolve_subject,
)
from services.input_limits import public_limits
from services import input_limits as L
from services.input_validation import (
    check_hymn_layout_overrides,
    check_hymn_overrides,
    check_string_list,
)
from services.private_files import (
    media_file_url,
    media_type_for,
    preview_file_url,
    resolve_under_root,
    upload_file_url,
)
from services.storage_assets import (
    delete_user_asset,
    download_user_asset,
    list_user_assets,
    signed_asset_url,
    storage_ready,
    upload_user_asset,
)
from routes.admin import register_admin_routes
from routes.auth import register_auth_routes
from routes.parish import register_parish_routes

# Optional outputs produced alongside mass_poster.png (Phase 3)
_BUNDLE_OPTIONAL = (
    "gospel_moment.png",
    "mass_poster_16x9.png",
    "mass_poster_instagram_square.png",
    "mass_poster_instagram_story.png",
    "mass_poster_open_graph.png",
)

_PROJECT = Path(__file__).resolve().parent
_OUTPUT_DIR = _PROJECT / "outputs"
_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
_PREVIEW_DIR = _OUTPUT_DIR / "preview_slides"
_PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

_UPLOAD_DIR = uploads_dir()
_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
_MASS_ASSET_DIR = _UPLOAD_DIR / "mass_assets"
_MASS_ASSET_DIR.mkdir(parents=True, exist_ok=True)
_SAVED_POSTER_DIR = _UPLOAD_DIR / "saved_posters"
_SAVED_POSTER_DIR.mkdir(parents=True, exist_ok=True)
_SAVED_MEDIA_MUSIC_DIR = _UPLOAD_DIR / "saved_media" / "music"
_SAVED_MEDIA_VIDEO_DIR = _UPLOAD_DIR / "saved_media" / "video"
_SAVED_MEDIA_MUSIC_DIR.mkdir(parents=True, exist_ok=True)
_SAVED_MEDIA_VIDEO_DIR.mkdir(parents=True, exist_ok=True)
_SAVED_MEDIA_MANIFEST_PATH = _UPLOAD_DIR / "saved_media" / "manifest.json"

_MAX_MUSIC_BYTES = 10_000_000
_MAX_VIDEO_BYTES = 11_000_000
_ALLOWED_MUSIC_TYPES = frozenset(("audio/mpeg", "audio/mp3", "application/octet-stream"))
_ALLOWED_VIDEO_TYPES = frozenset(("video/mp4", "application/octet-stream"))

_ALLOWED_LOGO_TYPES = frozenset(
    ("image/png", "image/jpeg", "image/webp", "image/gif", "image/x-png", "image/jpg")
)
_MAX_LOGO_BYTES = 2_500_000
_MAX_ASSET_BYTES = 8_000_000
_ALLOWED_IMAGE_TYPES = _ALLOWED_LOGO_TYPES

_BUNDLE_NAME = "mass_bundle.zip"  # legacy fallback only


def _bundle_zip_name(export_stem: str) -> str:
    stem = (export_stem or "mass_bundle").strip() or "mass_bundle"
    return f"{stem}_bundle.zip"


def _collect_generation_owned_paths(result: GenerationResult) -> list[str]:
    owned: list[str] = []
    for attr in ("pptx_path", "poster_path", "poster_ppt_path"):
        p = getattr(result, attr, None)
        if p and Path(p).is_file():
            owned.append(Path(p).name)
    bundle = _OUTPUT_DIR / _bundle_zip_name(result.export_stem)
    if bundle.is_file():
        owned.append(bundle.name)
    post_dir = _OUTPUT_DIR / "posters"
    if post_dir.is_dir() and result.export_stem:
        for child in post_dir.glob("*.png"):
            if child.is_file():
                owned.append(f"posters/{child.name}")
    if result.export_stem:
        gospel = _OUTPUT_DIR / f"{result.export_stem}_gospel_moment.png"
        if gospel.is_file():
            owned.append(gospel.name)
    return owned


def _resolve_child_file(parent: Path, basename: str) -> Optional[Path]:
    base = Path(basename or "").name
    if not base or base != basename.strip():
        return None
    cand = (parent / base).resolve()
    try:
        cand.relative_to(parent.resolve())
    except ValueError:
        return None
    return cand if cand.is_file() else None


def _safe_storage_leaf(basename: str, *, folder: str) -> str:
    base = Path(basename or "").name
    if not base or base != basename.strip():
        raise HTTPException(status_code=400, detail="Invalid asset name.")
    return f"{folder.strip('/')}/{base}"


def _materialize_storage_asset(
    session: Optional[AuthSession], storage_path: str, *, suffix: str = ".png"
) -> Optional[Path]:
    if not session or not storage_ready(session.token):
        return None
    key = (storage_path or "").strip()
    if not key:
        return None
    try:
        raw = download_user_asset(access_token=session.token, path=key)
    except Exception:
        return None
    tmp = tempfile.NamedTemporaryFile(prefix="verbum_asset_", suffix=suffix, delete=False)
    tmp.write(raw)
    tmp.flush()
    tmp.close()
    return Path(tmp.name)


async def _save_uploaded_image(file: UploadFile, dest_dir: Path, *, prefix: str) -> str:
    ctype = (file.content_type or "").split(";")[0].strip().lower()
    if ctype not in _ALLOWED_IMAGE_TYPES:
        raise HTTPException(status_code=400, detail="Use a PNG, JPEG, WebP, or GIF image.")
    raw = await file.read()
    if len(raw) > _MAX_ASSET_BYTES:
        raise HTTPException(status_code=400, detail="Image too large (max about 8 MB).")
    dest_dir.mkdir(parents=True, exist_ok=True)
    orig = Path(file.filename or "upload").name
    ext = Path(orig).suffix.lower()
    if ext not in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        ext = ".png"
    name = f"{prefix}_{uuid.uuid4().hex}{ext}"
    out = dest_dir / name
    out.write_bytes(raw)
    return name


async def _save_uploaded_media(
    file: UploadFile,
    dest_dir: Path,
    *,
    prefix: str,
    allowed_ext: str,
    allowed_types: frozenset[str],
    max_bytes: int,
) -> str:
    orig = Path(file.filename or "upload").name
    ext = Path(orig).suffix.lower()
    if ext != allowed_ext:
        raise HTTPException(status_code=400, detail=f"Use a {allowed_ext.lstrip('.').upper()} file only.")
    ctype = (file.content_type or "").split(";")[0].strip().lower()
    if ctype and ctype not in allowed_types and ctype != "application/octet-stream":
        raise HTTPException(status_code=400, detail=f"Use a {allowed_ext.lstrip('.').upper()} file only.")
    raw = await file.read()
    if len(raw) > max_bytes:
        mb = max(1, max_bytes // (1024 * 1024))
        raise HTTPException(status_code=400, detail=f"File too large (max about {mb} MB).")
    dest_dir.mkdir(parents=True, exist_ok=True)
    name = f"{prefix}_{uuid.uuid4().hex}{allowed_ext}"
    out = dest_dir / name
    out.write_bytes(raw)
    return name


def _manifest_user_key(session: Optional[AuthSession]) -> str:
    if session and session.user and session.user.user_id:
        return session.user.user_id
    return "_local"


def _load_saved_media_manifest() -> dict[str, Any]:
    if not _SAVED_MEDIA_MANIFEST_PATH.is_file():
        return {"users": {}}
    try:
        data = json.loads(_SAVED_MEDIA_MANIFEST_PATH.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {"users": {}}
    if not isinstance(data, dict):
        return {"users": {}}
    data.setdefault("users", {})
    return data


def _save_saved_media_manifest(data: dict[str, Any]) -> None:
    _SAVED_MEDIA_MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    _SAVED_MEDIA_MANIFEST_PATH.write_text(
        json.dumps(data, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def _register_saved_media_meta(
    session: Optional[AuthSession],
    basename: str,
    display_name: str,
    kind: str,
) -> None:
    data = _load_saved_media_manifest()
    users = data.setdefault("users", {})
    user_key = _manifest_user_key(session)
    users.setdefault(user_key, {})[basename] = {
        "display_name": display_name,
        "kind": kind,
    }
    _save_saved_media_manifest(data)


def _unregister_saved_media_meta(session: Optional[AuthSession], basename: str) -> None:
    data = _load_saved_media_manifest()
    users = data.get("users", {})
    user_key = _manifest_user_key(session)
    bucket = users.get(user_key)
    if not bucket or basename not in bucket:
        return
    del bucket[basename]
    _save_saved_media_manifest(data)


def _saved_media_display_name(
    session: Optional[AuthSession],
    basename: str,
    *,
    kind: str,
) -> str:
    user_key = _manifest_user_key(session)
    entry = _load_saved_media_manifest().get("users", {}).get(user_key, {}).get(basename, {})
    name = str(entry.get("display_name") or "").strip()
    if name:
        return name
    stem = Path(basename).stem
    if stem.startswith(f"{kind}_"):
        return stem[len(kind) + 1 :] + Path(basename).suffix
    return basename


def _enrich_saved_media_rows(
    rows: list[dict[str, str]],
    *,
    session: Optional[AuthSession],
    kind: str,
) -> list[dict[str, str]]:
    out: list[dict[str, str]] = []
    for row in rows:
        item = dict(row)
        item["display_name"] = _saved_media_display_name(
            session,
            item.get("basename", ""),
            kind=kind,
        )
        item["kind"] = kind
        out.append(item)
    return out


def _list_saved_media_rows(
    local_dir: Path,
    *,
    folder: str,
    allowed_ext: str,
    session: Optional[AuthSession],
) -> list[dict[str, str]]:
    local_items: list[dict[str, str]] = []
    if local_dir.is_dir():
        for p in sorted(local_dir.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True):
            if p.is_file() and p.suffix.lower() == allowed_ext:
                local_items.append({"basename": p.name, "url": upload_file_url(f"{folder}/{p.name}")})

    if session and storage_ready(session.token):
        try:
            rows = list_user_assets(
                user_id=session.user.user_id,
                access_token=session.token,
                prefix=folder,
            )
            remote_items = [
                {"basename": row["name"], "url": row["url"] or ""}
                for row in rows
                if Path(str(row.get("name") or "")).suffix.lower() == allowed_ext
            ]
            merged: dict[str, dict[str, str]] = {item["basename"]: item for item in local_items}
            for item in remote_items:
                merged[item["basename"]] = item
            out = list(merged.values())
            out.sort(key=lambda r: r["basename"], reverse=True)
            return out
        except Exception:
            logger.warning("Could not list %s from storage; using local files.", folder, exc_info=True)
    return local_items


def _delete_saved_media_item(
    session: Optional[AuthSession],
    *,
    basename: str,
    local_dir: Path,
    storage_folder: str,
    allowed_ext: str,
    not_found_detail: str,
) -> dict[str, Any]:
    base = Path(basename or "").name
    if not base or base != basename.strip():
        raise HTTPException(status_code=400, detail="Invalid asset name.")
    removed = False
    if session and storage_ready(session.token):
        try:
            rel = _safe_storage_leaf(base, folder=storage_folder)
            delete_user_asset(
                user_id=session.user.user_id,
                access_token=session.token,
                relative_path=rel,
            )
            removed = True
        except HTTPException:
            raise
        except Exception:
            logger.warning("Storage delete failed for %s", base, exc_info=True)
    p = _resolve_child_file(local_dir, base)
    if p and p.suffix.lower() == allowed_ext:
        p.unlink(missing_ok=True)
        removed = True
    if not removed:
        raise HTTPException(status_code=404, detail=not_found_detail)
    _unregister_saved_media_meta(session, base)
    return {"ok": True}


def _upload_saved_media_url(
    *,
    session: Optional[AuthSession],
    local_path: Path,
    relative_path: str,
    content_type: str,
    local_url: str,
) -> str:
    if not session or not storage_ready(session.token):
        return local_url
    try:
        stored = upload_user_asset(
            user_id=session.user.user_id,
            access_token=session.token,
            relative_path=relative_path,
            raw=local_path.read_bytes(),
            content_type=content_type,
            upsert=True,
        )
        return stored.signed_url
    except Exception:
        logger.warning(
            "Storage upload failed for %s; serving local URL instead.",
            relative_path,
            exc_info=True,
        )
        return local_url


def _ai_poster_download_urls() -> dict[str, str]:
    """Private URLs for Hugging Face layout exports."""
    base = _OUTPUT_DIR / "posters"
    out: dict[str, str] = {}
    mapping = {
        "instagram": "mass_poster_instagram.png",
        "story": "mass_poster_story.png",
        "facebook": "mass_poster_facebook.png",
    }
    for key, fname in mapping.items():
        p = base / fname
        if p.is_file():
            out[key] = media_file_url(f"posters/{fname}")
    return out


def _latest_pptx_path() -> Optional[Path]:
    """Most recently modified deck in ``outputs/`` (supports stem-based filenames)."""
    cands = [p for p in _OUTPUT_DIR.glob("*.pptx") if p.is_file()]
    if not cands:
        return None
    return max(cands, key=lambda p: p.stat().st_mtime)


def _write_mass_bundle_zip(result: GenerationResult) -> Path:
    """Pack generated PPT, posters, stem-based social PNGs, gospel art, and optional extras.

    Uses ZIP_STORED (no recompress) — .pptx/.png are already compressed, and
    DEFLATED was slow enough to leave the Mass Builder UI stuck after the deck
    was already written.
    """
    out = _OUTPUT_DIR / _bundle_zip_name(result.export_stem)
    entries: list[tuple[Path, str]] = []
    for attr in ("pptx_path", "poster_path", "poster_ppt_path"):
        p = getattr(result, attr, None)
        if p and Path(p).is_file():
            pp = Path(p)
            entries.append((pp, pp.name))
    if result.include_social_exports and result.poster_path:
        parent = Path(result.poster_path).parent
        stem = Path(result.poster_path).stem
        for child in sorted(parent.glob(f"{stem}_*.png")):
            if child.is_file() and all(o[0] != child for o in entries):
                entries.append((child, child.name))
        post_dir = _OUTPUT_DIR / "posters"
        if post_dir.is_dir():
            for child in sorted(post_dir.glob("*.png")):
                if child.is_file() and all(o[0] != child for o in entries):
                    entries.append((child, f"posters/{child.name}"))
    if result.export_stem:
        g = _OUTPUT_DIR / f"{result.export_stem}_gospel_moment.png"
        if g.is_file() and all(o[0] != g for o in entries):
            entries.append((g, g.name))
    if result.include_social_exports:
        for name in _BUNDLE_OPTIONAL:
            p = _OUTPUT_DIR / name
            if p.is_file() and all(o[0] != p for o in entries):
                entries.append((p, name))
    with zipfile.ZipFile(out, "w", zipfile.ZIP_STORED) as zf:
        for path, arc in entries:
            zf.write(path, arcname=arc)
    return out


def _finish_generation_side_effects(
    *,
    result: GenerationResult,
    user_id: Optional[str],
    access_token: Optional[str],
    mass_date: str,
    celebrant: str,
    skip_ownership: bool = False,
) -> None:
    """Zip + storage + analytics — never on the HTTP request path."""
    try:
        _write_mass_bundle_zip(result)
        print(f"[generate] background zip ready stem={result.export_stem}", flush=True)
    except Exception:
        logger.warning("background bundle zip failed", exc_info=True)

    if user_id and not skip_ownership:
        try:
            register_owned_files(user_id, _collect_generation_owned_paths(result))
        except Exception:
            logger.warning("register_owned_files failed", exc_info=True)

    # Re-register zip path after it exists so the ZIP download link works.
    if user_id and result.export_stem:
        try:
            bundle_name = _bundle_zip_name(result.export_stem)
            if (_OUTPUT_DIR / bundle_name).is_file():
                register_owned_files(user_id, [bundle_name])
        except Exception:
            logger.warning("register zip ownership failed", exc_info=True)

    if user_id and access_token and storage_ready(access_token):
        stem = result.export_stem or "latest"
        upload_items: list[tuple[str, str, str]] = []
        if result.pptx_path and result.pptx_path.is_file():
            upload_items.append(
                (
                    f"generated/{stem}/{result.pptx_path.name}",
                    str(result.pptx_path),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            )
        if result.poster_path and result.poster_path.is_file():
            upload_items.append(
                (f"generated/{stem}/{result.poster_path.name}", str(result.poster_path), "image/png")
            )
        poster_ppt = result.poster_ppt_path
        if poster_ppt and Path(poster_ppt).is_file():
            upload_items.append(
                (
                    f"generated/{stem}/{Path(poster_ppt).name}",
                    str(Path(poster_ppt)),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            )
        bundle = _OUTPUT_DIR / _bundle_zip_name(result.export_stem)
        if bundle.is_file():
            upload_items.append((f"generated/{stem}/{bundle.name}", str(bundle), "application/zip"))
        if upload_items:
            _upload_generation_assets_best_effort(
                user_id=user_id,
                access_token=access_token,
                export_stem=stem,
                items=upload_items,
            )

    if user_id and access_token and supabase_enabled():
        try:
            from services.supabase_client import record_generation

            record_generation(
                user_id,
                mass_date=mass_date,
                celebrant=celebrant,
                output_summary={
                    "title": result.title,
                    "slide_count": result.slide_count,
                    "export_stem": result.export_stem,
                },
                access_token=access_token,
            )
        except Exception:
            pass


def _upload_generation_assets_best_effort(
    *,
    user_id: str,
    access_token: str,
    export_stem: str,
    items: list[tuple[str, str, str]],
) -> None:
    """Best-effort Supabase uploads that must never block the HTTP response.

    ``items`` is ``(relative_path, absolute_file_path, content_type)``.
    """
    for rel, file_path_s, ctype in items:
        path = Path(file_path_s)
        if not path.is_file():
            continue
        try:
            upload_user_asset(
                user_id=user_id,
                access_token=access_token,
                relative_path=rel,
                raw=path.read_bytes(),
                content_type=ctype,
                upsert=True,
            )
        except Exception:
            logger.warning(
                "Background storage upload failed for %s (stem=%s).",
                path.name,
                export_stem or "latest",
                exc_info=True,
            )


def _spawn_daemon(target, *, name: str, kwargs: dict[str, Any]) -> None:
    """Fire-and-forget worker.

    Do not use FastAPI BackgroundTasks here: with BaseHTTPMiddleware the client
    does not receive the response until background tasks finish, which is what
    left the Mass Builder UI stuck at ~62% after the .pptx was already written.
    """
    threading.Thread(target=target, kwargs=kwargs, name=name, daemon=True).start()


app = FastAPI(title="LiturgyFlow")
app.add_middleware(GZipMiddleware, minimum_size=1000)
templates = Jinja2Templates(directory=str(_PROJECT / "templates"))
_STATIC_DIR = _PROJECT / "static"
_STATIC_DIR.mkdir(parents=True, exist_ok=True)
app.mount("/static", StaticFiles(directory=str(_STATIC_DIR)), name="static")


@app.get("/favicon.ico", include_in_schema=False)
def favicon() -> FileResponse:
    return FileResponse(
        _STATIC_DIR / "brand" / "favicon.ico",
        media_type="image/vnd.microsoft.icon",
    )

register_auth_routes(app, templates)
register_admin_routes(app)
register_parish_routes(app)
register_security_middleware(app)


@app.get("/api/files/media/{file_path:path}")
def api_serve_media_file(
    file_path: str,
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> FileResponse:
    session_may_access_media(file_path, session)
    path = resolve_under_root(_OUTPUT_DIR, file_path)
    return FileResponse(path, media_type=media_type_for(path), filename=path.name)


@app.get("/api/files/uploads/{file_path:path}")
def api_serve_upload_file(
    file_path: str,
    _session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> FileResponse:
    path = resolve_under_root(_UPLOAD_DIR, file_path)
    return FileResponse(path, media_type=media_type_for(path), filename=path.name)


@app.get("/api/files/preview/{filename}")
def api_serve_preview_file(
    filename: str,
    _session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> FileResponse:
    path = resolve_under_root(_PREVIEW_DIR, filename)
    return FileResponse(path, media_type=media_type_for(path), filename=path.name)


@app.on_event("startup")
def _bootstrap_superadmin_roles() -> None:
    try:
        ensure_practice_secret_configured()
    except RuntimeError as exc:
        if is_production_runtime():
            raise
        print(f"[Verbum] Practice secret check skipped: {exc}")
    if supabase_enabled():
        try:
            from services.supabase_client import bootstrap_superadmin_roles_from_env

            count = bootstrap_superadmin_roles_from_env()
            if count:
                print(f"[Verbum] Promoted {count} profile(s) to superadmin from SUPERADMIN_EMAILS.")
        except Exception as exc:
            print(f"[Verbum] Superadmin bootstrap skipped: {exc}")

    import threading
    from datetime import date, timedelta

    def _upcoming_sunday_iso() -> str:
        today = date.today()
        days = (6 - today.weekday()) % 7
        return (today + timedelta(days=days)).isoformat()

    def _warm() -> None:
        try:
            warm_readings_for_date(_upcoming_sunday_iso())
        except Exception as exc:
            print(f"[Verbum] Readings warm-up skipped: {exc}")
        try:
            from services.song_catalog import catalog_lite_response

            catalog_lite_response()
        except Exception as exc:
            print(f"[Verbum] Catalog warm-up skipped: {exc}")

    threading.Thread(target=_warm, daemon=True).start()


def _sync_church_profile_to_supabase(
    session: Optional[AuthSession],
    *,
    community_name: Optional[str] = None,
    logo_path: Optional[str] = None,
    celebrant_names: Optional[list[str]] = None,
) -> dict[str, Any]:
    if not session or not supabase_enabled():
        raise HTTPException(status_code=503, detail="Sign in is required to save church profile.")
    try:
        from services.supabase_client import upsert_church_profile
        from services.user_church_context import set_church_profile

        saved = upsert_church_profile(
            session.user.user_id,
            community_name=community_name,
            logo_path=logo_path,
            celebrant_names=celebrant_names,
            access_token=session.token,
        )
        set_church_profile(saved)
        return saved
    except HTTPException:
        raise
    except Exception as exc:
        detail = str(exc).strip() or "Could not save church profile to Supabase."
        raise HTTPException(status_code=502, detail=detail) from exc


def _preview_to_json(p: PreviewPayload) -> dict[str, Any]:
    lc = p.liturgical_color
    liturgical: Optional[dict[str, Any]] = None
    if lc:
        liturgical = {
            "color_name": lc.get("color_name"),
            "hex": lc.get("hex"),
            "season": lc.get("season"),
            "rgb": list(lc.get("rgb", ())),
        }
    return {
        "ok": p.ok,
        "error": p.error,
        "title": p.title,
        "gospel_reference": p.gospel_reference,
        "season": p.season,
        "lectionary_cycle": p.lectionary_cycle,
        "liturgical_color": liturgical,
        "gospel_text_length": p.gospel_text_length,
        "sentences": p.sentences,
        "sentence_count": len(p.sentences),
        "quote_attribution": p.quote_attribution,
        "songs_by_section": p.songs_by_section,
        "gospel_quote": p.gospel_quote,
        "default_song_selections": p.default_song_selections,
        "estimated_slide_count": p.estimated_slide_count,
        "first_reading_reference": p.first_reading_reference,
        "first_reading_excerpt": p.first_reading_excerpt,
        "second_reading_reference": p.second_reading_reference,
        "second_reading_excerpt": p.second_reading_excerpt,
        "psalm_text": p.psalm_text,
        "psalm_verses": p.psalm_verses,
        "psalm_reference": p.psalm_reference,
        "psalm_refrains": p.psalm_refrains,
        "gospel_text": p.gospel_text,
        "readings_complete": p.readings_complete,
    }


class ExtraSongSection(BaseModel):
    label: str = Field(..., min_length=1, max_length=L.SECTION_LABEL)
    song_id: str = Field(..., min_length=1, max_length=L.SONG_ID)


class SongSelection(BaseModel):
    entrance: Optional[str] = None
    offertory: Optional[str] = None
    communion_1: Optional[str] = None
    communion_2: Optional[str] = None
    recessional: Optional[str] = None
    meditation: Optional[str] = None
    extra_sections: Optional[list[ExtraSongSection]] = Field(
        None,
        max_length=L.MAX_EXTRA_SONG_SECTIONS,
        description="User-added Mass song sections (custom labels) beyond the five defaults.",
    )


class CommunityNameBody(BaseModel):
    community_name: str = Field(..., min_length=1, max_length=L.CHURCH_NAME)


class CommunityProfileBody(BaseModel):
    community_name: Optional[str] = Field(None, min_length=1, max_length=L.CHURCH_NAME)
    celebrant_names: Optional[list[str]] = Field(
        None,
        max_length=L.MAX_CELEBRANTS,
        description="Saved Mass celebrant names for the in-app picker.",
    )

    @model_validator(mode="after")
    def _validate_celebrant_names(self) -> CommunityProfileBody:
        if self.celebrant_names is not None:
            self.celebrant_names = check_string_list(
                self.celebrant_names,
                field="celebrant_names",
                max_items=L.MAX_CELEBRANTS,
                item_max_len=L.CELEBRANT_NAME,
            )
        return self


class GeminiApiKeyBody(BaseModel):
    api_key: str = Field(..., min_length=8, max_length=L.API_KEY)


class PreviewBody(BaseModel):
    date: str = Field(..., min_length=8, description="YYYY-MM-DD")
    readings_only: bool = Field(
        False,
        description="Skip hymn recommendations and web hymn discovery for faster readings-only UI.",
    )
    refresh: bool = Field(
        False,
        description="Bypass in-memory cache and retry live lectionary fetch.",
    )


class GenerateBody(BaseModel):
    date: str = Field(..., min_length=8)
    celebrant: str = Field(..., min_length=1, max_length=L.CELEBRANT_NAME)
    co_celebrant: str = Field(
        "",
        max_length=L.CELEBRANT_NAME,
        description="Optional co-celebrant name shown on the Mass divider; omitted when blank.",
    )
    sentence_index: Optional[int] = Field(None, ge=0)
    poster_template: str = Field(
        "liturgical_color",
        description="liturgical_color | classic_white",
    )
    include_social_exports: bool = Field(
        False,
        description="When true, also export 1080×1350 feed PNG and Instagram/Story/OG variants.",
    )
    export_pdf: bool = Field(
        False,
        description="Deprecated — PDF export was removed from /api/generate. Ignored if set.",
    )
    include_gospel_art: bool = Field(True)
    include_ai_mass_poster: bool = Field(
        False,
        description="Use AI (OpenAI or Gemini) for primary parish posters.",
    )
    ai_poster_backend: str = Field(
        "openai",
        description="openai | gemini — which API generates the hero art when include_ai_mass_poster is true.",
    )
    ai_poster_style: str = Field(
        "cinematic",
        max_length=L.AI_POSTER_STYLE,
        description="OpenAI hero art style key from data/styles.json (5 presets).",
    )
    reuse_existing_poster: bool = Field(
        False,
        description="Reuse the cached hero art for this date+style instead of re-calling the AI image API.",
    )
    community_name: Optional[str] = Field(None, max_length=L.CHURCH_NAME)
    songs: Optional[SongSelection] = None
    custom_theme: Optional[dict[str, Any]] = None
    divider_poster_basename: Optional[str] = Field(
        None,
        max_length=L.FILE_BASENAME,
        description="Basename of a file previously uploaded to mass_assets/.",
    )
    lotw_poster: str = Field(
        "lotw1",
        max_length=16,
        description="Liturgy of the Word divider poster design: lotw1 | lotw2 | lotw3 | lotw4.",
    )
    lote_poster: str = Field(
        "lote1",
        max_length=16,
        description="Liturgy of the Eucharist divider poster design: lote1 | lote2 | lote3 | lote4.",
    )
    announcement_basenames: list[str] = Field(default_factory=list)
    mass_collection_amount: Optional[str] = Field(None, max_length=L.COLLECTION_AMOUNT)
    mass_collection_currency: Optional[str] = Field(
        "PHP",
        description="Mass collection currency: PHP | KRW | MYR",
        max_length=L.CURRENCY_CODE,
    )
    mass_collection_date_label: Optional[str] = Field(None, max_length=L.COLLECTION_DATE_LABEL)
    food_sponsors: list[str] = Field(default_factory=list)
    psalm_text_override: Optional[str] = Field(None, max_length=L.PSALM_FULL)
    psalm_refrain_index: Optional[int] = Field(None, ge=0)
    psalm_response_override: Optional[str] = Field(
        None,
        max_length=L.PSALM_REFRAIN,
        description="Manual responsorial-psalm refrain; used when psalm_text_override is empty.",
    )
    gospel_quote_override: Optional[str] = Field(
        None,
        max_length=L.GOSPEL_QUOTE,
        description="Exact Gospel line for slides; overrides sentence_index when non-empty.",
    )
    hymn_typography: Optional[dict[str, Any]] = Field(
        None,
        description="Per-section hymn slide typography: entrance, communion, default, etc.",
    )
    include_church_logo: bool = Field(
        False,
        description="When false, omit parish logo from PowerPoint slides.",
    )
    include_church_name: bool = Field(
        False,
        description="When false, omit parish / community name from PowerPoint slides.",
    )
    include_footer: bool = Field(
        False,
        description="Developer option: when true, show the bottom community/section footer tag on slides.",
    )
    hymn_lyric_overrides: Optional[dict[str, dict[str, str]]] = Field(
        None,
        description="Per-section hymn lyric text overrides: { entrance: { song_id: lyrics } }.",
    )
    hymn_layout_overrides: Optional[dict[str, dict[str, str]]] = Field(
        None,
        description="Per-song hymn slide layout overrides: { entrance: { song_id: 'single'|'dual' } }.",
    )
    creed_choice: str = Field(
        "nicene",
        description="Creed for the Mass deck: nicene | apostles (only one is included).",
    )
    our_father_choice: str = Field(
        "english",
        description="Our Father language: english | malay | tagalog | visaya | korean.",
    )
    hymn_lyrics_layout: str = Field(
        "dual",
        description="Hymn slide layout: single (1 block/slide) | dual (2 blocks/slide).",
    )

    @model_validator(mode="after")
    def _validate_generate_lists(self) -> GenerateBody:
        self.food_sponsors = check_string_list(
            self.food_sponsors,
            field="food_sponsors",
            max_items=L.MAX_FOOD_SPONSORS,
            item_max_len=L.FOOD_SPONSOR,
        ) or []
        self.announcement_basenames = check_string_list(
            self.announcement_basenames,
            field="announcement_basenames",
            max_items=L.MAX_ANNOUNCEMENT_IMAGES,
            item_max_len=L.FILE_BASENAME,
        ) or []
        self.hymn_lyric_overrides = check_hymn_overrides(self.hymn_lyric_overrides)
        self.hymn_layout_overrides = check_hymn_layout_overrides(self.hymn_layout_overrides)
        lotw = str(self.lotw_poster or "").strip().lower()
        self.lotw_poster = lotw if lotw in {"lotw1", "lotw2", "lotw3", "lotw4"} else "lotw1"
        lote = str(self.lote_poster or "").strip().lower()
        self.lote_poster = lote if lote in {"lote1", "lote2", "lote3", "lote4"} else "lote1"
        return self


class RefreshSongsBody(BaseModel):
    date: str = Field(..., min_length=8, description="YYYY-MM-DD")
    section: str = Field(..., min_length=3, max_length=L.SECTION_KEY)
    current_ids: list[str] = Field(default_factory=list)


class RefreshAllSongsBody(BaseModel):
    date: str = Field(..., min_length=8, description="YYYY-MM-DD")
    current_ids: dict[str, list[str]] = Field(default_factory=dict)


class ImportSongsBody(BaseModel):
    entrance: list[str] = Field(default_factory=list)
    offertory: list[str] = Field(default_factory=list)
    communion: list[str] = Field(default_factory=list)
    recessional: list[str] = Field(default_factory=list)


class SongRowBody(BaseModel):
    title: str = Field(..., min_length=1, max_length=L.SONG_TITLE)
    language: str = Field("English", max_length=L.LANGUAGE)
    mass_part: list[str] = Field(default_factory=list)


class ImportSongRowsBody(BaseModel):
    songs: list[SongRowBody] = Field(default_factory=list, max_length=L.MAX_IMPORT_SONG_ROWS)


class LyricsSelectionBody(BaseModel):
    section: str = Field(..., min_length=3, max_length=L.SECTION_KEY)
    id: str = Field(..., min_length=1, max_length=L.SONG_ID)


class FetchLyricsBody(BaseModel):
    selections: list[LyricsSelectionBody] = Field(default_factory=list)


class SaveLyricsBody(BaseModel):
    title: str = Field(..., min_length=1, max_length=L.SONG_TITLE)
    lyrics: str = Field(..., min_length=1, max_length=L.LYRICS_FULL)
    sections: list[str] = Field(default_factory=list)
    language: str = Field("English", max_length=L.LANGUAGE)
    author: str = Field("", max_length=L.SONG_AUTHOR)
    gospel_moods: Optional[list[str]] = Field(
        None,
        max_length=L.MAX_GOSPEL_MOODS,
        description="Gospel mood tags: triumphant, solemn, mercy, journey, reverent.",
    )


class PriestSubmissionBody(BaseModel):
    name: str = Field(..., min_length=1, max_length=L.CELEBRANT_NAME)


class CatalogSongPatchBody(BaseModel):
    title: Optional[str] = Field(None, max_length=L.SONG_TITLE)
    author: Optional[str] = Field(None, max_length=L.SONG_AUTHOR)
    lyrics: Optional[str] = Field(None, max_length=L.LYRICS_FULL)
    language: Optional[str] = Field(None, max_length=L.LANGUAGE)
    gospel_moods: Optional[list[str]] = Field(
        None,
        max_length=L.MAX_GOSPEL_MOODS,
        description="Gospel mood tags: triumphant, solemn, mercy, journey, reverent.",
    )


class GenerateImageBody(BaseModel):
    prompt: str = Field(..., min_length=1, max_length=L.AI_PROMPT)


class HymnSlidePreviewChunkBody(BaseModel):
    text: str = Field("", max_length=L.LYRICS_FULL)
    block_kind: str = Field("verse", max_length=32)


class HymnSlidePreviewPlanBody(BaseModel):
    order: list[int] = Field(default_factory=list, max_length=128)
    disabled: list[int] = Field(default_factory=list, max_length=128)


class HymnSlidePreviewBody(BaseModel):
    hymn_title: str = Field("Hymn", max_length=L.SONG_TITLE)
    section: str = Field("default", max_length=L.SECTION_KEY)
    layout: str = Field("dual", max_length=16)
    hymn_typography: Optional[dict[str, Any]] = None
    chunks: list[HymnSlidePreviewChunkBody] = Field(default_factory=list, max_length=128)
    plan: Optional[HymnSlidePreviewPlanBody] = None


class PracticeShareSongBody(BaseModel):
    slot_key: str = Field("", max_length=64)
    slot_label: str = Field("", max_length=120)
    section: str = Field("", max_length=32)
    hymn_id: str = Field(..., max_length=120)
    title: str = Field(..., max_length=L.SONG_TITLE)
    author: str = Field("", max_length=240)
    language: str = Field("", max_length=32)
    lyrics: str = Field("", max_length=L.LYRICS_FULL)


class PracticeShareBody(BaseModel):
    mass_date: str = Field(..., max_length=16)
    mass_title: str = Field("", max_length=240)
    parish_name: str = Field("", max_length=L.CHURCH_NAME)
    celebrant: str = Field("", max_length=L.CELEBRANT_NAME)
    songs: list[PracticeShareSongBody] = Field(default_factory=list, max_length=24)
    ttl_days: int = Field(0, ge=0, le=7)
    optional_pin: str = Field(..., min_length=6, max_length=6)


class PracticeUnlockBody(BaseModel):
    pin: str = Field("", min_length=6, max_length=6)


class PracticeLeadBody(BaseModel):
    lead_token: str = Field(..., min_length=16, max_length=512)


class PracticeLyricBlockBody(BaseModel):
    id: str = Field("", max_length=64)
    kind: str = Field("verse", max_length=32)
    label: str = Field("", max_length=80)
    body: str = Field("", max_length=L.LYRICS_FULL)
    enabled: bool = True


class PracticeLyricSongUpdateBody(BaseModel):
    hymn_id: str = Field(..., max_length=120)
    slot_key: str = Field("", max_length=64)
    blocks: list[PracticeLyricBlockBody] = Field(default_factory=list, max_length=48)


class PracticeLyricsUpdateBody(BaseModel):
    lead_token: str = Field(..., min_length=16, max_length=512)
    songs: list[PracticeLyricSongUpdateBody] = Field(default_factory=list, max_length=24)


class GenerateImageResponse(BaseModel):
    image: str
    path: str


def _resolve_soffice_bin() -> Optional[str]:
    custom = shutil.which("soffice")
    if custom:
        return custom
    mac_bin = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if Path(mac_bin).is_file():
        return mac_bin
    return None


def _extract_ppt_text_slides(ppt: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    prs = Presentation(str(ppt))
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                txt = (shape.text or "").strip()
                if txt:
                    lines.append(txt)
        slides.append({"index": idx + 1, "text": "\n\n".join(lines)[:2000]})
    return slides


@app.post("/api/hymn-slide-preview")
def api_hymn_slide_preview(
    body: HymnSlidePreviewBody,
    _session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    """Return hymn slide typography specs using the same fit logic as PPTX generation."""
    chunks = [{"text": c.text, "block_kind": c.block_kind} for c in body.chunks]
    plan = body.plan.model_dump() if body.plan else None
    return build_hymn_slide_preview(
        hymn_title=body.hymn_title,
        section=body.section.strip().lower(),
        layout=body.layout,
        hymn_typography=body.hymn_typography,
        chunks=chunks,
        plan=plan,
    )


@app.post("/api/ppt-preview/refresh")
def api_ppt_preview_refresh(
    _session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    """Render PPT slides to PNG images for in-app visual preview."""
    ppt = _latest_pptx_path()
    if not ppt or not ppt.is_file():
        return {"ok": True, "mode": "text", "slides": [], "message": "Generate deck first."}
    soffice = _resolve_soffice_bin()
    if not soffice:
        return {
            "ok": True,
            "mode": "text",
            "slides": _extract_ppt_text_slides(ppt),
            "message": "Install LibreOffice for exact image preview. Showing text fallback.",
        }

    for p in _PREVIEW_DIR.glob("*"):
        if p.is_file():
            p.unlink(missing_ok=True)

    png_paths, pdf_msg = render_ppt_preview_pngs(ppt, _PREVIEW_DIR, soffice_bin=soffice)
    if not png_paths:
        return {
            "ok": True,
            "mode": "text",
            "slides": _extract_ppt_text_slides(ppt),
            "message": (pdf_msg or "Could not render slide images.") + " Showing text fallback.",
        }

    ts = int(png_paths[0].stat().st_mtime) if png_paths else 0
    slides = [
        {
            "index": i + 1,
            "image_url": preview_file_url(f.name) + f"?t={ts}&v={i}",
        }
        for i, f in enumerate(png_paths)
    ]
    msg = pdf_msg or ""
    if not msg.strip():
        msg = "Full-deck preview (PDF rasterization)."
    return {"ok": True, "mode": "image", "slides": slides, "message": msg}


@app.get("/api/input-limits")
def api_input_limits() -> dict[str, int]:
    return public_limits()


@app.get("/api/platform/announcement")
def api_platform_announcement() -> dict[str, Any]:
    from services.platform_announcements import get_active_announcement

    return get_active_announcement()


@app.get("/api/feature-flags")
def api_feature_flags(
    session: Optional[AuthSession] = Depends(optional_session),
) -> dict[str, Any]:
    from services.feature_flags import flags_payload
    from services.parish_store import get_user_parish_context

    parish_id: Optional[str] = None
    if session:
        try:
            ctx = get_user_parish_context(session.user.user_id)
            pid = (ctx or {}).get("parish_id")
            parish_id = str(pid).strip() if pid else None
        except Exception:
            parish_id = None
    return flags_payload(parish_id=parish_id)


@app.get("/health")
async def health() -> dict[str, str]:
    # Async so a wedged sync thread pool (e.g. hung Supabase/httpx) cannot
    # make the health check itself hang.
    return {"status": "ok"}


@app.get("/go/google-search", response_class=HTMLResponse)
def go_google_search(q: str = "") -> HTMLResponse:
    """Open Google search via a same-origin page so mobile stays in the browser.

    Direct links to google.com often trigger the Google app (Universal Links /
    App Links). Loading this LiturgyFlow URL in a new tab, then navigating with
    location.replace, keeps the search in Safari/Chrome.
    """
    query = " ".join(str(q or "").split())[:500]
    if not query:
        raise HTTPException(status_code=400, detail="Missing search query.")
    google_url = "https://www.google.com/search?" + urlencode({"q": query, "hl": "en"})
    # Escape for HTML attribute + JS string
    href = google_url.replace("&", "&amp;").replace('"', "&quot;")
    js_url = json.dumps(google_url)
    html = (
        "<!DOCTYPE html><html lang=\"en\"><head>"
        "<meta charset=\"utf-8\"/>"
        "<meta name=\"viewport\" content=\"width=device-width, initial-scale=1\"/>"
        "<meta name=\"referrer\" content=\"no-referrer\"/>"
        f"<meta http-equiv=\"refresh\" content=\"0;url={href}\"/>"
        "<title>Opening Google search…</title>"
        f"<script>location.replace({js_url});</script>"
        "</head><body style=\"font:16px/1.4 system-ui,sans-serif;padding:24px;color:#1f2937\">"
        "<p>Opening Google search in this tab…</p>"
        f"<p><a href=\"{href}\" rel=\"noopener noreferrer\">Continue to Google</a></p>"
        "</body></html>"
    )
    return HTMLResponse(
        content=html,
        headers={
            "Cache-Control": "no-store",
            "Referrer-Policy": "no-referrer",
        },
    )


@app.get("/api/image-quota")
def api_image_quota(
    request: Request,
    session: Optional[AuthSession] = Depends(optional_session),
) -> dict[str, Any]:
    return quota_status_payload(session, request)


@app.get("/api/poster-exists")
def api_poster_exists(date: str, style: str = "cinematic") -> dict[str, bool]:
    """Report whether a cached AI hero already exists for this date + style.

    Used by the UI to offer reusing a previously generated Sunday poster
    instead of re-calling (and re-paying for) the AI image API.
    """
    from services.ai_styles import resolve_ai_image_style

    iso = (date or "").strip()
    if not re.fullmatch(r"\d{4}-\d{2}-\d{2}", iso):
        return {"exists": False}
    resolved_style = resolve_ai_image_style((style or "cinematic").strip())
    hero_path = _OUTPUT_DIR / "images" / f"{iso}_{resolved_style}_hero.png"
    return {"exists": hero_path.is_file()}


def _enforce_ai_image_quota(
    session: Optional[AuthSession],
    request: Request,
    *,
    source: str,
) -> dict[str, Any]:
    subject = resolve_subject(session, request)
    return reserve_daily_image_generation(subject, source=source)


@app.post("/generate-image", response_model=GenerateImageResponse)
def generate_image(
    body: GenerateImageBody,
    request: Request,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> GenerateImageResponse:
    import base64

    from generators.ai_image_generator import generate_openai_poster

    if not (os.getenv("OPENAI_API_KEY") or "").strip():
        raise HTTPException(status_code=503, detail="OPENAI_API_KEY is not configured.")

    prompt = body.prompt.strip()
    if not prompt:
        raise HTTPException(status_code=400, detail="prompt is required.")

    _enforce_ai_image_quota(session, request, source="generate-image")

    poster_path = _OUTPUT_DIR / "poster.png"
    try:
        saved = generate_openai_poster(prompt, output_path=poster_path)
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Image generation failed: {exc}") from exc

    raw = saved.read_bytes()
    return GenerateImageResponse(
        image=base64.b64encode(raw).decode("ascii"),
        path=str(saved),
    )


def _community_api_payload(session: Optional[AuthSession] = None) -> dict[str, Any]:
    profile = load_community()
    logo_exists = logo_file_absolute().is_file()
    celebrants = profile.get("celebrant_names")
    if not isinstance(celebrants, list):
        celebrants = get_celebrant_names()
    church_ctx = get_church_profile_context()
    user = session.user if session else None
    membership = membership_payload(church_ctx, user=user)
    logo_url: Optional[str] = None
    storage_logo = str((church_ctx or {}).get("logo_path") or "").strip()
    if session and storage_ready(session.token) and storage_logo:
        try:
            logo_url = signed_asset_url(access_token=session.token, path=storage_logo)
        except Exception:
            logo_url = None
    if not logo_url and logo_exists:
        logo_url = upload_file_url("community_logo.png")
    return {
        "ok": True,
        "community_name": get_community_name(),
        "celebrant_names": celebrants,
        "logo_url": logo_url,
        **membership,
    }


@app.get("/api/community")
def api_community(
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    return _community_api_payload(session)


@app.post("/api/community")
def api_set_community_name(
    body: CommunityNameBody,
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    return api_submit_parish_name(body, session)


@app.post("/api/community/submit-parish")
def api_submit_parish_name(
    body: CommunityNameBody,
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    if not session or not supabase_enabled():
        raise HTTPException(status_code=503, detail="Sign in is required to submit parish name.")
    from services.supabase_client import submit_parish_name

    saved = submit_parish_name(
        session.user.user_id,
        body.community_name.strip(),
        access_token=session.token,
    )
    from services.user_church_context import set_church_profile

    set_church_profile(saved)
    return _community_api_payload(session)


@app.post("/api/community/profile")
def api_set_community_profile(
    body: CommunityProfileBody,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    kwargs: dict[str, Any] = {}
    church_ctx = get_church_profile_context()
    if body.community_name is not None and not parish_name_is_locked(church_ctx):
        kwargs["community_name"] = body.community_name.strip()
    if body.celebrant_names is not None:
        kwargs["celebrant_names"] = body.celebrant_names
    if not kwargs:
        raise HTTPException(status_code=400, detail="No profile fields to update.")
    if session and supabase_enabled():
        _sync_church_profile_to_supabase(
            session,
            community_name=kwargs.get("community_name"),
            celebrant_names=kwargs.get("celebrant_names"),
        )
    else:
        update_community(**kwargs)
    return _community_api_payload(session)


@app.get("/api/settings/gemini-api-key")
def api_get_gemini_api_key_status(
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    from services.env_config import gemini_api_key_configured, gemini_api_key_hint

    configured = gemini_api_key_configured()
    hint = gemini_api_key_hint() if configured else None
    return {"configured": configured, "key_hint": hint}


@app.post("/api/settings/gemini-api-key")
def api_save_gemini_api_key(
    body: GeminiApiKeyBody,
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    from services.env_config import gemini_api_key_hint, save_gemini_api_key

    try:
        save_gemini_api_key(body.api_key.strip())
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    return {
        "ok": True,
        "configured": True,
        "key_hint": gemini_api_key_hint(),
        "key_format_warning": (
            None
            if body.api_key.strip().startswith("AIza")
            else (
                "This key does not look like a standard Google AI Studio key (AIza…). "
                "Image generation may fail with quota errors."
            )
        ),
    }


@app.post("/api/upload-logo")
async def api_upload_logo(
    file: UploadFile = File(...),
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    if session and supabase_enabled():
        church_ctx = get_church_profile_context()
        if not can_edit_logo(church_ctx):
            raise HTTPException(
                status_code=409,
                detail="Parish logo is locked and cannot be changed.",
            )
    ctype = (file.content_type or "").split(";")[0].strip().lower()
    if ctype not in _ALLOWED_LOGO_TYPES:
        raise HTTPException(
            status_code=400,
            detail="Use a PNG, JPEG, WebP, or GIF image.",
        )
    raw = await file.read()
    if len(raw) > _MAX_LOGO_BYTES:
        raise HTTPException(status_code=400, detail="Image must be at most about 2.5 MB.")

    try:
        from PIL import Image
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="Pillow is required for image upload.") from exc

    try:
        im = Image.open(io.BytesIO(raw))
        im.load()
    except OSError as exc:
        raise HTTPException(status_code=400, detail="Could not read image file.") from exc

    if im.mode in ("RGBA", "LA") or (im.mode == "P" and "transparency" in getattr(im, "info", {})):
        im = im.convert("RGBA")
    else:
        im = im.convert("RGB")

    _UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    out_path = logo_file_absolute()
    im.save(out_path, format="PNG", optimize=True)

    logo_url = upload_file_url("community_logo.png")
    if session and supabase_enabled():
        raw_png = out_path.read_bytes()
        stored = upload_user_asset(
            user_id=session.user.user_id,
            access_token=session.token,
            relative_path="logo/community_logo.png",
            raw=raw_png,
            content_type="image/png",
            upsert=True,
        )
        _sync_church_profile_to_supabase(session, logo_path=stored.path)
        logo_url = stored.signed_url
        church_ctx = get_church_profile_context()
        if parish_name_is_locked(church_ctx) and not (church_ctx or {}).get("logo_locked_at"):
            from services.supabase_client import lock_church_logo

            lock_church_logo(session.user.user_id, access_token=session.token)
    else:
        update_community(logo_path=LOGO_RELATIVE)
    out: dict[str, Any] = {
        "ok": True,
        "logo_url": logo_url,
        "message": "Logo saved. It will appear on the next generated poster.",
    }
    if session:
        out.update(membership_payload(get_church_profile_context(), user=session.user))
    return out


@app.post("/api/upload/mass-divider")
async def api_upload_mass_divider(
    file: UploadFile = File(...),
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    name = await _save_uploaded_image(file, _MASS_ASSET_DIR, prefix="divider")
    url = upload_file_url(f"mass_assets/{name}")
    if session and storage_ready(session.token):
        raw = (_MASS_ASSET_DIR / name).read_bytes()
        stored = upload_user_asset(
            user_id=session.user.user_id,
            access_token=session.token,
            relative_path=f"mass_assets/{name}",
            raw=raw,
            content_type=(file.content_type or "image/png"),
            upsert=True,
        )
        url = stored.signed_url
    return {"ok": True, "basename": name, "url": url}


@app.post("/api/upload/announcement-slide")
async def api_upload_announcement_slide(
    file: UploadFile = File(...),
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    name = await _save_uploaded_image(file, _MASS_ASSET_DIR, prefix="announce")
    url = upload_file_url(f"mass_assets/{name}")
    if session and storage_ready(session.token):
        raw = (_MASS_ASSET_DIR / name).read_bytes()
        stored = upload_user_asset(
            user_id=session.user.user_id,
            access_token=session.token,
            relative_path=f"mass_assets/{name}",
            raw=raw,
            content_type=(file.content_type or "image/png"),
            upsert=True,
        )
        url = stored.signed_url
    return {"ok": True, "basename": name, "url": url}


@app.post("/api/upload/saved-poster")
async def api_upload_saved_poster(
    file: UploadFile = File(...),
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    name = await _save_uploaded_image(file, _SAVED_POSTER_DIR, prefix="poster")
    url = upload_file_url(f"saved_posters/{name}")
    if session and storage_ready(session.token):
        raw = (_SAVED_POSTER_DIR / name).read_bytes()
        stored = upload_user_asset(
            user_id=session.user.user_id,
            access_token=session.token,
            relative_path=f"saved_posters/{name}",
            raw=raw,
            content_type=(file.content_type or "image/png"),
            upsert=True,
        )
        url = stored.signed_url
    return {"ok": True, "basename": name, "url": url}


@app.get("/api/saved-posters")
def api_list_saved_posters(
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    if session and storage_ready(session.token):
        rows = list_user_assets(
            user_id=session.user.user_id,
            access_token=session.token,
            prefix="saved_posters",
        )
        posters = [
            {"basename": row["name"], "url": row["url"] or ""}
            for row in rows
            if Path(str(row.get("name") or "")).suffix.lower() in {".png", ".jpg", ".jpeg", ".webp", ".gif"}
        ]
        posters.sort(key=lambda r: r["basename"], reverse=True)
        return {"ok": True, "posters": posters}
    if not _SAVED_POSTER_DIR.is_dir():
        return {"ok": True, "posters": []}
    rows = []
    for p in sorted(_SAVED_POSTER_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True):
        if p.is_file() and p.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
            rows.append({"basename": p.name, "url": upload_file_url(f"saved_posters/{p.name}")})
    return {"ok": True, "posters": rows}


@app.delete("/api/saved-posters/{basename}")
def api_delete_saved_poster(
    basename: str,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    if session and storage_ready(session.token):
        rel = _safe_storage_leaf(basename, folder="saved_posters")
        delete_user_asset(
            user_id=session.user.user_id,
            access_token=session.token,
            relative_path=rel,
        )
        return {"ok": True}
    p = _resolve_child_file(_SAVED_POSTER_DIR, basename)
    if not p:
        raise HTTPException(status_code=404, detail="Poster not found.")
    p.unlink(missing_ok=True)
    return {"ok": True}


@app.post("/api/upload/saved-media/music")
async def api_upload_saved_media_music(
    file: UploadFile = File(...),
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    name = await _save_uploaded_media(
        file,
        _SAVED_MEDIA_MUSIC_DIR,
        prefix="music",
        allowed_ext=".mp3",
        allowed_types=_ALLOWED_MUSIC_TYPES,
        max_bytes=_MAX_MUSIC_BYTES,
    )
    display_name = Path(file.filename or "upload").name
    _register_saved_media_meta(session, name, display_name, "music")
    url = upload_file_url(f"saved_media/music/{name}")
    url = _upload_saved_media_url(
        session=session,
        local_path=_SAVED_MEDIA_MUSIC_DIR / name,
        relative_path=f"saved_media/music/{name}",
        content_type="audio/mpeg",
        local_url=url,
    )
    return {"ok": True, "basename": name, "url": url, "display_name": display_name}


@app.post("/api/upload/saved-media/video")
async def api_upload_saved_media_video(
    file: UploadFile = File(...),
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    name = await _save_uploaded_media(
        file,
        _SAVED_MEDIA_VIDEO_DIR,
        prefix="video",
        allowed_ext=".mp4",
        allowed_types=_ALLOWED_VIDEO_TYPES,
        max_bytes=_MAX_VIDEO_BYTES,
    )
    display_name = Path(file.filename or "upload").name
    _register_saved_media_meta(session, name, display_name, "video")
    url = upload_file_url(f"saved_media/video/{name}")
    url = _upload_saved_media_url(
        session=session,
        local_path=_SAVED_MEDIA_VIDEO_DIR / name,
        relative_path=f"saved_media/video/{name}",
        content_type="video/mp4",
        local_url=url,
    )
    return {"ok": True, "basename": name, "url": url, "display_name": display_name}


@app.get("/api/saved-media")
def api_list_saved_media(
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    music = _list_saved_media_rows(
        _SAVED_MEDIA_MUSIC_DIR,
        folder="saved_media/music",
        allowed_ext=".mp3",
        session=session,
    )
    video = _list_saved_media_rows(
        _SAVED_MEDIA_VIDEO_DIR,
        folder="saved_media/video",
        allowed_ext=".mp4",
        session=session,
    )
    return {
        "ok": True,
        "music": _enrich_saved_media_rows(music, session=session, kind="music"),
        "video": _enrich_saved_media_rows(video, session=session, kind="video"),
    }


@app.delete("/api/saved-media/music/{basename}")
def api_delete_saved_media_music(
    basename: str,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    return _delete_saved_media_item(
        session,
        basename=basename,
        local_dir=_SAVED_MEDIA_MUSIC_DIR,
        storage_folder="saved_media/music",
        allowed_ext=".mp3",
        not_found_detail="Audio file not found.",
    )


@app.delete("/api/saved-media/video/{basename}")
def api_delete_saved_media_video(
    basename: str,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    return _delete_saved_media_item(
        session,
        basename=basename,
        local_dir=_SAVED_MEDIA_VIDEO_DIR,
        storage_folder="saved_media/video",
        allowed_ext=".mp4",
        not_found_detail="Video file not found.",
    )


@app.get("/api/catalog/songs")
def api_catalog_songs(
    lite: bool = True,
    if_none_match: Optional[str] = Header(None, alias="If-None-Match"),
    _session: Optional[AuthSession] = Depends(require_approved_membership),
) -> Response:
    if lite:
        body, etag = catalog_lite_response()
        headers = {
            "Cache-Control": "private, max-age=300",
            "ETag": etag,
            "Vary": "Authorization",
        }
        if if_none_match and if_none_match.strip() == etag:
            return Response(status_code=304, headers=headers)
        return Response(content=body, media_type="application/json", headers=headers)
    payload = {
        "ok": True,
        "catalog": catalog_for_api(include_inferred_moods=True),
    }
    return JSONResponse(payload, headers={"Cache-Control": "private, max-age=120"})


@app.get("/api/catalog/songs/{section}/{hymn_id:path}")
def api_get_catalog_song(
    section: str,
    hymn_id: str,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> JSONResponse:
    uid = session.user.user_id if session else None
    hid = (hymn_id or "").strip()
    requested_sec = (section or "").strip().lower()
    allowed, retry_after = check_catalog_lyric_fetch_allowed(uid, hid)
    if not allowed:
        raise HTTPException(
            status_code=429,
            detail="Too many hymn lyric requests. Please slow down.",
            headers={"Retry-After": str(max(1, retry_after))},
        )
    row = get_hymn(requested_sec, hid)
    resolved_section = requested_sec
    by_id_sec, by_id_row = find_catalog_row_by_id(hid)
    if by_id_row is not None and by_id_sec:
        row = by_id_row
        resolved_section = by_id_sec
    if not row:
        normalized = fetch_song_from_normalized(hid)
        if normalized:
            row = normalized
            resolved_section = str(normalized.get("section") or requested_sec).strip().lower() or requested_sec
    if not row:
        raise HTTPException(status_code=404, detail="Song not found.")
    lyrics = fetch_lyrics_from_normalized(hid) or str(row.get("lyrics") or "")
    catalog_lyrics = lyrics
    parish_version = False
    if session and session.user and session.user.user_id:
        try:
            from services.parish_hymn_overrides import get_override
            from services.parish_store import get_user_parish_context

            parish_ctx = get_user_parish_context(session.user.user_id) or {}
            parish_id = str(parish_ctx.get("parish_id") or "").strip()
            if parish_id:
                ov = get_override(parish_id, hymn_id=hid, section=resolved_section)
                if not ov:
                    ov = get_override(parish_id, hymn_id=hid)
                if ov and str(ov.get("lyrics") or "").strip():
                    lyrics = str(ov.get("lyrics") or "")
                    parish_version = True
                    if ov.get("section"):
                        resolved_section = str(ov.get("section") or resolved_section)
        except Exception:
            pass
    log_lyric_read(
        user_id=uid,
        hymn_id=hid,
        section=resolved_section,
        source="catalog_api",
    )
    payload = {
        "ok": True,
        "section": resolved_section,
        "song": {
            "id": str(row.get("id") or hid),
            "title": str(row.get("title") or ""),
            "author": str(row.get("author") or ""),
            "language": str(row.get("language") or "").strip(),
            "lyrics": lyrics,
            "catalog_lyrics": catalog_lyrics,
            "parish_version": parish_version,
            "gospel_moods": gospel_moods_for_song(row),
        },
    }
    return JSONResponse(payload, headers={"Cache-Control": "private, no-store"})


@app.patch("/api/catalog/songs/{section}/{hymn_id:path}")
def api_patch_catalog_song(
    section: str,
    hymn_id: str,
    body: CatalogSongPatchBody,
    session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    res = update_catalog_song(
        section=section,
        hymn_id=hymn_id,
        title=body.title,
        author=body.author,
        lyrics=body.lyrics,
        language=body.language,
        gospel_moods=body.gospel_moods,
        updated_by=session.user.user_id if session else None,
    )
    if not res.get("ok"):
        raise HTTPException(status_code=400, detail=res.get("error") or "Update failed.")
    return res


@app.delete("/api/catalog/songs/{section}/{hymn_id:path}")
def api_delete_catalog_song(
    section: str,
    hymn_id: str,
    session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    res = delete_catalog_song(
        section=section,
        hymn_id=hymn_id,
        updated_by=session.user.user_id if session else None,
    )
    if not res.get("ok"):
        raise HTTPException(status_code=400, detail=res.get("error") or "Delete failed.")
    return res


@app.post("/api/design/analyze-template")
async def api_design_analyze_template(
    file: UploadFile = File(...),
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    ctype = (file.content_type or "").split(";")[0].strip().lower()
    if ctype not in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/octet-stream",
    } and not (file.filename or "").lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Upload a .pptx file.")
    raw = await file.read()
    if len(raw) > 25_000_000:
        raise HTTPException(status_code=400, detail="Presentation too large.")
    tmp = _UPLOAD_DIR / f"_analyze_{uuid.uuid4().hex}.pptx"
    try:
        tmp.write_bytes(raw)
        result = analyze_pptx_theme(tmp)
    finally:
        tmp.unlink(missing_ok=True)
    if not result.get("ok"):
        raise HTTPException(status_code=400, detail=result.get("error") or "Analysis failed.")
    return result


@app.get("/", response_class=HTMLResponse)
def index(request: Request) -> Any:
    # Public marketing landing page. Signed-in visitors are bounced to /home
    # client-side (see landing.html); the app itself lives at /home and friends.
    return templates.TemplateResponse(
        request,
        "landing.html",
        {
            "title": "LiturgyFlow",
            "auth_enabled": auth_enabled(),
            "invite_contact_email": invite_contact_email(),
        },
    )


@app.get("/home", response_class=HTMLResponse)
@app.get("/notifications", response_class=HTMLResponse)
@app.get("/today", response_class=HTMLResponse)
@app.get("/radio", response_class=HTMLResponse)
@app.get("/mass/builder", response_class=HTMLResponse)
@app.get("/mass/calendar", response_class=HTMLResponse)
@app.get("/media/posters", response_class=HTMLResponse)
@app.get("/media/presentation", response_class=HTMLResponse)
@app.get("/media/history", response_class=HTMLResponse)
@app.get("/library/songs", response_class=HTMLResponse)
@app.get("/library/collections", response_class=HTMLResponse)
@app.get("/design/theme-lab", response_class=HTMLResponse)
@app.get("/design/templates", response_class=HTMLResponse)
@app.get("/settings/church", response_class=HTMLResponse)
@app.get("/settings/app", response_class=HTMLResponse)
@app.get("/settings/team", response_class=HTMLResponse)
@app.get("/superadmin", response_class=HTMLResponse)
@app.get("/lyrics-dashboard", response_class=HTMLResponse)
@app.get("/theme-dashboard", response_class=HTMLResponse)
@app.get("/mass-flow-dashboard", response_class=HTMLResponse)
def dashboard(request: Request) -> Any:
    return templates.TemplateResponse(
        request,
        "index.html",
        {"title": "LiturgyFlow", "auth_enabled": auth_enabled()},
    )


def _practice_share_url(request: Request, token: str) -> str:
    base = str(request.base_url).rstrip("/")
    return f"{base}/practice/{token}"


@app.get("/practice/{token}", response_class=HTMLResponse)
def practice_page(request: Request, token: str) -> Any:
    lead = (request.query_params.get("lead") or "").strip()
    return templates.TemplateResponse(
        request,
        "practice.html",
        {
            "title": "Choir practice",
            "token": (token or "").strip(),
            "lead_token": lead,
        },
    )


@app.get("/api/practice/qr/{token}")
def api_practice_qr(token: str, request: Request) -> Response:
    """QR code image for a valid practice share link."""
    row = get_practice_share_by_token(token)
    if not row:
        raise HTTPException(status_code=404, detail="Practice link not found.")
    url = _practice_share_url(request, token.strip())
    data_url = practice_qr_data_url(url)
    if not data_url:
        raise HTTPException(status_code=503, detail="QR generation unavailable.")
    header, encoded = data_url.split(",", 1)
    raw = base64.b64decode(encoded)
    media = "image/svg+xml" if "svg" in header else "image/png"
    return Response(content=raw, media_type=media, headers={"Cache-Control": "private, max-age=3600"})


def _practice_rate_limit_response(retry_after: int) -> JSONResponse:
    return JSONResponse(
        status_code=429,
        content={"detail": "Too many requests. Please slow down."},
        headers={"Retry-After": str(max(1, retry_after)), **practice_no_store_headers()},
    )


@app.get("/api/practice/{token}")
def api_practice_share(
    token: str,
    request: Request,
) -> Any:
    """Public read-only lyrics snapshot for choir rehearsal (token is the secret)."""
    if token.strip().lower() == "qr":
        raise HTTPException(status_code=404, detail="Not found.")
    tok = token.strip()
    allowed, retry_after = check_practice_fetch_allowed(request, tok)
    if not allowed:
        return _practice_rate_limit_response(retry_after)
    allowed_token, retry_token = check_practice_token_allowed(tok)
    if not allowed_token:
        return _practice_rate_limit_response(retry_token)
    row = get_practice_share_by_token(tok)
    if not row:
        raise HTTPException(status_code=404, detail="This practice link is invalid or has expired.")
    lead = (request.query_params.get("lead") or "").strip()
    can_edit = bool(lead) and verify_lead_token(tok, lead, row.get("expires_at"))
    unlocked = is_unlocked(request, tok, row.get("expires_at")) or can_edit
    payload = fetch_practice_share(tok, unlocked=unlocked, can_edit=can_edit)
    if not payload.get("ok"):
        raise HTTPException(status_code=404, detail=payload.get("error") or "Not found.")
    return JSONResponse(payload, headers=practice_no_store_headers())


@app.post("/api/practice/{token}/unlock")
def api_practice_unlock(
    token: str,
    body: PracticeUnlockBody,
    request: Request,
) -> Any:
    """Verify PIN and issue a device-bound unlock cookie (PIN never goes in the URL)."""
    if token.strip().lower() == "qr":
        raise HTTPException(status_code=404, detail="Not found.")
    tok = token.strip()
    allowed, retry_after = check_pin_unlock_allowed(request, tok)
    if not allowed:
        return _practice_rate_limit_response(retry_after)
    allowed_token, retry_token = check_practice_token_allowed(tok)
    if not allowed_token:
        return _practice_rate_limit_response(retry_token)
    result = verify_practice_share_pin(tok, body.pin or "")
    if not result.get("ok"):
        status = 401 if result.get("requires_pin") else 404
        return JSONResponse(
            result,
            status_code=status,
            headers=practice_no_store_headers(),
        )
    row = get_practice_share_by_token(tok)
    if not row:
        raise HTTPException(status_code=404, detail="This practice link is invalid or has expired.")
    response = JSONResponse(result, headers=practice_no_store_headers())
    issue_unlock_cookie(request, response, tok, row.get("expires_at"))
    return response


@app.post("/api/practice/{token}/lead")
def api_practice_lead(
    token: str,
    body: PracticeLeadBody,
    request: Request,
) -> Any:
    """Leader entry — unlock + edit mode via signed lead token (no PIN)."""
    if token.strip().lower() == "qr":
        raise HTTPException(status_code=404, detail="Not found.")
    tok = token.strip()
    allowed, retry_after = check_practice_fetch_allowed(request, tok)
    if not allowed:
        return _practice_rate_limit_response(retry_after)
    row = get_practice_share_by_token(tok)
    if not row:
        raise HTTPException(status_code=404, detail="This practice link is invalid or has expired.")
    if not verify_lead_token(tok, body.lead_token, row.get("expires_at")):
        raise HTTPException(status_code=401, detail="Leader link is invalid or expired.")
    payload = fetch_practice_share(tok, unlocked=True, can_edit=True)
    response = JSONResponse(payload, headers=practice_no_store_headers())
    issue_unlock_cookie(request, response, tok, row.get("expires_at"))
    return response


@app.post("/api/practice/{token}/lyrics")
def api_practice_lyrics_update(
    token: str,
    body: PracticeLyricsUpdateBody,
    request: Request,
) -> Any:
    """Practice-only lyric edits for the share leader. Never writes to the hymn catalog."""
    if token.strip().lower() == "qr":
        raise HTTPException(status_code=404, detail="Not found.")
    tok = token.strip()
    allowed, retry_after = check_practice_fetch_allowed(request, tok)
    if not allowed:
        return _practice_rate_limit_response(retry_after)
    row = get_practice_share_by_token(tok)
    if not row:
        raise HTTPException(status_code=404, detail="This practice link is invalid or has expired.")
    if not verify_lead_token(tok, body.lead_token, row.get("expires_at")):
        raise HTTPException(status_code=401, detail="Leader link is invalid or expired.")
    songs = [s.model_dump() for s in body.songs]
    result = update_practice_share_lyrics(tok, songs_update=songs)
    if not result.get("ok"):
        raise HTTPException(status_code=400, detail=result.get("error") or "Could not save.")
    return JSONResponse(result, headers=practice_no_store_headers())


@app.post("/api/practice/share")
def api_create_practice_share(
    body: PracticeShareBody,
    request: Request,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> dict[str, Any]:
    parish_id: Optional[str] = None
    created_by: Optional[str] = None
    if session:
        created_by = session.user.user_id
        try:
            from services.parish_store import get_user_parish_context

            ctx = get_user_parish_context(session.user.user_id)
            pid = (ctx or {}).get("parish_id")
            parish_id = str(pid).strip() if pid else None
        except Exception:
            parish_id = None

    parish_name = (body.parish_name or "").strip()
    if not parish_name:
        parish_name = get_community_name()

    actor_key = session.user.user_id if session else "anon"
    allowed, retry_after = check_practice_share_create_allowed(request, actor_key)
    if not allowed:
        return _practice_rate_limit_response(retry_after)

    from services.feature_flags import resolve_flags

    if not resolve_flags(parish_id=parish_id).get("choir_practice_shares", True):
        raise HTTPException(
            status_code=403,
            detail="Choir practice shares are temporarily disabled.",
        )

    songs = [s.model_dump() for s in body.songs]
    try:
        result = create_practice_share(
            created_by_user_id=created_by,
            parish_id=parish_id,
            mass_date=body.mass_date.strip(),
            mass_title=body.mass_title,
            parish_name=parish_name,
            celebrant=body.celebrant,
            songs=songs,
            ttl_days=body.ttl_days,
            optional_pin=body.optional_pin or None,
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except RuntimeError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    token = str(result.get("token") or "")
    share_url = _practice_share_url(request, token)
    lead = issue_lead_token(token, result.get("expires_at"))
    # Return qr_url for async image load; skip embedding a data URL so create stays fast.
    return {
        **result,
        "url": share_url,
        "lead_token": lead,
        "lead_url": f"{share_url}?lead={lead}",
        "qr_url": f"/api/practice/qr/{token}",
        "qr_data_url": None,
    }


@app.get("/api/catholic-news")
def api_catholic_news(
    vatican: bool = True,
    cna: bool = True,
    limit: int = 6,
    offset: int = 0,
    max_age_days: int = 3,
) -> Any:
    """Headlines from Vatican News and Catholic News Agency RSS (fresh each request)."""
    cap = max(1, min(int(limit), 15))
    off = max(0, int(offset))
    age = max(0, min(int(max_age_days), 14))
    return fetch_catholic_headlines(
        include_vatican=vatican,
        include_cna=cna,
        max_items=cap,
        offset=off,
        max_age_days=age,
    )


@app.get("/api/wyd-news")
def api_wyd_news(limit: int = 6) -> Any:
    """World Youth Day Seoul 2027 announcements (official site) plus countdown metadata."""
    from datetime import date

    from services.wyd_news import fetch_wyd_announcements

    cap = max(1, min(int(limit), 15))
    errors: list[str] = []

    # Primary source: the official WYD announcement board (server-rendered).
    items: list[Any] = []
    try:
        items = fetch_wyd_announcements(limit=cap)
    except Exception as exc:  # pragma: no cover - defensive
        errors.append(f"WYD site: {exc}")

    # Fallback: keyword-filtered Catholic news if the official board is empty.
    if not items:
        feed = fetch_catholic_headlines(
            include_vatican=True,
            include_cna=True,
            max_items=cap,
            offset=0,
            max_age_days=0,
            keywords=["world youth day", "wyd", "seoul 2027", "jornada mundial"],
        )
        items = feed.get("items", [])
        errors.extend(feed.get("errors", []))

    event_start = date(2027, 8, 3)
    event_end = date(2027, 8, 8)
    days_until = (event_start - date.today()).days
    return {
        "ok": bool(items) or not errors,
        "items": items,
        "errors": errors,
        "event_start": event_start.isoformat(),
        "event_end": event_end.isoformat(),
        "location": "Seoul, South Korea",
        "days_until": days_until,
        "official_url": "https://wydseoul.org/en",
    }


@app.get("/api/ewtn/radio")
def api_ewtn_radio() -> Any:
    """Catholic live radio stations (EWTN feeds + curated networks)."""
    return radio_catalog()


@app.get("/api/calendar/month")
def api_calendar_month(year: int, month: int) -> Any:
    """Lightweight per-day summaries for the liturgical calendar grid."""
    try:
        return fetch_calendar_month(year, month)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.get("/api/readings/{date}")
def api_readings(
    date: str,
    refresh: bool = False,
    _session: Optional[AuthSession] = Depends(optional_session),
) -> JSONResponse:
    """Fast lectionary readings for dashboard cards (no song discovery).

    Public liturgical data — anonymous-friendly, mirroring ``POST /api/preview``
    so the home dashboard loads readings before/without sign-in.
    """
    payload, from_cache = readings_snapshot(date.strip(), force_refresh=refresh)
    if not payload.get("ok"):
        raise HTTPException(
            status_code=400,
            detail=payload.get("error") or "Unable to load readings.",
        )
    complete = bool(payload.get("readings_complete"))
    max_age = 3600 if from_cache and complete else 15
    return JSONResponse(
        payload,
        headers={"Cache-Control": f"private, max-age={max_age}"},
    )


@app.get("/api/gospel-image/{date}")
def api_gospel_image(date: str) -> JSONResponse:
    """Gospel-matched background image (Openverse / Creative Commons) for the home card."""
    from services.gospel_image_search import fetch_gospel_background

    payload = fetch_gospel_background(date.strip())
    max_age = 21600 if payload.get("ok") else 300
    return JSONResponse(payload, headers={"Cache-Control": f"public, max-age={max_age}"})


@app.post("/api/preview")
async def api_preview(
    body: PreviewBody,
    _session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> Any:
    p = await run_in_threadpool(
        fetch_preview,
        body.date.strip(),
        readings_only=body.readings_only,
        force_refresh=body.refresh,
    )
    return _preview_to_json(p)


@app.post("/api/generate")
def api_generate(
    body: GenerateBody,
    request: Request,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> Any:
    """Build Mass media and return local download URLs.

    Sync ``def`` (not async): FastAPI runs this in a worker thread so
    BaseHTTPMiddleware cannot deadlock against ``run_in_threadpool`` while the
    UI sits at ~62%. PDF export was removed from this path entirely.
    """
    print(f"[generate] start date={body.date!r}", flush=True)
    song_map = body.songs.model_dump(exclude_none=True) if body.songs else None
    temp_assets: list[Path] = []
    divider_path = None
    if body.divider_poster_basename and str(body.divider_poster_basename).strip():
        raw_divider = str(body.divider_poster_basename).strip()
        divider_path = _resolve_child_file(_MASS_ASSET_DIR, raw_divider)
        if not divider_path and session and storage_ready(session.token):
            key = _safe_storage_leaf(raw_divider, folder="mass_assets")
            divider_path = _materialize_storage_asset(session, f"{session.user.user_id}/{key}")
            if divider_path:
                temp_assets.append(divider_path)
    ann_paths: list[Path] = []
    for raw in body.announcement_basenames or []:
        bn = str(raw).strip()
        if not bn:
            continue
        p = _resolve_child_file(_MASS_ASSET_DIR, bn)
        if not p and session and storage_ready(session.token):
            key = _safe_storage_leaf(bn, folder="mass_assets")
            p = _materialize_storage_asset(session, f"{session.user.user_id}/{key}")
            if p:
                temp_assets.append(p)
        if p:
            ann_paths.append(p)
        if len(ann_paths) >= 24:
            break
    sponsors = [str(s).strip() for s in (body.food_sponsors or []) if str(s).strip()][:24]
    psalm_override = (body.psalm_text_override or "").strip() or None
    gospel_override = (body.gospel_quote_override or "").strip() or None

    if body.include_ai_mass_poster:
        backend = (body.ai_poster_backend or "openai").strip().lower()
        _enforce_ai_image_quota(
            session,
            request,
            source=f"mass-poster:{backend}",
        )

    try:
        print("[generate] building media…", flush=True)
        result = generate_mass_media(
            body.date.strip(),
            body.celebrant.strip(),
            co_celebrant=(body.co_celebrant or "").strip(),
            sentence_index=body.sentence_index,
            poster_template=body.poster_template,
            include_social_exports=body.include_social_exports,
            include_gospel_art=body.include_gospel_art,
            include_ai_mass_poster=body.include_ai_mass_poster,
            ai_poster_backend=(body.ai_poster_backend or "openai").strip().lower(),
            ai_poster_style=body.ai_poster_style.strip() or "cinematic",
            reuse_existing_poster=body.reuse_existing_poster,
            community_name=body.community_name.strip() if body.community_name else None,
            song_selections=song_map,
            custom_theme=body.custom_theme,
            divider_poster_path=divider_path,
            lotw_poster=body.lotw_poster,
            lote_poster=body.lote_poster,
            announcement_image_paths=ann_paths or None,
            mass_collection_amount=body.mass_collection_amount.strip() if body.mass_collection_amount else None,
            mass_collection_date_label=body.mass_collection_date_label.strip()
            if body.mass_collection_date_label
            else None,
            mass_collection_currency=body.mass_collection_currency.strip().upper()
            if body.mass_collection_currency
            else "PHP",
            food_sponsors=sponsors or None,
            psalm_text_override=psalm_override,
            psalm_refrain_index=body.psalm_refrain_index,
            psalm_response_override=(body.psalm_response_override or "").strip() or None,
            gospel_quote_override=gospel_override,
            hymn_typography=body.hymn_typography,
            include_church_logo=body.include_church_logo,
            include_church_name=body.include_church_name,
            include_footer=body.include_footer,
            hymn_lyric_overrides=body.hymn_lyric_overrides,
            creed_choice=body.creed_choice,
            our_father_choice=body.our_father_choice,
            hymn_lyrics_layout=body.hymn_lyrics_layout,
            hymn_layout_overrides=body.hymn_layout_overrides,
        )
    finally:
        for p in temp_assets:
            p.unlink(missing_ok=True)
    if not result.ok:
        raise HTTPException(status_code=400, detail=result.error or "Generation failed.")

    print("[generate] ppt ready — returning URLs (zip/upload in background)", flush=True)
    bundle_rel = _bundle_zip_name(result.export_stem)
    # Prefer an existing zip from a prior run; never block the response on zipping.
    zip_ready = (_OUTPUT_DIR / bundle_rel).is_file()

    lc = result.liturgical_color
    liturgical_payload: Optional[dict[str, Any]] = None
    if lc:
        liturgical_payload = {
            "color_name": lc.get("color_name"),
            "hex": lc.get("hex"),
            "season": lc.get("season"),
            "rgb": list(lc.get("rgb", ())),
        }
    poster_ppt = result.poster_ppt_path
    ppt_url = media_file_url(result.pptx_path.name) if result.pptx_path and result.pptx_path.is_file() else None
    poster_url = media_file_url(result.poster_path.name) if result.poster_path and result.poster_path.is_file() else None
    poster_ppt_url = (
        media_file_url(poster_ppt.name) if poster_ppt and Path(poster_ppt).is_file() else None
    )
    zip_url = media_file_url(bundle_rel) if zip_ready else None

    # Ownership must be registered before the response — the browser downloads
    # immediately and /api/files/media/* 404s without it (UI stuck at 100%).
    if session:
        try:
            print("[generate] registering ownership…", flush=True)
            register_owned_files(session.user.user_id, _collect_generation_owned_paths(result))
            print("[generate] ownership ok", flush=True)
        except Exception:
            logger.warning("register_owned_files failed", exc_info=True)

    uid = session.user.user_id if session else None
    token = session.token if session else None
    _spawn_daemon(
        _finish_generation_side_effects,
        name="mass-gen-finish",
        kwargs={
            "result": result,
            "user_id": uid,
            "access_token": token,
            "mass_date": body.date.strip(),
            "celebrant": body.celebrant.strip(),
            "skip_ownership": True,
        },
    )

    out: dict[str, Any] = {
        "ok": True,
        "title": result.title,
        "gospel_reference": result.gospel_reference,
        "slide_excerpt": result.slide_line_preview,
        "gospel_quote": result.gospel_quote,
        "liturgical_color": liturgical_payload,
        "selected_songs": result.selected_songs,
        "slide_count": result.slide_count,
        "export_stem": result.export_stem,
        "ai_poster_urls": _ai_poster_download_urls(),
    }
    if ppt_url:
        out["pptx_url"] = ppt_url
    if poster_url:
        out["poster_url"] = poster_url
    if poster_ppt_url:
        out["poster_ppt_url"] = poster_ppt_url
    if zip_url:
        out["zip_url"] = zip_url
    print(
        f"[generate] done stem={result.export_stem} slides={result.slide_count} pptx={bool(ppt_url)}",
        flush=True,
    )
    return out


@app.post("/api/regenerate-pptx")
async def api_regenerate_pptx(
    body: GenerateBody,
    session: Optional[AuthSession] = Depends(require_approved_membership),
) -> Any:
    """Rebuild the Mass PowerPoint with current hymn typography (overwrites existing .pptx)."""
    song_map = body.songs.model_dump(exclude_none=True) if body.songs else None
    temp_assets: list[Path] = []
    divider_path = None
    if body.divider_poster_basename and str(body.divider_poster_basename).strip():
        raw_divider = str(body.divider_poster_basename).strip()
        divider_path = _resolve_child_file(_MASS_ASSET_DIR, raw_divider)
        if not divider_path and session and storage_ready(session.token):
            key = _safe_storage_leaf(raw_divider, folder="mass_assets")
            divider_path = _materialize_storage_asset(session, f"{session.user.user_id}/{key}")
            if divider_path:
                temp_assets.append(divider_path)
    ann_paths: list[Path] = []
    for raw in body.announcement_basenames or []:
        bn = str(raw).strip()
        if not bn:
            continue
        p = _resolve_child_file(_MASS_ASSET_DIR, bn)
        if not p and session and storage_ready(session.token):
            key = _safe_storage_leaf(bn, folder="mass_assets")
            p = _materialize_storage_asset(session, f"{session.user.user_id}/{key}")
            if p:
                temp_assets.append(p)
        if p:
            ann_paths.append(p)
        if len(ann_paths) >= 24:
            break
    sponsors = [str(s).strip() for s in (body.food_sponsors or []) if str(s).strip()][:24]
    psalm_override = (body.psalm_text_override or "").strip() or None
    gospel_override = (body.gospel_quote_override or "").strip() or None
    try:
        result = await run_in_threadpool(
            regenerate_mass_pptx,
            body.date.strip(),
            body.celebrant.strip(),
            co_celebrant=(body.co_celebrant or "").strip(),
            sentence_index=body.sentence_index,
            song_selections=song_map,
            custom_theme=body.custom_theme,
            hymn_typography=body.hymn_typography,
            divider_poster_path=divider_path,
            lotw_poster=body.lotw_poster,
            lote_poster=body.lote_poster,
            announcement_image_paths=ann_paths or None,
            mass_collection_amount=body.mass_collection_amount.strip()
            if body.mass_collection_amount
            else None,
            mass_collection_date_label=body.mass_collection_date_label.strip()
            if body.mass_collection_date_label
            else None,
            mass_collection_currency=body.mass_collection_currency.strip().upper()
            if body.mass_collection_currency
            else "PHP",
            food_sponsors=sponsors or None,
            psalm_text_override=psalm_override,
            psalm_refrain_index=body.psalm_refrain_index,
            psalm_response_override=(body.psalm_response_override or "").strip() or None,
            gospel_quote_override=gospel_override,
            include_church_logo=body.include_church_logo,
            include_church_name=body.include_church_name,
            include_footer=body.include_footer,
            hymn_lyric_overrides=body.hymn_lyric_overrides,
            creed_choice=body.creed_choice,
            our_father_choice=body.our_father_choice,
            hymn_lyrics_layout=body.hymn_lyrics_layout,
            hymn_layout_overrides=body.hymn_layout_overrides,
        )
    finally:
        for p in temp_assets:
            p.unlink(missing_ok=True)
    if not result.ok:
        raise HTTPException(status_code=400, detail=result.error or "PowerPoint rebuild failed.")
    ppt_url = (
        media_file_url(result.pptx_path.name)
        if result.pptx_path and result.pptx_path.is_file()
        else None
    )
    if session and storage_ready(session.token) and result.pptx_path and result.pptx_path.is_file():
        stem = result.export_stem or "latest"
        _spawn_daemon(
            _upload_generation_assets_best_effort,
            name="mass-regen-upload",
            kwargs={
                "user_id": session.user.user_id,
                "access_token": session.token,
                "export_stem": stem,
                "items": [
                    (
                        f"generated/{stem}/{result.pptx_path.name}",
                        str(result.pptx_path),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                ],
            },
        )
    if session and result.pptx_path and result.pptx_path.is_file():
        register_owned_files(session.user.user_id, [result.pptx_path.name])
    return {
        "ok": True,
        "slide_count": result.slide_count,
        "export_stem": result.export_stem,
        "pptx_url": ppt_url,
        "title": result.title,
    }


@app.post("/api/songs/refresh")
def api_refresh_songs(
    body: RefreshSongsBody,
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    sec = body.section.strip().lower()
    if sec not in {"entrance", "offertory", "communion", "recessional", "meditation"}:
        raise HTTPException(status_code=400, detail="Invalid section.")
    songs = refresh_song_section(
        date=body.date.strip(),
        section=sec,
        current_ids=[str(x) for x in (body.current_ids or [])],
        limit=10,
    )
    return {"ok": True, "section": sec, "songs": songs}


@app.post("/api/songs/refresh-all")
def api_refresh_all_songs(
    body: RefreshAllSongsBody,
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    songs = refresh_all_song_sections(
        date=body.date.strip(),
        current_ids=body.current_ids or {},
        limit=10,
    )
    return {"ok": True, "songs_by_section": songs}


@app.post("/api/songs/import")
def api_import_songs(
    body: ImportSongsBody,
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    summary = import_titles(
        {
            "entrance": body.entrance,
            "offertory": body.offertory,
            "communion": body.communion,
            "recessional": body.recessional,
        }
    )
    return {"ok": True, **summary}


@app.post("/api/songs/import-list")
def api_import_song_rows(
    body: ImportSongRowsBody,
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    summary = import_song_rows([s.model_dump() for s in body.songs])
    return {"ok": True, **summary}


@app.post("/api/songs/fetch-lyrics")
def api_fetch_lyrics(
    body: FetchLyricsBody,
    _session: Optional[AuthSession] = Depends(require_superadmin),
) -> dict[str, Any]:
    if not body.selections:
        return {"ok": True, "updated": 0, "skipped": 0, "results": []}
    results = [fetch_and_store_for_selection({"section": s.section, "id": s.id}) for s in body.selections]
    updated = sum(1 for r in results if r.get("ok") and r.get("reason") == "fetched")
    skipped = len(results) - updated
    return {"ok": True, "updated": updated, "skipped": skipped, "results": results}


@app.post("/api/lyrics/save")
def api_save_lyrics(
    body: SaveLyricsBody,
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    if auth_enabled() and session and not is_superadmin_user(session.user):
        from services.parish_hymn_overrides import save_override
        from services.parish_store import get_user_parish_context
        from services.song_catalog import find_catalog_matches_by_title, format_song_title_case

        clean_title = format_song_title_case(body.title or "")
        exact = [
            m
            for m in find_catalog_matches_by_title(clean_title)
            if m.get("match") == "exact"
        ]
        parish_ctx = get_user_parish_context(session.user.user_id) or {}
        parish_id = str(parish_ctx.get("parish_id") or "").strip()
        # Existing catalog song → parish short version (does not change global SoT).
        if exact and parish_id:
            hit = exact[0]
            result = save_override(
                parish_id,
                hymn_id=str(hit.get("id") or ""),
                section=str((body.sections or [hit.get("section") or "meditation"])[0]),
                lyrics=body.lyrics,
                title=clean_title,
                updated_by=session.user.user_id,
            )
            if not result.get("ok"):
                raise HTTPException(status_code=400, detail=result.get("error") or "Could not save parish lyrics.")
            return result
        result = submit_pending_song(
            session,
            title=body.title,
            lyrics=body.lyrics,
            sections=body.sections,
            language=body.language,
            author=body.author,
        )
        if not result.get("ok"):
            raise HTTPException(
                status_code=409 if result.get("duplicate") else 400,
                detail=result.get("error") or "Could not submit song.",
            )
        return result
    result = save_lyrics_song(
        title=body.title,
        lyrics=body.lyrics,
        sections=body.sections,
        language=body.language,
        author=body.author,
        gospel_moods=body.gospel_moods,
        updated_by=session.user.user_id if session else None,
    )
    if not result.get("ok"):
        raise HTTPException(status_code=400, detail=result.get("error") or "Could not save lyrics.")
    return result


@app.post("/api/submissions/priest")
def api_submit_priest(
    body: PriestSubmissionBody,
    session: Optional[AuthSession] = Depends(require_session_when_auth),
) -> dict[str, Any]:
    if not session:
        raise HTTPException(status_code=401, detail="Sign in required.")
    if is_superadmin_user(session.user):
        from services import community_store

        names = community_store.append_celebrant_name(body.name.strip())
        if supabase_enabled():
            from services.pending_submissions import sync_celebrants_to_supabase_profiles

            sync_celebrants_to_supabase_profiles()
        else:
            update_community(celebrant_names=names)
        payload = _community_api_payload(session)
        return {"ok": True, "celebrant_names": names, **payload}
    result = submit_pending_priest(session, name=body.name)
    if not result.get("ok"):
        raise HTTPException(status_code=400, detail=result.get("error") or "Could not submit priest.")
    return result
